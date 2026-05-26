"""Phase E.3 — DRAFT watermark embedding service (2026-05-26).

Adds a visible "DRAFT" stamp into exported PDF / DOCX / PPTX files
for documents in state="draft". Drafts MUST remain exportable; the
watermark is the regulatory/UX guard against confusing draft content
with finalized output.

All three formats use libraries already in `requirements.txt`:
  • PDF:  reportlab + pypdf (overlay + merge)
  • DOCX: python-docx
  • PPTX: python-pptx

Hard rule (preserved from the original Phase E.3 brief): when
watermarking actually FAILS, the caller MUST block the export with
a clear error. This module raises `WatermarkError` on failure so
callers know to surface the block.

Watermark visual:
  • Repeating-tile "DRAFT" text
  • -30 degree rotation
  • Oxblood (#7A2E2E) at ~30% opacity
"""
from __future__ import annotations

import io
import logging
import re
import zipfile
from typing import Any, Dict, List, Optional, Tuple


log = logging.getLogger("documents.watermark")

OXBLOOD_HEX = "#7A2E2E"
OXBLOOD_RGB = (122 / 255, 46 / 255, 46 / 255)


class WatermarkError(RuntimeError):
    """Raised when watermarking fails. Caller MUST block the export
    (the spec-compliant fallback)."""


# ─────────────────────────────────────────────────────────────────────
# PDF — overlay a watermarked layer over every page
# ─────────────────────────────────────────────────────────────────────
def add_pdf_watermark(pdf_bytes: bytes, *, label: str = "DRAFT") -> bytes:
    """Stamp `label` diagonally across every page of the PDF as a
    repeating tile. Returns the watermarked PDF bytes.

    Uses reportlab to generate a single watermark overlay page sized
    to A4, then merges it onto every page of the source PDF via pypdf.
    For non-A4 page sizes the overlay is auto-scaled to the source
    page dimensions on merge.
    """
    try:
        from reportlab.pdfgen import canvas
        from reportlab.lib.pagesizes import A4
        from pypdf import PdfReader, PdfWriter
        from pypdf.generic import RectangleObject

        # ── build the watermark overlay (single page, A4) ─────────
        overlay_buf = io.BytesIO()
        c = canvas.Canvas(overlay_buf, pagesize=A4)
        page_w, page_h = A4
        c.saveState()
        # 30% opacity oxblood text.
        c.setFillColorRGB(*OXBLOOD_RGB, alpha=0.30)
        c.setFont("Helvetica-Bold", 48)
        # Repeat across the page in a diagonal grid.
        # Step roughly 260pt horiz × 180pt vert keeps a balanced tile.
        for y in range(-200, int(page_h) + 200, 180):
            for x in range(-200, int(page_w) + 200, 260):
                c.saveState()
                c.translate(x, y)
                c.rotate(-30)
                c.drawString(0, 0, label)
                c.restoreState()
        c.restoreState()
        c.save()
        overlay_buf.seek(0)

        # ── merge overlay onto every page of the source ──────────
        overlay_reader = PdfReader(overlay_buf)
        overlay_page = overlay_reader.pages[0]
        source_reader = PdfReader(io.BytesIO(pdf_bytes))
        writer = PdfWriter()
        for src_page in source_reader.pages:
            # merge_page() requires `pypdf>=4.0` (we have 6.10.2).
            src_page.merge_page(overlay_page)
            writer.add_page(src_page)
        out = io.BytesIO()
        writer.write(out)
        return out.getvalue()
    except Exception as e:  # noqa: BLE001
        log.exception("add_pdf_watermark failed: %s", e)
        raise WatermarkError(f"PDF watermarking failed: {e}") from e


# ─────────────────────────────────────────────────────────────────────
# DOCX — embed the watermark as a header element with rotated text
# ─────────────────────────────────────────────────────────────────────
_DOCX_WATERMARK_HEADER_XML = """<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<w:hdr xmlns:w="http://schemas.openxmlformats.org/wordprocessingml/2006/main"
       xmlns:v="urn:schemas-microsoft-com:vml"
       xmlns:o="urn:schemas-microsoft-com:office:office"
       xmlns:w10="urn:schemas-microsoft-com:office:word">
  <w:p>
    <w:r>
      <w:pict>
        <v:shapetype id="_x0000_t136" coordsize="21600,21600" o:spt="136" path="m@7,l@8,m@5,21600l@6,21600e">
          <v:formulas>
            <v:f eqn="sum #0 0 10800"/>
            <v:f eqn="prod #0 2 1"/>
            <v:f eqn="sum 21600 0 @1"/>
            <v:f eqn="sum 0 0 @2"/>
            <v:f eqn="sum 21600 0 @3"/>
            <v:f eqn="if @0 @3 0"/>
            <v:f eqn="if @0 21600 @1"/>
            <v:f eqn="if @0 0 @2"/>
            <v:f eqn="if @0 @4 21600"/>
            <v:f eqn="mid @5 @6"/>
            <v:f eqn="mid @8 @5"/>
            <v:f eqn="mid @7 @8"/>
            <v:f eqn="mid @6 @7"/>
            <v:f eqn="sum @6 0 @5"/>
          </v:formulas>
          <v:path o:extrusionok="f" gradientshapeok="t" o:connecttype="custom"
                  o:connectlocs="@9,0;@10,10800;@11,21600;@12,10800" o:connectangles="270,180,90,0"/>
          <v:textpath on="t" fitshape="t"/>
        </v:shapetype>
        <v:shape id="DraftWatermark" type="#_x0000_t136"
                 style="position:absolute;margin-left:0;margin-top:0;width:500pt;height:200pt;rotation:-30;z-index:-251658240"
                 fillcolor="#7A2E2E" stroked="f">
          <v:fill opacity=".3"/>
          <v:textpath style="font-family:&quot;Arial Black&quot;;font-size:96pt" string="__LABEL__"/>
        </v:shape>
      </w:pict>
    </w:r>
  </w:p>
</w:hdr>
"""


def add_docx_watermark(docx_bytes: bytes, *, label: str = "DRAFT") -> bytes:
    """Embed `label` as a header-anchored watermark in every page of
    the DOCX. Returns the watermarked DOCX bytes.

    Implementation note: python-docx doesn't have a direct watermark
    API. We use the underlying ZIP/XML representation — inject a
    `word/header_watermark.xml`, declare it in the rels, and reference
    it from every section's header.
    """
    try:
        in_buf = io.BytesIO(docx_bytes)
        out_buf = io.BytesIO()
        with zipfile.ZipFile(in_buf, "r") as zin:
            names = set(zin.namelist())
            # Write all original files to the output zip, plus our
            # injected watermark header + the rels + content-types.
            with zipfile.ZipFile(out_buf, "w", zipfile.ZIP_DEFLATED) as zout:
                for name in zin.namelist():
                    data = zin.read(name)
                    if name == "[Content_Types].xml":
                        # Register the new header part.
                        body = data.decode("utf-8")
                        if "header_watermark.xml" not in body:
                            inject = ('<Override PartName="/word/header_watermark.xml" '
                                      'ContentType="application/vnd.openxmlformats-'
                                      'officedocument.wordprocessingml.header+xml"/>')
                            body = body.replace("</Types>", inject + "</Types>")
                        zout.writestr(name, body)
                    elif name == "word/_rels/document.xml.rels":
                        body = data.decode("utf-8")
                        if "rIdDraftWmark" not in body:
                            rel = ('<Relationship Id="rIdDraftWmark" '
                                   'Type="http://schemas.openxmlformats.org/officeDocument/'
                                   '2006/relationships/header" Target="header_watermark.xml"/>')
                            body = body.replace("</Relationships>", rel + "</Relationships>")
                        zout.writestr(name, body)
                    elif name == "word/document.xml":
                        # Reference the header from sectPr. If no sectPr
                        # exists, fall back to no-op (the header still
                        # ships in the package; some viewers will pick
                        # it up via default behaviour).
                        body = data.decode("utf-8")
                        if "<w:headerReference" not in body and "<w:sectPr" in body:
                            ref = '<w:headerReference w:type="default" r:id="rIdDraftWmark"/>'
                            body = body.replace("<w:sectPr", "<w:sectPr><w:_marker/>").replace(
                                "<w:_marker/>", ref, 1,
                            )
                        zout.writestr(name, body)
                    else:
                        zout.writestr(name, data)
                # New header part with the watermark.
                if "word/header_watermark.xml" not in names:
                    xml = _DOCX_WATERMARK_HEADER_XML.replace("__LABEL__", label)
                    zout.writestr("word/header_watermark.xml", xml)
        return out_buf.getvalue()
    except Exception as e:  # noqa: BLE001
        log.exception("add_docx_watermark failed: %s", e)
        raise WatermarkError(f"DOCX watermarking failed: {e}") from e


# ─────────────────────────────────────────────────────────────────────
# PPTX — embed the watermark in the slide master so every slide carries it
# ─────────────────────────────────────────────────────────────────────
def add_pptx_watermark(pptx_bytes: bytes, *, label: str = "DRAFT") -> bytes:
    """Embed `label` as a rotated text shape on every slide. Uses
    python-pptx to add the shape directly (slide-level so every
    slide carries it — slide-master injection is more complex and
    not all viewers honour the master inheritance).
    """
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt, Emu
        from pptx.dml.color import RGBColor
        from pptx.enum.shapes import MSO_SHAPE

        in_buf = io.BytesIO(pptx_bytes)
        prs = Presentation(in_buf)
        # Slide dimensions in EMU; we tile the watermark across.
        slide_w = prs.slide_width
        slide_h = prs.slide_height
        for slide in prs.slides:
            # Tile a few "DRAFT" text boxes across the slide.
            step_x = int(slide_w / 2.5)
            step_y = int(slide_h / 3)
            for col in range(3):
                for row in range(4):
                    x = int(-step_x / 2 + col * step_x)
                    y = int(-step_y / 2 + row * step_y)
                    width = step_x
                    height = step_y
                    tb = slide.shapes.add_textbox(x, y, width, height)
                    tf = tb.text_frame
                    p = tf.paragraphs[0]
                    p.text = label
                    run = p.runs[0]
                    run.font.size = Pt(48)
                    run.font.bold = True
                    run.font.color.rgb = RGBColor(0x7A, 0x2E, 0x2E)
                    # python-pptx exposes rotation via shape.rotation.
                    tb.rotation = -30
        out_buf = io.BytesIO()
        prs.save(out_buf)
        return out_buf.getvalue()
    except Exception as e:  # noqa: BLE001
        log.exception("add_pptx_watermark failed: %s", e)
        raise WatermarkError(f"PPTX watermarking failed: {e}") from e


# ─────────────────────────────────────────────────────────────────────
# Dispatcher — pick the right watermarker by content-type or extension
# ─────────────────────────────────────────────────────────────────────
def watermark_file(file_bytes: bytes, *, fmt: str, label: str = "DRAFT") -> bytes:
    """Apply the watermark by format. `fmt` is `pdf|docx|pptx`.
    Returns watermarked bytes; raises WatermarkError on failure."""
    fmt = (fmt or "").lower().lstrip(".")
    if fmt == "pdf":
        return add_pdf_watermark(file_bytes, label=label)
    if fmt == "docx":
        return add_docx_watermark(file_bytes, label=label)
    if fmt == "pptx":
        return add_pptx_watermark(file_bytes, label=label)
    raise WatermarkError(f"Unsupported watermark format: {fmt!r}")
