"""Phase C.2 — Work Studio export rendering pipeline.

Deterministic-template-first per D-002. Three native renderers built
programmatically with python-docx and python-pptx; a fourth Jinja+
WeasyPrint path produces the PDF variant. Same input bytes always
produce the same output bytes.

Public surface (the only functions the router calls):
    render_brief_docx(content, ctx_meta)   -> (bytes, sha256, filename)
    render_deck_pptx(content, ctx_meta)    -> (bytes, sha256, filename)
    render_report_docx(content, ctx_meta)  -> (bytes, sha256, filename)
    render_brief_pdf(content, ctx_meta)    -> (bytes, sha256, filename)
    render_report_pdf(content, ctx_meta)   -> (bytes, sha256, filename)
    render_deck_pdf(...)                   -> NotImplementedError (soft fork)
    scan_for_banned_words(text)            -> Optional[str]
    validate_content(content, kind)        -> raises ContentValidationError on bad shape

`content` is the JSON content_dict produced by the LLM stage (schema in
`backend/templates/work_studio/SCHEMAS.md`). `ctx_meta` is the small
metadata bag built by the router.
"""
from __future__ import annotations

import hashlib
import io
import logging
import re
from datetime import datetime, timezone
from typing import Any, Dict, List, Optional, Tuple

from docx import Document
from docx.shared import Inches, Pt, RGBColor
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.oxml.ns import qn
from docx.oxml import OxmlElement
from pptx import Presentation
from pptx.util import Inches as PInches, Pt as PPt
from pptx.dml.color import RGBColor as PRGBColor
from pptx.enum.text import PP_ALIGN
from jinja2 import Template
from weasyprint import HTML

from services import two_pass as _tp

logger = logging.getLogger("akki.work_studio_export")


# =============================================================================
# Shared design tokens (mirrors design_guidelines.json)
# =============================================================================
class Palette:
    CREAM = RGBColor(0xF7, 0xF3, 0xEA)
    INK = RGBColor(0x11, 0x18, 0x27)
    NAVY = RGBColor(0x0A, 0x1F, 0x44)
    OXBLOOD = RGBColor(0x8B, 0x2E, 0x2B)
    MUTED = RGBColor(0x4B, 0x55, 0x63)
    RULE = RGBColor(0xE5, 0xE0, 0xD3)


class PPalette:  # python-pptx wants its own RGBColor type
    CREAM = PRGBColor(0xF7, 0xF3, 0xEA)
    INK = PRGBColor(0x11, 0x18, 0x27)
    NAVY = PRGBColor(0x0A, 0x1F, 0x44)
    OXBLOOD = PRGBColor(0x8B, 0x2E, 0x2B)
    MUTED = PRGBColor(0x4B, 0x55, 0x63)


CLASSIFICATION_COLOURS = {
    "Public":       Palette.MUTED,
    "Internal":     Palette.NAVY,
    "Confidential": Palette.OXBLOOD,
    "Restricted":   Palette.OXBLOOD,
}


VALID_CLASSIFICATIONS = {"Public", "Internal", "Confidential", "Restricted"}


# =============================================================================
# Content validation — runs before any render call
# =============================================================================
class ContentValidationError(ValueError):
    """Raised when the LLM-produced content_dict fails schema enforcement."""


def _normalise_classification(c: Any) -> str:
    s = (c or "Internal").strip().capitalize()
    if s not in VALID_CLASSIFICATIONS:
        return "Internal"
    return s


def validate_content(content: Dict[str, Any], kind: str) -> Dict[str, Any]:
    """Validate + normalise. Returns the cleaned content_dict.
    Raises ContentValidationError on hard schema violations.

    Citation enforcement is the key invariant: every `cites` index must
    point to an entry in the `citations` manifest. We DO NOT silently
    drop bad indices — bad indices indicate the LLM fabricated a source.
    """
    if kind not in ("brief", "deck", "report"):
        raise ContentValidationError(f"Unknown kind: {kind}")

    title = (content.get("title") or "").strip()
    if not title:
        raise ContentValidationError("`title` is required.")
    if len(title) > 140:
        title = title[:140]

    subtitle = (content.get("subtitle") or "").strip()[:200] or None
    classification = _normalise_classification(content.get("classification"))
    period = (content.get("period") or "").strip() or "—"
    generated_for = (content.get("generated_for") or "—").strip()

    exec_summary = (content.get("executive_summary") or "").strip()
    if not exec_summary:
        raise ContentValidationError("`executive_summary` is required.")
    # Normalise whitespace in the exec summary.
    exec_summary = re.sub(r"\s+", " ", exec_summary)

    citations_raw = content.get("citations") or []
    if not isinstance(citations_raw, list):
        raise ContentValidationError("`citations` must be a list.")
    citations: List[Dict[str, Any]] = []
    for c in citations_raw:
        if not isinstance(c, dict):
            continue
        doc_id = (c.get("doc_id") or "").strip()
        doc_name = (c.get("doc_name") or "Source").strip()
        if not doc_id:
            continue
        citations.append({
            "doc_id": doc_id,
            "doc_name": doc_name,
            "paragraph_anchor": c.get("paragraph_anchor") or None,
        })

    sections_raw = content.get("sections") or []
    if not isinstance(sections_raw, list) or not sections_raw:
        raise ContentValidationError("`sections` must be a non-empty list.")

    sections: List[Dict[str, Any]] = []
    n_cit = len(citations)
    for s in sections_raw:
        if not isinstance(s, dict):
            continue
        heading = (s.get("heading") or "").strip()
        if not heading:
            continue
        cites = s.get("cites") or []
        # Normalise + validate citations indices.
        norm_cites = []
        for idx in cites:
            try:
                i = int(idx)
            except (TypeError, ValueError):
                continue
            if 1 <= i <= n_cit:
                norm_cites.append(i)
            else:
                # Fabricated citation index — this is a hard fail.
                raise ContentValidationError(
                    f"Section `{heading[:60]}` references citation [{i}] "
                    f"but only {n_cit} citations are declared."
                )
        if kind == "deck":
            bullets = [b for b in (s.get("bullets") or []) if isinstance(b, str) and b.strip()]
            sections.append({
                "heading": heading,
                "bullets": [re.sub(r"\s+", " ", b.strip()) for b in bullets][:6],
                "callout": (s.get("callout") or "").strip() or None,
                "cites": norm_cites,
            })
        else:  # brief / report
            paras = [p for p in (s.get("paragraphs") or []) if isinstance(p, str) and p.strip()]
            sections.append({
                "heading": heading,
                "subheading": (s.get("subheading") or "").strip() or None,
                "paragraphs": [re.sub(r"\s+", " ", p.strip()) for p in paras][:6],
                "pullquote": (s.get("pullquote") or "").strip() or None,
                "cites": norm_cites,
            })

    if not sections:
        raise ContentValidationError("`sections` produced no usable rows.")

    out = {
        "title": title,
        "subtitle": subtitle,
        "classification": classification,
        "period": period,
        "generated_for": generated_for,
        "executive_summary": exec_summary,
        "sections": sections,
        "citations": citations,
    }
    if kind == "report":
        out["recommendations"] = [
            r.strip() for r in (content.get("recommendations") or [])
            if isinstance(r, str) and r.strip()
        ][:6]
    if kind == "deck":
        out["conclusion"] = (content.get("conclusion") or "").strip() or None
    return out


# =============================================================================
# Banned-word scanner (post-render)
# =============================================================================
def scan_for_banned_words(text: str) -> Optional[str]:
    """Scan rendered/visible text for banned words. Returns the first
    hit (lowercased) or None. Reuses the union list from `two_pass`."""
    return _tp.find_banned_word(text)


# =============================================================================
# DOCX helpers
# =============================================================================
def _set_cell_shading(cell, hex_colour: str) -> None:
    """python-docx exposes no public API for cell shading; manipulate the XML directly."""
    tc_pr = cell._tc.get_or_add_tcPr()
    shd = OxmlElement("w:shd")
    shd.set(qn("w:val"), "clear")
    shd.set(qn("w:color"), "auto")
    shd.set(qn("w:fill"), hex_colour)
    tc_pr.append(shd)


def _add_classification_chip(doc: Document, classification: str) -> None:
    """Add a small 'Classification: X' bar at the top of the cover page."""
    p = doc.add_paragraph()
    p.alignment = WD_ALIGN_PARAGRAPH.RIGHT
    run = p.add_run(f"  CLASSIFICATION  ·  {classification.upper()}  ")
    run.font.name = "Calibri"
    run.font.size = Pt(8)
    run.font.bold = True
    run.font.color.rgb = Palette.CREAM
    # Colour-code the chip background via the run's shading.
    pPr = run._element.get_or_add_rPr()
    shd = OxmlElement("w:shd")
    fill = "8B2E2B" if classification in ("Confidential", "Restricted") else "0A1F44"
    shd.set(qn("w:val"), "clear")
    shd.set(qn("w:color"), "auto")
    shd.set(qn("w:fill"), fill)
    pPr.append(shd)


def _set_doc_defaults(doc: Document) -> None:
    """Set the body 'Normal' style font to Georgia / 11pt / Ink."""
    style = doc.styles["Normal"]
    style.font.name = "Georgia"
    style.font.size = Pt(11)
    style.font.color.rgb = Palette.INK
    rPr = style.element.get_or_add_rPr()
    # East-Asian + ASCII fall-throughs so Georgia is selected on most renderers.
    rFonts = rPr.find(qn("w:rFonts")) or OxmlElement("w:rFonts")
    rFonts.set(qn("w:ascii"), "Georgia")
    rFonts.set(qn("w:hAnsi"), "Georgia")
    rFonts.set(qn("w:eastAsia"), "Georgia")
    rPr.append(rFonts)


def _add_heading_calibri(doc: Document, text: str, level: int = 1, ink: RGBColor = Palette.NAVY) -> None:
    p = doc.add_paragraph()
    p.paragraph_format.space_before = Pt(20 if level == 1 else 14)
    p.paragraph_format.space_after = Pt(6)
    run = p.add_run(text)
    run.font.name = "Calibri"
    run.font.size = Pt(20 if level == 1 else 14)
    run.font.bold = True
    run.font.color.rgb = ink


def _add_paragraph_georgia(doc: Document, text: str, *, italic: bool = False, color: Optional[RGBColor] = None) -> None:
    p = doc.add_paragraph()
    p.paragraph_format.space_after = Pt(6)
    run = p.add_run(text)
    run.font.name = "Georgia"
    run.font.size = Pt(11)
    run.italic = italic
    if color is not None:
        run.font.color.rgb = color
    else:
        run.font.color.rgb = Palette.INK


def _add_pullquote(doc: Document, text: str) -> None:
    p = doc.add_paragraph()
    p.paragraph_format.left_indent = Inches(0.5)
    p.paragraph_format.right_indent = Inches(0.5)
    p.paragraph_format.space_before = Pt(8)
    p.paragraph_format.space_after = Pt(8)
    run = p.add_run("\u201C" + text.strip() + "\u201D")
    run.font.name = "Georgia"
    run.font.size = Pt(13)
    run.italic = True
    run.font.color.rgb = Palette.NAVY


def _cite_marker(cites: List[int]) -> str:
    """Render '[1, 3]' style inline marker; empty if no cites."""
    if not cites:
        return ""
    return " " + ", ".join(f"[{i}]" for i in cites)


def _add_citations_appendix(doc: Document, citations: List[Dict[str, Any]]) -> None:
    if not citations:
        return
    _add_heading_calibri(doc, "Sources", level=1)
    for i, c in enumerate(citations, start=1):
        p = doc.add_paragraph()
        p.paragraph_format.space_after = Pt(3)
        head = p.add_run(f"[{i}]  ")
        head.font.name = "Calibri"
        head.font.size = Pt(10)
        head.font.bold = True
        head.font.color.rgb = Palette.OXBLOOD
        body = p.add_run(c["doc_name"])
        body.font.name = "Georgia"
        body.font.size = Pt(10)
        body.font.color.rgb = Palette.INK


def _bytes_and_hash(buf: io.BytesIO) -> Tuple[bytes, str]:
    data = buf.getvalue()
    return data, hashlib.sha256(data).hexdigest()


def _safe_filename(stem: str, ext: str) -> str:
    s = re.sub(r"[^a-zA-Z0-9]+", "-", stem.strip()).strip("-").lower() or "akki-export"
    return f"{s[:60]}.{ext}"


# =============================================================================
# Brief — DOCX
# =============================================================================
def render_brief_docx(content: Dict[str, Any], ctx_meta: Dict[str, Any]) -> Tuple[bytes, str, str]:
    c = validate_content(content, "brief")
    doc = Document()
    _set_doc_defaults(doc)

    # Margins
    for section in doc.sections:
        section.top_margin = Inches(1.0)
        section.bottom_margin = Inches(1.0)
        section.left_margin = Inches(1.1)
        section.right_margin = Inches(1.1)

    # Cover page
    _add_classification_chip(doc, c["classification"])
    title_p = doc.add_paragraph()
    title_p.paragraph_format.space_before = Pt(60)
    title_p.paragraph_format.space_after = Pt(8)
    tr = title_p.add_run(c["title"])
    tr.font.name = "Georgia"
    tr.font.size = Pt(28)
    tr.font.bold = True
    tr.font.color.rgb = Palette.NAVY
    if c["subtitle"]:
        sub = doc.add_paragraph()
        sub.paragraph_format.space_after = Pt(4)
        sr = sub.add_run(c["subtitle"])
        sr.font.name = "Georgia"
        sr.font.size = Pt(14)
        sr.italic = True
        sr.font.color.rgb = Palette.MUTED
    meta_p = doc.add_paragraph()
    meta_p.paragraph_format.space_before = Pt(12)
    mr = meta_p.add_run(f"Prepared for {c['generated_for']}  ·  {c['period']}  ·  {ctx_meta.get('generated_at_human','')}")
    mr.font.name = "Calibri"
    mr.font.size = Pt(10)
    mr.font.color.rgb = Palette.MUTED
    doc.add_page_break()

    # Executive summary
    _add_heading_calibri(doc, "Executive summary", level=1)
    _add_paragraph_georgia(doc, c["executive_summary"])
    doc.add_paragraph()

    # Body sections
    for s in c["sections"]:
        _add_heading_calibri(doc, s["heading"], level=2)
        if s.get("subheading"):
            _add_paragraph_georgia(doc, s["subheading"], italic=True, color=Palette.MUTED)
        for i, para in enumerate(s["paragraphs"]):
            tail = _cite_marker(s["cites"]) if i == len(s["paragraphs"]) - 1 else ""
            _add_paragraph_georgia(doc, para + tail)
        if s.get("pullquote"):
            _add_pullquote(doc, s["pullquote"])

    # Sources appendix
    doc.add_page_break()
    _add_citations_appendix(doc, c["citations"])

    # Footer (per-section properties — apply to every section in the doc).
    for section in doc.sections:
        footer_p = section.footer.paragraphs[0]
        footer_p.alignment = WD_ALIGN_PARAGRAPH.CENTER
        fr = footer_p.add_run(f"  AKKI  ·  {c['generated_for']}  ·  {c['classification']}  ")
        fr.font.name = "Calibri"
        fr.font.size = Pt(8)
        fr.font.color.rgb = Palette.MUTED

    buf = io.BytesIO()
    doc.save(buf)
    data, sha = _bytes_and_hash(buf)
    fn = _safe_filename(c["title"], "docx")
    return data, sha, fn


# =============================================================================
# Report — DOCX (heavier than Brief)
# =============================================================================
def render_report_docx(content: Dict[str, Any], ctx_meta: Dict[str, Any]) -> Tuple[bytes, str, str]:
    c = validate_content(content, "report")
    doc = Document()
    _set_doc_defaults(doc)

    for section in doc.sections:
        section.top_margin = Inches(1.0)
        section.bottom_margin = Inches(1.0)
        section.left_margin = Inches(1.1)
        section.right_margin = Inches(1.1)

    # Cover
    _add_classification_chip(doc, c["classification"])
    p = doc.add_paragraph()
    p.paragraph_format.space_before = Pt(80)
    r = p.add_run(c["title"])
    r.font.name = "Georgia"
    r.font.size = Pt(32)
    r.font.bold = True
    r.font.color.rgb = Palette.NAVY
    if c["subtitle"]:
        sub = doc.add_paragraph()
        sr = sub.add_run(c["subtitle"])
        sr.font.name = "Georgia"
        sr.font.size = Pt(15)
        sr.italic = True
        sr.font.color.rgb = Palette.MUTED
    meta = doc.add_paragraph()
    meta.paragraph_format.space_before = Pt(20)
    mr = meta.add_run(f"Prepared for {c['generated_for']}  ·  {c['period']}  ·  {ctx_meta.get('generated_at_human','')}")
    mr.font.name = "Calibri"
    mr.font.size = Pt(10)
    mr.font.color.rgb = Palette.MUTED
    doc.add_page_break()

    # Executive summary
    _add_heading_calibri(doc, "Executive summary", level=1)
    _add_paragraph_georgia(doc, c["executive_summary"])
    doc.add_page_break()

    # Index
    _add_heading_calibri(doc, "Index", level=1)
    for i, s in enumerate(c["sections"], start=1):
        ip = doc.add_paragraph()
        ir = ip.add_run(f"{i:02d}.  {s['heading']}")
        ir.font.name = "Calibri"
        ir.font.size = Pt(11)
        ir.font.color.rgb = Palette.INK
    doc.add_page_break()

    # Body
    for s in c["sections"]:
        _add_heading_calibri(doc, s["heading"], level=1)
        if s.get("subheading"):
            _add_paragraph_georgia(doc, s["subheading"], italic=True, color=Palette.MUTED)
        for i, para in enumerate(s["paragraphs"]):
            tail = _cite_marker(s["cites"]) if i == len(s["paragraphs"]) - 1 else ""
            _add_paragraph_georgia(doc, para + tail)
        if s.get("pullquote"):
            _add_pullquote(doc, s["pullquote"])
        doc.add_paragraph()

    # Recommendations
    if c.get("recommendations"):
        doc.add_page_break()
        _add_heading_calibri(doc, "Recommendations", level=1)
        for i, rec in enumerate(c["recommendations"], start=1):
            rp = doc.add_paragraph()
            rp.paragraph_format.space_after = Pt(4)
            rr = rp.add_run(f"{i}.  ")
            rr.font.name = "Calibri"
            rr.font.size = Pt(11)
            rr.font.bold = True
            rr.font.color.rgb = Palette.OXBLOOD
            br = rp.add_run(rec)
            br.font.name = "Georgia"
            br.font.size = Pt(11)
            br.font.color.rgb = Palette.INK

    # Sources
    doc.add_page_break()
    _add_citations_appendix(doc, c["citations"])

    for section in doc.sections:
        footer_p = section.footer.paragraphs[0]
        footer_p.alignment = WD_ALIGN_PARAGRAPH.CENTER
        fr = footer_p.add_run(f"  AKKI  ·  {c['generated_for']}  ·  {c['classification']}  ")
        fr.font.name = "Calibri"
        fr.font.size = Pt(8)
        fr.font.color.rgb = Palette.MUTED

    buf = io.BytesIO()
    doc.save(buf)
    data, sha = _bytes_and_hash(buf)
    return data, sha, _safe_filename(c["title"], "docx")


# =============================================================================
# Deck — PPTX
# =============================================================================
def _pptx_set_slide_bg(slide, rgb: PRGBColor) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = rgb


def _pptx_add_text(slide, left, top, width, height, text, *, font="Calibri", size=18, bold=False, color=PPalette.INK, align=PP_ALIGN.LEFT):
    tx = slide.shapes.add_textbox(left, top, width, height)
    tf = tx.text_frame
    tf.word_wrap = True
    tf.text = text
    p = tf.paragraphs[0]
    p.alignment = align
    for run in p.runs:
        run.font.name = font
        run.font.size = PPt(size)
        run.font.bold = bold
        run.font.color.rgb = color
    return tx


def render_deck_pptx(content: Dict[str, Any], ctx_meta: Dict[str, Any]) -> Tuple[bytes, str, str]:
    c = validate_content(content, "deck")
    prs = Presentation()
    prs.slide_width = PInches(13.333)   # widescreen 16:9
    prs.slide_height = PInches(7.5)
    blank = prs.slide_layouts[6]        # blank layout

    # Slide 1 — Title
    s1 = prs.slides.add_slide(blank)
    _pptx_set_slide_bg(s1, PPalette.CREAM)
    _pptx_add_text(s1, PInches(0.6), PInches(0.4), PInches(12), PInches(0.4),
                   f"CLASSIFICATION  ·  {c['classification'].upper()}",
                   font="Calibri", size=10, bold=True,
                   color=PPalette.OXBLOOD if c["classification"] in ("Confidential", "Restricted") else PPalette.NAVY)
    _pptx_add_text(s1, PInches(0.6), PInches(2.0), PInches(12), PInches(2.5),
                   c["title"], font="Georgia", size=42, bold=True, color=PPalette.NAVY)
    if c["subtitle"]:
        _pptx_add_text(s1, PInches(0.6), PInches(3.7), PInches(12), PInches(1.0),
                       c["subtitle"], font="Georgia", size=20, color=PPalette.MUTED)
    _pptx_add_text(s1, PInches(0.6), PInches(6.6), PInches(12), PInches(0.4),
                   f"Prepared for {c['generated_for']}  ·  {c['period']}  ·  {ctx_meta.get('generated_at_human','')}",
                   font="Calibri", size=11, color=PPalette.MUTED)

    # Slide 2 — Executive summary
    s2 = prs.slides.add_slide(blank)
    _pptx_set_slide_bg(s2, PPalette.CREAM)
    _pptx_add_text(s2, PInches(0.6), PInches(0.5), PInches(12), PInches(0.5),
                   "Executive summary", font="Calibri", size=18, bold=True, color=PPalette.NAVY)
    _pptx_add_text(s2, PInches(0.6), PInches(1.4), PInches(12), PInches(5.5),
                   c["executive_summary"], font="Georgia", size=20, color=PPalette.INK)

    # Slides 3..N — Body
    for s in c["sections"]:
        sl = prs.slides.add_slide(blank)
        _pptx_set_slide_bg(sl, PPalette.CREAM)
        _pptx_add_text(sl, PInches(0.6), PInches(0.5), PInches(12), PInches(0.7),
                       s["heading"], font="Calibri", size=24, bold=True, color=PPalette.NAVY)
        # Bullets
        bulleted = ""
        for b in s.get("bullets", [])[:5]:
            bulleted += f"\u2022  {b}\n"
        if s.get("cites"):
            bulleted += f"\n{', '.join(f'[{i}]' for i in s['cites'])}"
        if bulleted:
            _pptx_add_text(sl, PInches(0.6), PInches(1.7), PInches(12), PInches(4.5),
                           bulleted.strip(), font="Georgia", size=18, color=PPalette.INK)
        if s.get("callout"):
            _pptx_add_text(sl, PInches(0.6), PInches(6.3), PInches(12), PInches(0.7),
                           "\u201C" + s["callout"] + "\u201D",
                           font="Georgia", size=14, color=PPalette.OXBLOOD)

    # Slide N+1 — Conclusion
    sN1 = prs.slides.add_slide(blank)
    _pptx_set_slide_bg(sN1, PPalette.NAVY)  # navy enterprise band
    _pptx_add_text(sN1, PInches(0.6), PInches(0.5), PInches(12), PInches(0.7),
                   "What this means", font="Calibri", size=22, bold=True, color=PPalette.CREAM)
    conclusion = c.get("conclusion") or c["executive_summary"].split(".")[0] + "."
    _pptx_add_text(sN1, PInches(0.6), PInches(1.7), PInches(12), PInches(4.5),
                   conclusion, font="Georgia", size=22, color=PPalette.CREAM)

    # Slide N+2 — Sources
    sS = prs.slides.add_slide(blank)
    _pptx_set_slide_bg(sS, PPalette.CREAM)
    _pptx_add_text(sS, PInches(0.6), PInches(0.5), PInches(12), PInches(0.7),
                   "Sources", font="Calibri", size=22, bold=True, color=PPalette.NAVY)
    if c["citations"]:
        body = ""
        for i, cit in enumerate(c["citations"], start=1):
            body += f"[{i}]  {cit['doc_name']}\n"
        _pptx_add_text(sS, PInches(0.6), PInches(1.7), PInches(12), PInches(5.0),
                       body.rstrip(), font="Calibri", size=14, color=PPalette.INK)
    else:
        _pptx_add_text(sS, PInches(0.6), PInches(1.7), PInches(12), PInches(5.0),
                       "No sources cited.", font="Calibri", size=14, color=PPalette.MUTED)

    buf = io.BytesIO()
    prs.save(buf)
    data, sha = _bytes_and_hash(buf)
    return data, sha, _safe_filename(c["title"], "pptx")


# =============================================================================
# PDF — Jinja → WeasyPrint
# =============================================================================
_HTML_TEMPLATE = """<!doctype html>
<html><head><meta charset="utf-8"><title>{{ c.title }}</title>
<style>
  @page { size: A4; margin: 22mm 20mm; }
  body { font-family: Georgia, serif; color: #111827; font-size: 11pt; line-height: 1.55; }
  .chip { display: inline-block; padding: 2px 8px; font-family: Calibri, sans-serif; font-weight: 700;
          font-size: 8pt; letter-spacing: 0.16em; color: #F7F3EA;
          background: {{ chip_bg }}; }
  h1.cover { font-family: Georgia, serif; font-size: 28pt; color: #0A1F44; margin: 60px 0 8px; }
  h1.cover.report { font-size: 32pt; }
  .subtitle { font-style: italic; color: #4B5563; font-size: 14pt; margin-bottom: 12px; }
  .meta { font-family: Calibri, sans-serif; color: #4B5563; font-size: 10pt; margin-top: 18px; }
  h2 { font-family: Calibri, sans-serif; color: #0A1F44; font-size: 18pt; margin: 22px 0 6px; }
  h3 { font-family: Calibri, sans-serif; color: #0A1F44; font-size: 14pt; margin: 14px 0 4px; }
  .pullquote { color: #0A1F44; font-style: italic; margin: 8pt 0 8pt 18pt; padding-left: 10pt;
               border-left: 2pt solid #8B2E2B; font-size: 13pt; }
  .citelist { padding-left: 0; list-style: none; }
  .citelist li { font-family: Calibri, sans-serif; font-size: 10pt; margin-bottom: 4pt; }
  .citelist li b { color: #8B2E2B; }
  .footer { position: running(footer); text-align: center; font-family: Calibri, sans-serif;
            font-size: 8pt; color: #4B5563; }
  @page { @bottom-center { content: element(footer); } }
  .pagebreak { page-break-after: always; }
  .index-row { font-family: Calibri, sans-serif; font-size: 11pt; margin: 4pt 0; }
  .recommend { margin-bottom: 6pt; }
  .recommend b { color: #8B2E2B; font-family: Calibri, sans-serif; }
</style></head>
<body>
  <div class="footer">  AKKI  ·  {{ c.generated_for }}  ·  {{ c.classification }}  </div>

  <div style="text-align:right;"><span class="chip">CLASSIFICATION  ·  {{ c.classification|upper }}</span></div>
  <h1 class="cover {{ kind }}">{{ c.title }}</h1>
  {% if c.subtitle %}<div class="subtitle">{{ c.subtitle }}</div>{% endif %}
  <div class="meta">Prepared for {{ c.generated_for }}  ·  {{ c.period }}  ·  {{ generated_at_human }}</div>
  <div class="pagebreak"></div>

  <h2>Executive summary</h2>
  <p>{{ c.executive_summary }}</p>

  {% if kind == "report" %}
    <div class="pagebreak"></div>
    <h2>Index</h2>
    {% for s in c.sections %}
      <div class="index-row">{{ "%02d" | format(loop.index) }}.  {{ s.heading }}</div>
    {% endfor %}
    <div class="pagebreak"></div>
  {% endif %}

  {% for s in c.sections %}
    <h2>{{ s.heading }}</h2>
    {% if s.subheading %}<h3>{{ s.subheading }}</h3>{% endif %}
    {% for p in s.paragraphs %}
      <p>{{ p }}{% if loop.last and s.cites %} {{ s.cites | map("string") | map("regex_replace", "^(.*)$", "[\\1]") | join(", ") }}{% endif %}</p>
    {% endfor %}
    {% if s.pullquote %}<div class="pullquote">\u201c{{ s.pullquote }}\u201d</div>{% endif %}
  {% endfor %}

  {% if kind == "report" and c.recommendations %}
    <div class="pagebreak"></div>
    <h2>Recommendations</h2>
    {% for r in c.recommendations %}
      <p class="recommend"><b>{{ loop.index }}.</b>  {{ r }}</p>
    {% endfor %}
  {% endif %}

  <div class="pagebreak"></div>
  <h2>Sources</h2>
  {% if c.citations %}
    <ul class="citelist">
      {% for cit in c.citations %}
        <li><b>[{{ loop.index }}]</b>  {{ cit.doc_name }}</li>
      {% endfor %}
    </ul>
  {% else %}
    <p>No sources cited.</p>
  {% endif %}
</body></html>
"""


def _regex_replace(value, pattern, repl):
    return re.sub(pattern, repl, str(value))


def _render_pdf(content: Dict[str, Any], ctx_meta: Dict[str, Any], kind: str) -> Tuple[bytes, str, str]:
    c = validate_content(content, kind)
    chip_bg = "#8B2E2B" if c["classification"] in ("Confidential", "Restricted") else "#0A1F44"
    tmpl = Template(_HTML_TEMPLATE)
    tmpl.environment.filters["regex_replace"] = _regex_replace
    html = tmpl.render(
        c=c, kind=kind, chip_bg=chip_bg,
        generated_at_human=ctx_meta.get("generated_at_human", ""),
    )
    pdf_bytes = HTML(string=html).write_pdf()
    sha = hashlib.sha256(pdf_bytes).hexdigest()
    return pdf_bytes, sha, _safe_filename(c["title"], "pdf")


def render_brief_pdf(content: Dict[str, Any], ctx_meta: Dict[str, Any]) -> Tuple[bytes, str, str]:
    return _render_pdf(content, ctx_meta, "brief")


def render_report_pdf(content: Dict[str, Any], ctx_meta: Dict[str, Any]) -> Tuple[bytes, str, str]:
    return _render_pdf(content, ctx_meta, "report")


def render_deck_pdf(content: Dict[str, Any], ctx_meta: Dict[str, Any]) -> Tuple[bytes, str, str]:
    """Soft fork — PDF rendering for decks requires headless office or
    image-based slide rasterisation. The dev pod has neither libreoffice
    nor a slide-image rasteriser available. The router treats this as
    a `bad output_format` for kind=deck and forces 'pptx' instead.
    Documented in the C.2 acceptance reply.
    """
    raise NotImplementedError(
        "Deck PDF export is not supported in the current dev pod "
        "(libreoffice and image-rasteriser fallback are both absent). "
        "Use output_format='pptx' for decks."
    )


# =============================================================================
# Plain-text scrape (post-render banned-word check)
# =============================================================================
def scrape_docx_text(docx_bytes: bytes) -> str:
    doc = Document(io.BytesIO(docx_bytes))
    parts: List[str] = []
    for p in doc.paragraphs:
        parts.append(p.text)
    for t in doc.tables:
        for row in t.rows:
            for cell in row.cells:
                parts.append(cell.text)
    return "\n".join(parts)


def scrape_pptx_text(pptx_bytes: bytes) -> str:
    prs = Presentation(io.BytesIO(pptx_bytes))
    parts: List[str] = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for p in shape.text_frame.paragraphs:
                    parts.append("".join(r.text for r in p.runs))
    return "\n".join(parts)


def scrape_text(file_bytes: bytes, ext: str) -> str:
    if ext == "docx":
        return scrape_docx_text(file_bytes)
    if ext == "pptx":
        return scrape_pptx_text(file_bytes)
    if ext == "pdf":
        # Re-render the source HTML's text equivalent via the validated
        # content_dict — we don't parse PDFs back. The router calls
        # scrape_text on the parallel docx/pptx representation when the
        # output is PDF; for PDF-only outputs (rare here) we fall back
        # to scanning the content_dict text directly via a sibling util.
        return ""
    return ""


def scrape_content_text(content: Dict[str, Any]) -> str:
    """Fallback text source — concatenate all visible strings from the
    validated content_dict. Used by the post-render scan when the output
    is a PDF or when the parallel docx/pptx scrape is unavailable."""
    parts: List[str] = []
    parts.append(content.get("title") or "")
    parts.append(content.get("subtitle") or "")
    parts.append(content.get("executive_summary") or "")
    parts.append(content.get("conclusion") or "")
    for s in content.get("sections") or []:
        parts.append(s.get("heading") or "")
        parts.append(s.get("subheading") or "")
        parts.extend(s.get("paragraphs") or [])
        parts.extend(s.get("bullets") or [])
        parts.append(s.get("callout") or "")
        parts.append(s.get("pullquote") or "")
    parts.extend(content.get("recommendations") or [])
    return "\n".join(p for p in parts if p)


# =============================================================================
# Generated-at human stamp helper
# =============================================================================
def now_human() -> str:
    return datetime.now(timezone.utc).strftime("%d %b %Y · %H:%M UTC")
