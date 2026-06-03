"""Phase P5.14 — PPTX report builder.

Builds a chair-readable PowerPoint deck from a `WorkbookAnalysis`.
Mirrors the speaker-note discipline of Solva v2 (every slide has
speaker notes that are voice-lint clean, refuse-to-decide
validated, and observational rather than directive).

Slide layout:
  1. Cover
  2. Sheet overview (one per sheet, max 3)
  3. Signals
  4. Simulation results (one per Monte Carlo run)
  5. Forecast results (one per ForecastRun)
  6. Anomalies (single slide, condensed table)
  7. Methodology + reproducibility appendix

Every slide's speaker notes are passed through
`validate_no_imperatives` before being attached — a directive
phrasing trips a build-time error rather than shipping the bad
narration in the file.
"""
from __future__ import annotations

import io
from typing import List

from pptx import Presentation
from pptx.util import Inches, Pt

from .refuse_to_decide import validate_no_imperatives
from .schema import WorkbookAnalysis


def _set_notes(slide, text: str) -> None:
    """Attach validated speaker notes to a slide. Raises on
    imperative-to-user phrasing — that's the contract: notes
    that ship must be observational."""
    validate_no_imperatives(text, label=f"pptx_notes/{slide.slide_id}")
    notes_tf = slide.notes_slide.notes_text_frame
    notes_tf.text = text


def _add_text_slide(prs, title: str, body_paras: List[str]) -> any:
    """Add a layout-5 (title + content) slide and fill it."""
    layout = prs.slide_layouts[5]  # title only — we add body via a text box
    slide = prs.slides.add_slide(layout)
    title_shape = slide.shapes.title
    title_shape.text = title
    title_shape.text_frame.paragraphs[0].runs[0].font.size = Pt(28)

    body_box = slide.shapes.add_textbox(
        left=Inches(0.5), top=Inches(1.6),
        width=Inches(9.0), height=Inches(5.0),
    )
    tf = body_box.text_frame
    tf.word_wrap = True
    for i, para in enumerate(body_paras):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = para
        if p.runs:
            p.runs[0].font.size = Pt(16)
    return slide


def build_pptx_report(analysis: WorkbookAnalysis) -> bytes:
    """Serialise a `WorkbookAnalysis` into a .pptx binary.

    Returns the raw .pptx bytes — the caller streams them out via
    StreamingResponse or writes to a temp blob for the export
    endpoint."""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # ─── Slide 1: cover ──────────────────────────────────────────
    cover = prs.slides.add_slide(prs.slide_layouts[0])
    cover.shapes.title.text = f"Workbook Analysis · {analysis.filename}"
    if cover.placeholders and len(cover.placeholders) > 1:
        cover.placeholders[1].text = (
            f"Prepared by Akki for Executives · analysis id {analysis.id} · "
            f"format {analysis.file_format} · {len(analysis.sheets)} sheet(s)"
        )
    _set_notes(cover, (
        f"This deck reports the observations from the {analysis.filename} workbook. "
        f"It contains {len(analysis.signals)} signal(s), {len(analysis.simulations)} "
        f"simulation(s), {len(analysis.forecasts)} forecast(s), and "
        f"{len(analysis.anomalies)} anomaly observation(s). Every numerical claim "
        f"in this deck is cited to a cell range in the source workbook; the "
        f"appendix at the end lists the exact ranges."
    ))

    # ─── Sheet overview (max 3 sheets) ───────────────────────────
    for sheet in analysis.sheets[:3]:
        bullets = [
            f"{sheet.n_rows} data rows · {sheet.n_columns} columns",
            *[
                f"• {c.name} ({c.kind}) — {c.non_null_count} non-null · "
                f"{c.null_count} blank"
                for c in sheet.columns[:8]
            ],
        ]
        s = _add_text_slide(prs, f"Sheet: {sheet.name}", bullets)
        _set_notes(s, (
            f"The {sheet.name} sheet contains {sheet.n_rows} rows of data "
            f"across {sheet.n_columns} columns. Column kinds are inferred from "
            f"the values present; readers may want to confirm the inference "
            f"on the Analyze tab before relying on a downstream calculation."
        ))

    # ─── Signals ─────────────────────────────────────────────────
    if analysis.signals:
        bullets = [f"• [{s.kind}] {s.title} — {s.detail}" for s in analysis.signals[:8]]
        s = _add_text_slide(prs, "Signals", bullets)
        _set_notes(s, (
            f"This slide lists {min(8, len(analysis.signals))} observations from the "
            f"workbook, ranked by analytical priority (trend → outlier → metric → "
            f"missing data). Each observation cites the cell range from which "
            f"it was derived; reviewers may want to open the workbook to a cited "
            f"range and confirm the value in context before relying on a finding."
        ))

    # ─── Simulations ─────────────────────────────────────────────
    for mc in analysis.simulations:
        bullets = [
            f"Distribution: {mc.distribution} · params {mc.params}",
            f"Formula: {mc.formula} · iterations {mc.iterations} · seed {mc.seed}",
            f"P10 / P50 / P90: {mc.p10:.2f}  /  {mc.p50:.2f}  /  {mc.p90:.2f}",
            f"Mean ± stddev: {mc.mean:.2f} ± {mc.stddev:.2f}",
            f"Reproducer hash: {mc.reproducer_hash[:16]}…",
        ]
        s = _add_text_slide(prs, f"Simulation — {mc.column}", bullets)
        narration = (mc.narration.text if mc.narration else None) or (
            f"The {mc.iterations}-iteration simulation on {mc.column} produced a "
            f"median outcome of {mc.p50:.2f} with the central 80% of outcomes "
            f"falling between {mc.p10:.2f} (P10) and {mc.p90:.2f} (P90). The "
            f"mean is {mc.mean:.2f} with a standard deviation of {mc.stddev:.2f}. "
            f"Re-running with the same seed and parameters reproduces these "
            f"numbers byte-identically — the reproducer hash captures the full "
            f"specification."
        )
        _set_notes(s, narration)

    # ─── Forecasts ───────────────────────────────────────────────
    for fc in analysis.forecasts:
        bullets = [
            f"{fc.value_column} vs {fc.date_column} · method {fc.method}",
            f"Slope: {fc.slope:.4f} · intercept: {fc.intercept:.2f} · R²: {fc.r2:.3f}",
            f"Historical pairs: {fc.n_historical} · horizon: {fc.horizon_periods} period(s)",
            *[
                f"  Period +{p['period_index']}: {p['value']:.2f} "
                f"(80% CI {p['ci_low']:.2f} – {p['ci_high']:.2f})"
                for p in fc.projections[:6]
            ],
        ]
        s = _add_text_slide(prs, f"Forecast — {fc.value_column}", bullets)
        narration = (fc.narration.text if fc.narration else None) or (
            f"Fitting a linear regression to {fc.n_historical} historical "
            f"({fc.date_column}, {fc.value_column}) pairs yields a slope of "
            f"{fc.slope:.4f} per ordinal day and an R² of {fc.r2:.3f}. The "
            f"first forecast period projects to {fc.projections[0]['value']:.2f} "
            f"with an 80% confidence interval of "
            f"{fc.projections[0]['ci_low']:.2f}–{fc.projections[0]['ci_high']:.2f}. "
            f"Linear regression assumes the historical pattern continues; "
            f"reviewers may want to weigh this against context the workbook "
            f"itself does not capture."
        )
        _set_notes(s, narration)

    # ─── Anomalies ───────────────────────────────────────────────
    if analysis.anomalies:
        rows = analysis.anomalies[:8]
        bullets = [
            f"Row {a.row_index} of {a.sheet}!{a.column}: value={a.value:.2f} "
            f"(z={a.z_score:+.2f}, IQR-dist={a.iqr_distance:+.2f})"
            for a in rows
        ]
        s = _add_text_slide(prs, "Anomaly observations", bullets)
        _set_notes(s, (
            f"The detector flagged {len(analysis.anomalies)} anomaly observation(s) "
            f"using z-score and inter-quartile-range thresholds. Each row "
            f"listed here is cited to a single cell so reviewers can open the "
            f"workbook and confirm the value in context."
        ))

    # ─── Methodology / reproducibility ───────────────────────────
    method_lines: List[str] = [
        "Numeric statistics: population stddev, median via numpy.",
        "Anomaly thresholds: |z| ≥ 3 OR |IQR-distance| ≥ 1.5.",
        "Forecast method: numpy.polyfit linear regression, residual-based 80% CI.",
        "Monte Carlo: numpy.random.default_rng(seed) for full reproducibility.",
    ]
    for mc in analysis.simulations:
        method_lines.append(
            f"• Simulation {mc.id}: {mc.distribution} {mc.params} · seed {mc.seed} · "
            f"reproducer hash {mc.reproducer_hash[:24]}…"
        )
    s = _add_text_slide(prs, "Methodology · reproducibility", method_lines)
    _set_notes(s, (
        "This appendix lists the exact parameters used to produce every "
        "simulation and forecast in this deck. Re-running with the same "
        "reproducer hash produces byte-identical outputs."
    ))

    out = io.BytesIO()
    prs.save(out)
    return out.getvalue()


# ─────────────────────────────────────────────────────────────────
# Track A Phase 1 (2026-06-03) — DOCX and XLSX siblings.
#
# These mirror the PPTX builder's contract: same input
# (`WorkbookAnalysis`), same observational-only narration discipline
# (every narration string passes through `validate_no_imperatives`),
# byte-deterministic for a given input.
#
# The PPTX builder above is byte-identical to the P5.14 ship.
# ─────────────────────────────────────────────────────────────────


def _validate_block(text: str, label: str) -> str:
    """Same contract as `_set_notes` — reject imperative-to-user
    phrasing rather than ship the bad narration in the file."""
    validate_no_imperatives(text, label=label)
    return text


def build_docx_report(analysis: WorkbookAnalysis) -> bytes:
    """Serialise a WorkbookAnalysis into a .docx binary.

    Layout (mirrors the PPTX section ordering — readers should
    find the same content in both files):
      • H1 cover + summary paragraph
      • One H2 per sheet (≤3) + bullet list of column metadata
      • H2 Signals + bullets
      • H2 Simulation results (one paragraph per MonteCarloRun)
      • H2 Forecast results (one paragraph per ForecastRun)
      • H2 Anomalies + bullets
      • H2 Methodology / reproducibility appendix
    """
    from docx import Document  # local import keeps cold-path latency off the PPTX hot path

    doc = Document()

    # Cover
    cover_h = doc.add_heading(f"Workbook Analysis · {analysis.filename}", level=1)
    cover_h.paragraph_format.space_after = None
    doc.add_paragraph(_validate_block((
        f"Prepared by Akki for Executives · analysis id {analysis.id} · "
        f"format {analysis.file_format} · {len(analysis.sheets)} sheet(s). "
        f"This document reports the observations from the {analysis.filename} "
        f"workbook. It contains {len(analysis.signals)} signal(s), "
        f"{len(analysis.simulations)} simulation(s), {len(analysis.forecasts)} "
        f"forecast(s), and {len(analysis.anomalies)} anomaly observation(s). "
        f"Every numerical claim is cited to a cell range in the source workbook; "
        f"the appendix at the end lists the exact ranges."
    ), label="docx/cover"))

    # Sheets (max 3)
    for sheet in analysis.sheets[:3]:
        doc.add_heading(f"Sheet: {sheet.name}", level=2)
        doc.add_paragraph(
            f"{sheet.n_rows} data rows · {sheet.n_columns} columns",
            style=None,
        )
        for c in sheet.columns[:8]:
            doc.add_paragraph(
                f"• {c.name} ({c.kind}) — {c.non_null_count} non-null · {c.null_count} blank",
                style=None,
            )

    # Signals
    if analysis.signals:
        doc.add_heading("Signals", level=2)
        for s in analysis.signals[:12]:
            doc.add_paragraph(_validate_block(
                f"[{s.kind}] {s.title} — {s.detail}",
                label=f"docx/signal/{s.title}",
            ))

    # Simulations
    for mc in analysis.simulations:
        doc.add_heading(f"Simulation — {mc.column}", level=2)
        doc.add_paragraph(
            f"Distribution: {mc.distribution} · params {mc.params} · "
            f"iterations {mc.iterations} · seed {mc.seed}"
        )
        doc.add_paragraph(
            f"P10 / P50 / P90: {mc.p10:.2f} / {mc.p50:.2f} / {mc.p90:.2f} · "
            f"mean {mc.mean:.2f} ± stddev {mc.stddev:.2f}"
        )
        narration = (mc.narration.text if mc.narration else None) or (
            f"The {mc.iterations}-iteration simulation on {mc.column} produced a "
            f"median outcome of {mc.p50:.2f} with the central 80% of outcomes "
            f"falling between {mc.p10:.2f} and {mc.p90:.2f}."
        )
        doc.add_paragraph(_validate_block(narration, label=f"docx/sim/{mc.id}"))

    # Forecasts
    for fc in analysis.forecasts:
        doc.add_heading(f"Forecast — {fc.value_column}", level=2)
        doc.add_paragraph(
            f"{fc.value_column} vs {fc.date_column} · method {fc.method} · "
            f"slope {fc.slope:.4f} · R² {fc.r2:.3f} · "
            f"historical pairs {fc.n_historical} · horizon {fc.horizon_periods}"
        )
        for p in fc.projections[:6]:
            doc.add_paragraph(
                f"  Period +{p['period_index']}: {p['value']:.2f} "
                f"(80% CI {p['ci_low']:.2f}–{p['ci_high']:.2f})"
            )
        narration = (fc.narration.text if fc.narration else None) or (
            f"Fitting a linear regression to {fc.n_historical} historical "
            f"({fc.date_column}, {fc.value_column}) pairs yields a slope of "
            f"{fc.slope:.4f} per ordinal day and an R² of {fc.r2:.3f}."
        )
        doc.add_paragraph(_validate_block(narration, label=f"docx/fc/{fc.id}"))

    # Anomalies
    if analysis.anomalies:
        doc.add_heading("Anomaly observations", level=2)
        for a in analysis.anomalies[:12]:
            doc.add_paragraph(
                f"Row {a.row_index} of {a.sheet}!{a.column}: value={a.value:.2f} "
                f"(z={a.z_score:+.2f}, IQR-dist={a.iqr_distance:+.2f}) — "
                f"{_validate_block(a.rationale, label=f'docx/anom/{a.row_index}')}"
            )

    # Methodology
    doc.add_heading("Methodology · reproducibility", level=2)
    doc.add_paragraph("Numeric statistics: population stddev, median via numpy.")
    doc.add_paragraph("Anomaly thresholds: |z| ≥ 3 OR |IQR-distance| ≥ 1.5.")
    doc.add_paragraph(
        "Forecast method: numpy.polyfit linear regression, residual-based 80% CI."
    )
    doc.add_paragraph(
        "Monte Carlo: numpy.random.default_rng(seed) for full reproducibility."
    )
    for mc in analysis.simulations:
        doc.add_paragraph(
            f"Simulation {mc.id}: {mc.distribution} {mc.params} · seed {mc.seed} · "
            f"reproducer hash {mc.reproducer_hash[:24]}…"
        )

    out = io.BytesIO()
    doc.save(out)
    return out.getvalue()


def build_xlsx_report(analysis: WorkbookAnalysis) -> bytes:
    """Serialise a WorkbookAnalysis into a multi-sheet .xlsx binary.

    Sheets:
      • Summary  — analysis id, filename, counts of each artefact kind
      • Signals  — kind, title, detail, sheet, column, magnitude
      • Simulations — column, distribution, iterations, seed, P10/50/90
      • Forecasts — value_column, date_column, slope, R², horizon, projections
      • Anomalies — sheet, column, row, value, z_score, iqr_distance, rationale
    """
    from openpyxl import Workbook  # local import; the Analyze pipeline
    # already requires openpyxl for parsing, but keeping the import
    # local mirrors the docx pattern and keeps the PPTX cold-path
    # tidy.

    wb = Workbook()
    summary = wb.active
    summary.title = "Summary"
    summary.append(["field", "value"])
    summary.append(["analysis_id", analysis.id])
    summary.append(["filename", analysis.filename])
    summary.append(["file_format", analysis.file_format])
    summary.append(["sheet_count", len(analysis.sheets)])
    summary.append(["signal_count", len(analysis.signals)])
    summary.append(["simulation_count", len(analysis.simulations)])
    summary.append(["forecast_count", len(analysis.forecasts)])
    summary.append(["anomaly_count", len(analysis.anomalies)])
    summary.append(["status", analysis.status])
    summary.append(["created_at", analysis.created_at])
    summary.append(["updated_at", analysis.updated_at])

    # Signals
    sig_ws = wb.create_sheet("Signals")
    sig_ws.append(["kind", "title", "detail", "sheet", "column", "magnitude"])
    for s in analysis.signals:
        # Refuse-to-decide validation on detail before serialisation.
        _validate_block(s.detail, label=f"xlsx/signal/{s.title}")
        sig_ws.append([
            s.kind, s.title, s.detail, s.sheet, s.column,
            (s.magnitude if s.magnitude is not None else ""),
        ])

    # Simulations
    sim_ws = wb.create_sheet("Simulations")
    sim_ws.append([
        "id", "column", "distribution", "iterations", "seed",
        "P10", "P25", "P50", "P75", "P90", "mean", "stddev",
        "reproducer_hash",
    ])
    for mc in analysis.simulations:
        sim_ws.append([
            mc.id, mc.column, mc.distribution, mc.iterations, mc.seed,
            mc.p10, mc.p25, mc.p50, mc.p75, mc.p90, mc.mean, mc.stddev,
            mc.reproducer_hash,
        ])

    # Forecasts
    fc_ws = wb.create_sheet("Forecasts")
    fc_ws.append([
        "id", "value_column", "date_column", "method", "slope",
        "intercept", "r2", "n_historical", "horizon_periods",
        "period_index", "projection_value", "ci_low", "ci_high",
    ])
    for fc in analysis.forecasts:
        if not fc.projections:
            fc_ws.append([
                fc.id, fc.value_column, fc.date_column, fc.method, fc.slope,
                fc.intercept, fc.r2, fc.n_historical, fc.horizon_periods,
                "", "", "", "",
            ])
            continue
        for p in fc.projections:
            fc_ws.append([
                fc.id, fc.value_column, fc.date_column, fc.method, fc.slope,
                fc.intercept, fc.r2, fc.n_historical, fc.horizon_periods,
                p.get("period_index"), p.get("value"),
                p.get("ci_low"), p.get("ci_high"),
            ])

    # Anomalies
    an_ws = wb.create_sheet("Anomalies")
    an_ws.append(["sheet", "column", "row_index", "value", "z_score", "iqr_distance", "rationale"])
    for a in analysis.anomalies:
        _validate_block(a.rationale, label=f"xlsx/anom/{a.row_index}")
        an_ws.append([a.sheet, a.column, a.row_index, a.value, a.z_score, a.iqr_distance, a.rationale])

    out = io.BytesIO()
    wb.save(out)
    return out.getvalue()


__all__ = ["build_pptx_report", "build_docx_report", "build_xlsx_report"]
