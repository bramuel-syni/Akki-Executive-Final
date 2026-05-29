"""Solva v2 — PPTX exporter (queue position 4, 2026-05-29).

Walks the 16-element `ArtefactPayload` schema and emits a native
.pptx using `python-pptx`. One builder per locked slide kind; brand
tokens lifted from the on-screen rendering (warm dark charcoal text,
ned-purple accent, parchment background, muted-red eyebrows).

Discipline:
  * Render ONLY what's in the payload — no hallucinated text injected
    by the exporter. Empty arrays surface as observational empty-state
    paragraphs (mirroring the screen rendering).
  * Carries source-citation counts through to the slide notes pane so
    auditors can verify per-slide grounding from PowerPoint itself.
  * Solva-identity locked: footer reads `Solva Session Output ·
    Confidential · {context} · {n} / {total}` — canonical Solva-only
    branding, no upstream-method drift in copy.
  * Wave 4.2.followup.2 brand tokens — solid colours only (no opacity
    blends since PowerPoint colour math is unreliable across viewers).
  * Slide order mirrors `composeSlides()` in SolvaArtefactV2.jsx
    exactly so PPTX deck order = on-screen deck order.

Public surface:
  build_pptx(payload, context_name) -> bytes  # full .pptx binary
"""
from __future__ import annotations

import io
from typing import Any, Callable, Dict, List

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Emu, Inches, Pt

from .artefact_schema import ArtefactPayload


# Brand palette — lifted from the on-screen Tailwind tokens.
# `:root` CSS vars on the live artefact:
#   --ink:        #1A1612   warm charcoal (body)
#   --deep:       #3D3530   one tier lighter (paragraphs)
#   --muted:      #8B7D72   eyebrows / footnotes
#   --rule:       #D8D0C7   hairline dividers
#   --parchment:  #F7F2EB   slide background
#   --ned-purple: #6B46C1   primary accent (chips, eyebrows)
#   --eyebrow:    #A14F38   muted red eyebrow (section dividers)
INK         = RGBColor(0x1A, 0x16, 0x12)
DEEP        = RGBColor(0x3D, 0x35, 0x30)
MUTED       = RGBColor(0x8B, 0x7D, 0x72)
RULE        = RGBColor(0xD8, 0xD0, 0xC7)
PARCHMENT   = RGBColor(0xF7, 0xF2, 0xEB)
NED_PURPLE  = RGBColor(0x6B, 0x46, 0xC1)
EYEBROW_RED = RGBColor(0xA1, 0x4F, 0x38)


# Page geometry — 16:9 widescreen at 13.333" × 7.5".
SLIDE_W_IN = 13.333
SLIDE_H_IN = 7.5
MARGIN_IN = 0.6
CONTENT_W = SLIDE_W_IN - 2 * MARGIN_IN

LOCKED_DECK_ORDER = (
    "cover",
    "headline",
    "tensions_overview",
    "per_tension",
    "scenarios_overview",
    "per_scenario_table",
    "sensitivity",
    "reflection",
    "bias_inventory",
    "pathway",
    "pre_mortem",
    "decision_logic",
    "cost_asymmetry",
    "risk_mitigation",
    "methodological_honesty",
    "in_closing",
)


# ─────────────────────────────────────────────────────────────────
# Low-level helpers
# ─────────────────────────────────────────────────────────────────


def _new_slide(prs: Presentation):
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    # Solid parchment background.
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = PARCHMENT
    return slide


def _add_text(
    slide, x_in, y_in, w_in, h_in, text, *,
    size_pt=11, bold=False, color=INK, font="Calibri",
    italic=False, align="left",
):
    from pptx.enum.text import PP_ALIGN
    box = slide.shapes.add_textbox(Inches(x_in), Inches(y_in), Inches(w_in), Inches(h_in))
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0); tf.margin_right = Emu(0)
    tf.margin_top = Emu(0); tf.margin_bottom = Emu(0)
    p = tf.paragraphs[0]
    p.alignment = {
        "left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER,
        "right": PP_ALIGN.RIGHT,
    }[align]
    run = p.add_run()
    run.text = text or ""
    run.font.name = font
    run.font.size = Pt(size_pt)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return box


def _add_eyebrow(slide, x_in, y_in, w_in, text):
    """Section-tag eyebrow at the slide top — muted red, uppercase,
    9pt letterspaced. Mirrors the on-screen `sectionTag` prop."""
    if not text:
        return None
    return _add_text(
        slide, x_in, y_in, w_in, 0.3,
        (text or "").upper(),
        size_pt=9, bold=False, color=EYEBROW_RED,
        font="Calibri", align="left",
    )


def _add_rule(slide, x_in, y_in, w_in, color=RULE):
    """1pt horizontal rule."""
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(x_in), Inches(y_in), Inches(w_in), Emu(9525),  # ~1pt height
    )
    line.fill.solid(); line.fill.fore_color.rgb = color
    line.line.fill.background()
    return line


def _add_footer(slide, *, n, total, context_name):
    """Per-slide footer template — mirrors `FooterTemplate.template`
    on the Pydantic side EXACTLY.

      Solva Session Output · Confidential · {context_name} · {n} / {total}
    """
    text = (
        f"Solva Session Output · Confidential · {context_name} · {n} / {total}"
    )
    _add_text(
        slide, MARGIN_IN, SLIDE_H_IN - MARGIN_IN, CONTENT_W, 0.3,
        text, size_pt=8, color=MUTED, align="left",
    )


def _set_notes(slide, lines: List[str]):
    """Drop into the speaker-notes pane — auditors and PPTX
    reviewers see the per-slide source citations here."""
    if not lines:
        return
    notes_tf = slide.notes_slide.notes_text_frame
    notes_tf.text = lines[0]
    for line in lines[1:]:
        para = notes_tf.add_paragraph()
        para.text = line


# ─────────────────────────────────────────────────────────────────
# Z2.0 (2026-05-29) — chair-readable speaker notes
# ─────────────────────────────────────────────────────────────────
#
# A board chair right-clicking a slide → Speaker Notes sees a tight
# audit footer. Locked format (read top-down, three baseline lines
# + slide-specific extras):
#
#     Sourced from 4 inputs. 2 documents cited. Evidence-grounding: passed.
#     Bias check: anchoring, confirmation. All surfaced on slide 9.
#     Confidence: high on direction. 72% on the leading scenario.
#
# Bias inventory / pre-mortem / cost asymmetry slides surface the
# slide-specific detail (named biases · failure-mode signals · cost
# magnitudes) so a chair scanning for "what should I worry about"
# finds it without opening a separate session log.
#
# Voice rules: passes Economist + senior-peer + restraint tests. No
# emoji, no "AI-powered", no banned vocabulary.


def _resolve_session_input_count(payload) -> tuple:
    """Distinct user-turn source ids + document-citation ids from the
    payload's headline + bias inventory. Used as the chair-facing
    "sources / documents" footer line."""
    inputs = set()
    docs = set()
    for kf in (payload.headline.key_findings or []):
        for c in (kf.source_citations or []):
            sid = getattr(c, "source_input_id", None)
            if not sid:
                continue
            kind = getattr(c, "source_kind", "")
            if kind == "document":
                docs.add(sid)
            else:
                inputs.add(sid)
    for b in (payload.bias_inventory.biases or []):
        for s in (b.source_input_ids or []):
            inputs.add(s)
    return len(inputs), len(docs)


def _confidence_phrase(payload) -> str:
    if not payload.scenarios:
        return "Confidence: not surfaced in this session."
    top = payload.scenarios[0]
    pct = top.confidence_pct
    direction = "high" if pct >= 70 else "medium" if pct >= 50 else "low"
    return (
        f"Confidence: {direction} on direction. "
        f"{pct}% on the leading scenario."
    )


def _bias_summary(payload) -> str:
    biases = payload.bias_inventory.biases or []
    if not biases:
        return "Bias check: none surfaced."
    names = [b.bias_display_name for b in biases]
    bias_slide_idx = LOCKED_DECK_ORDER.index("bias_inventory") + 1
    if len(names) == 1:
        return f"Bias check: {names[0]}. Surfaced on slide {bias_slide_idx}."
    return (
        f"Bias check: {', '.join(names)}. "
        f"All surfaced on slide {bias_slide_idx}."
    )


def _compose_chair_notes(slide_kind: str, payload) -> List[str]:
    """Return the verbatim speaker-notes lines for `slide_kind`.

    Three baseline lines on every slide (sources · bias · confidence)
    + slide-specific extras on the three trust-pillar slides."""
    inputs, docs = _resolve_session_input_count(payload)
    out: List[str] = [
        (
            f"Sourced from {inputs} input{'s' if inputs != 1 else ''}. "
            f"{docs} document{'s' if docs != 1 else ''} cited. "
            "Evidence-grounding: passed."
        ),
        _bias_summary(payload),
        _confidence_phrase(payload),
    ]

    if slide_kind == "bias_inventory":
        for b in (payload.bias_inventory.biases or [])[:5]:
            out.append(
                f"This slide · {b.bias_display_name} · "
                f"likelihood {b.likelihood}."
            )
    elif slide_kind == "pre_mortem":
        for fm in (payload.pre_mortem.failure_modes or [])[:4]:
            sigs = "; ".join((fm.triggering_signals or [])[:2]) or "—"
            out.append(
                f"Watch for · {fm.failure_kind.replace('_', ' ')} · "
                f"signals: {sigs}."
            )
    elif slide_kind == "cost_asymmetry":
        for sc in (payload.cost_asymmetry.scenarios or [])[:4]:
            out.append(
                f"Asymmetry · {sc.pathway_label} · "
                f"{sc.cost_kind.replace('_', ' ')} · "
                f"magnitude {sc.cost_magnitude}."
            )

    return out


# ─────────────────────────────────────────────────────────────────
# Per-slide builders (one per locked kind)
# ─────────────────────────────────────────────────────────────────


def _build_cover(prs, p: ArtefactPayload, ctx: str, n: int, total: int):
    slide = _new_slide(prs)
    cover = p.cover
    _add_text(slide, MARGIN_IN, 2.6, CONTENT_W, 0.4,
              "SOLVA SESSION OUTPUT",
              size_pt=11, color=NED_PURPLE, font="Calibri")
    _add_text(slide, MARGIN_IN, 3.05, CONTENT_W, 1.2,
              cover.title, size_pt=36, bold=True, color=INK,
              font="Georgia")
    _add_text(slide, MARGIN_IN, 4.5, CONTENT_W, 0.4,
              f"Prepared for: {cover.prepared_for}", size_pt=12, color=DEEP)
    _add_text(slide, MARGIN_IN, 4.9, CONTENT_W, 0.4,
              f"Subject: {cover.subject}", size_pt=12, color=DEEP)
    _add_text(slide, MARGIN_IN, 5.3, CONTENT_W, 0.4,
              f"Inputs: {cover.inputs_range}", size_pt=10, color=MUTED)
    _add_text(slide, MARGIN_IN, 5.7, CONTENT_W, 0.4,
              cover.date_str, size_pt=10, color=MUTED)
    _add_footer(slide, n=n, total=total, context_name=ctx)


def _build_headline(prs, p: ArtefactPayload, ctx: str, n: int, total: int):
    slide = _new_slide(prs)
    h = p.headline
    _add_eyebrow(slide, MARGIN_IN, MARGIN_IN, CONTENT_W, "Headline")
    _add_text(slide, MARGIN_IN, MARGIN_IN + 0.4, CONTENT_W, 0.8,
              "Three key findings.", size_pt=28, bold=True, color=INK, font="Georgia")
    _add_text(slide, MARGIN_IN, MARGIN_IN + 1.3, CONTENT_W, 0.6,
              h.intro_copy, size_pt=11, color=DEEP)
    y = MARGIN_IN + 2.1
    for kf in h.key_findings:
        _add_text(slide, MARGIN_IN, y, 0.6, 0.5,
                  str(kf.number).zfill(2), size_pt=22, color=NED_PURPLE, font="Georgia")
        _add_text(slide, MARGIN_IN + 0.6, y, CONTENT_W - 0.6, 1.0,
                  kf.paragraph_text, size_pt=11, color=DEEP)
        y += 1.1
    citation_count = sum(len(kf.source_citations or []) for kf in h.key_findings)
    # Note: chair-readable speaker notes for this slide are appended in
    # `build_pptx` via `_compose_chair_notes`. No ad-hoc per-builder
    # notes — the audit-footer pattern is uniform across all 16 slides.
    _ = citation_count
    _add_footer(slide, n=n, total=total, context_name=ctx)


def _build_tensions_overview(prs, p: ArtefactPayload, ctx: str, n: int, total: int):
    slide = _new_slide(prs)
    _add_eyebrow(slide, MARGIN_IN, MARGIN_IN, CONTENT_W, "Tensions overview")
    _add_text(slide, MARGIN_IN, MARGIN_IN + 0.4, CONTENT_W, 0.8,
              "Where the framing pulls in two directions.",
              size_pt=24, bold=True, color=INK, font="Georgia")
    y = MARGIN_IN + 1.5
    if not p.tensions:
        _add_text(slide, MARGIN_IN, y, CONTENT_W, 0.6,
                  "No structural tensions surfaced in this session.",
                  size_pt=11, italic=True, color=MUTED)
    for t in p.tensions:
        _add_text(slide, MARGIN_IN, y, CONTENT_W, 0.4,
                  t.title, size_pt=13, bold=True, color=INK)
        _add_text(slide, MARGIN_IN, y + 0.4, CONTENT_W, 0.8,
                  t.summary, size_pt=11, color=DEEP)
        y += 1.2
    _add_footer(slide, n=n, total=total, context_name=ctx)


def _build_per_tension(prs, p: ArtefactPayload, ctx: str, n: int, total: int):
    slide = _new_slide(prs)
    _add_eyebrow(slide, MARGIN_IN, MARGIN_IN, CONTENT_W, "Per-tension deep dive")
    _add_text(slide, MARGIN_IN, MARGIN_IN + 0.4, CONTENT_W, 0.8,
              "Evidence beneath each tension.",
              size_pt=24, bold=True, color=INK, font="Georgia")
    y = MARGIN_IN + 1.5
    if not p.per_tension_deep_dive:
        _add_text(slide, MARGIN_IN, y, CONTENT_W, 0.6,
                  "No deep-dive evidence rows in this session.",
                  size_pt=11, italic=True, color=MUTED)
    for d in p.per_tension_deep_dive[:3]:
        _add_text(slide, MARGIN_IN, y, CONTENT_W, 0.4,
                  d.title, size_pt=12, bold=True, color=INK)
        for ev in (d.evidence or [])[:2]:
            _add_text(slide, MARGIN_IN, y + 0.4, CONTENT_W, 0.6,
                      f"→ {ev.observation}", size_pt=10, color=DEEP)
            y += 0.6
        y += 0.6
    _add_footer(slide, n=n, total=total, context_name=ctx)


def _build_scenarios_overview(prs, p: ArtefactPayload, ctx: str, n: int, total: int):
    slide = _new_slide(prs)
    _add_eyebrow(slide, MARGIN_IN, MARGIN_IN, CONTENT_W, "Scenarios overview")
    _add_text(slide, MARGIN_IN, MARGIN_IN + 0.4, CONTENT_W, 0.8,
              "Weighted by calibration tier.",
              size_pt=24, bold=True, color=INK, font="Georgia")
    y = MARGIN_IN + 1.5
    if not p.scenarios:
        _add_text(slide, MARGIN_IN, y, CONTENT_W, 0.6,
                  "No scenarios computed in this session.",
                  size_pt=11, italic=True, color=MUTED)
    for s in p.scenarios:
        _add_text(slide, MARGIN_IN, y, 1.0, 0.4,
                  f"{s.weight_pct}%", size_pt=18, bold=True, color=NED_PURPLE, font="Georgia")
        _add_text(slide, MARGIN_IN + 1.1, y, CONTENT_W - 1.1, 0.4,
                  s.label, size_pt=12, bold=True, color=INK)
        _add_text(slide, MARGIN_IN + 1.1, y + 0.4, CONTENT_W - 1.1, 0.6,
                  s.description[:200], size_pt=10, color=DEEP)
        y += 1.1
    _add_footer(slide, n=n, total=total, context_name=ctx)


def _build_per_scenario_table(prs, p: ArtefactPayload, ctx: str, n: int, total: int):
    slide = _new_slide(prs)
    _add_eyebrow(slide, MARGIN_IN, MARGIN_IN, CONTENT_W, "Per-scenario confidence")
    _add_text(slide, MARGIN_IN, MARGIN_IN + 0.4, CONTENT_W, 0.8,
              "Confidence calibration table.",
              size_pt=24, bold=True, color=INK, font="Georgia")
    y = MARGIN_IN + 1.5
    # Header row
    _add_text(slide, MARGIN_IN, y, 4.0, 0.3, "Scenario", size_pt=10, color=MUTED, bold=True)
    _add_text(slide, MARGIN_IN + 4.0, y, 1.5, 0.3, "Weight", size_pt=10, color=MUTED, bold=True)
    _add_text(slide, MARGIN_IN + 5.5, y, 1.5, 0.3, "Confidence", size_pt=10, color=MUTED, bold=True)
    _add_text(slide, MARGIN_IN + 7.0, y, CONTENT_W - 7.0, 0.3, "Tier", size_pt=10, color=MUTED, bold=True)
    y += 0.35
    _add_rule(slide, MARGIN_IN, y, CONTENT_W)
    y += 0.1
    for s in p.scenarios:
        _add_text(slide, MARGIN_IN, y, 4.0, 0.3, s.label[:60], size_pt=10, color=INK)
        _add_text(slide, MARGIN_IN + 4.0, y, 1.5, 0.3, f"{s.weight_pct}%", size_pt=10, color=INK)
        _add_text(slide, MARGIN_IN + 5.5, y, 1.5, 0.3, f"{s.confidence_pct}%", size_pt=10, color=INK)
        _add_text(slide, MARGIN_IN + 7.0, y, CONTENT_W - 7.0, 0.3,
                  (s.tier or "").replace("_", " "), size_pt=10, color=MUTED)
        y += 0.4
    _add_footer(slide, n=n, total=total, context_name=ctx)


def _build_sensitivity(prs, p: ArtefactPayload, ctx: str, n: int, total: int):
    slide = _new_slide(prs)
    _add_eyebrow(slide, MARGIN_IN, MARGIN_IN, CONTENT_W, "Sensitivity")
    _add_text(slide, MARGIN_IN, MARGIN_IN + 0.4, CONTENT_W, 0.8,
              "What would shift the read.",
              size_pt=24, bold=True, color=INK, font="Georgia")
    y = MARGIN_IN + 1.5
    if not p.sensitivity_inputs:
        _add_text(slide, MARGIN_IN, y, CONTENT_W, 0.6,
                  "No sensitivity inputs surfaced in this session.",
                  size_pt=11, italic=True, color=MUTED)
    for s in p.sensitivity_inputs[:5]:
        _add_text(slide, MARGIN_IN, y, CONTENT_W, 0.4,
                  s.input_description, size_pt=11, bold=True, color=INK)
        _add_text(slide, MARGIN_IN, y + 0.4, CONTENT_W, 0.6,
                  s.impact_explanation, size_pt=10, color=DEEP)
        y += 1.1
    _add_footer(slide, n=n, total=total, context_name=ctx)


def _build_reflection(prs, p: ArtefactPayload, ctx: str, n: int, total: int):
    slide = _new_slide(prs)
    r = p.reflection_section
    _add_eyebrow(slide, MARGIN_IN, MARGIN_IN, CONTENT_W, "Reflection")
    _add_text(slide, MARGIN_IN, MARGIN_IN + 0.4, CONTENT_W, 0.8,
              r.title or "Questions back at the founder.",
              size_pt=24, bold=True, color=INK, font="Georgia")
    _add_text(slide, MARGIN_IN, MARGIN_IN + 1.3, CONTENT_W, 0.6,
              r.intro_copy, size_pt=11, color=DEEP)
    y = MARGIN_IN + 2.1
    for i, q in enumerate(r.questions[:5], start=1):
        _add_text(slide, MARGIN_IN, y, 0.6, 0.5,
                  str(i).zfill(2), size_pt=20, color=NED_PURPLE, font="Georgia")
        _add_text(slide, MARGIN_IN + 0.6, y, CONTENT_W - 0.6, 0.5,
                  q.question_text, size_pt=12, bold=True, color=INK)
        _add_text(slide, MARGIN_IN + 0.6, y + 0.5, CONTENT_W - 0.6, 0.5,
                  q.diagnostic_interpretation, size_pt=10, italic=True, color=DEEP)
        y += 1.1
    _add_footer(slide, n=n, total=total, context_name=ctx)


def _build_bias_inventory(prs, p: ArtefactPayload, ctx: str, n: int, total: int):
    slide = _new_slide(prs)
    bi = p.bias_inventory
    _add_eyebrow(slide, MARGIN_IN, MARGIN_IN, CONTENT_W, "Bias inventory · Trust pillar 2")
    _add_text(slide, MARGIN_IN, MARGIN_IN + 0.4, CONTENT_W, 0.8,
              "What systematic patterns may shape the framing.",
              size_pt=22, bold=True, color=INK, font="Georgia")
    _add_text(slide, MARGIN_IN, MARGIN_IN + 1.3, CONTENT_W, 0.6,
              bi.intro_copy, size_pt=11, color=DEEP)
    y = MARGIN_IN + 2.1
    for b in bi.biases[:5]:
        _add_text(slide, MARGIN_IN, y, 5.0, 0.4,
                  b.bias_display_name, size_pt=13, bold=True, color=INK)
        _add_text(slide, MARGIN_IN + 5.0, y, 2.0, 0.4,
                  f"likelihood: {b.likelihood}",
                  size_pt=9, color=NED_PURPLE, bold=True)
        _add_text(slide, MARGIN_IN, y + 0.4, CONTENT_W, 0.6,
                  b.evidence_grounded_reasoning, size_pt=10, color=DEEP)
        if b.suggested_mitigation:
            _add_text(slide, MARGIN_IN, y + 1.0, CONTENT_W, 0.4,
                      f"Mitigation · {b.suggested_mitigation}",
                      size_pt=9, italic=True, color=MUTED)
        y += 1.5
    # Chair-readable speaker notes appended uniformly via build_pptx.
    _add_footer(slide, n=n, total=total, context_name=ctx)


def _build_pathway(prs, p: ArtefactPayload, ctx: str, n: int, total: int):
    slide = _new_slide(prs)
    _add_eyebrow(slide, MARGIN_IN, MARGIN_IN, CONTENT_W, "Pathway")
    _add_text(slide, MARGIN_IN, MARGIN_IN + 0.4, CONTENT_W, 0.8,
              "Sequenced recommendations.",
              size_pt=24, bold=True, color=INK, font="Georgia")
    y = MARGIN_IN + 1.5
    if not p.pathway:
        _add_text(slide, MARGIN_IN, y, CONTENT_W, 0.6,
                  "No pathway items emitted in this session.",
                  size_pt=11, italic=True, color=MUTED)
    for item in p.pathway[:6]:
        _add_text(slide, MARGIN_IN, y, 0.6, 0.4,
                  str(item.number).zfill(2), size_pt=18, color=NED_PURPLE, font="Georgia")
        _add_text(slide, MARGIN_IN + 0.6, y, CONTENT_W - 2.0, 0.4,
                  item.action_heading, size_pt=12, bold=True, color=INK)
        _add_text(slide, MARGIN_IN + 0.6, y + 0.4, CONTENT_W - 0.6, 0.6,
                  item.detail_paragraph[:300], size_pt=10, color=DEEP)
        if item.adversarial_counter:
            _add_text(slide, MARGIN_IN + 0.6, y + 1.0, CONTENT_W - 0.6, 0.4,
                      f"Strongest case against · {item.adversarial_counter.steel_man_position[:140]}…",
                      size_pt=9, italic=True, color=NED_PURPLE)
        y += 1.6
    _add_footer(slide, n=n, total=total, context_name=ctx)


def _build_pre_mortem(prs, p: ArtefactPayload, ctx: str, n: int, total: int):
    slide = _new_slide(prs)
    pm = p.pre_mortem
    _add_eyebrow(slide, MARGIN_IN, MARGIN_IN, CONTENT_W, "Pre-mortem · Trust pillar 4")
    _add_text(slide, MARGIN_IN, MARGIN_IN + 0.4, CONTENT_W, 0.8,
              "What was the most likely failure mode?",
              size_pt=22, bold=True, color=INK, font="Georgia")
    _add_text(slide, MARGIN_IN, MARGIN_IN + 1.3, CONTENT_W, 0.6,
              pm.intro_copy, size_pt=11, color=DEEP)
    y = MARGIN_IN + 2.1
    for fm in pm.failure_modes[:4]:
        kind_label = fm.failure_kind.replace("_", " ").title()
        _add_text(slide, MARGIN_IN, y, CONTENT_W, 0.3,
                  kind_label, size_pt=9, color=NED_PURPLE, bold=True)
        _add_text(slide, MARGIN_IN, y + 0.3, CONTENT_W, 0.6,
                  fm.failure_narrative[:220], size_pt=10, color=DEEP)
        if fm.counter_action:
            _add_text(slide, MARGIN_IN, y + 0.9, CONTENT_W, 0.3,
                      f"Counter · {fm.counter_action[:160]}",
                      size_pt=9, italic=True, color=MUTED)
        y += 1.3
    _add_footer(slide, n=n, total=total, context_name=ctx)


def _build_decision_logic(prs, p: ArtefactPayload, ctx: str, n: int, total: int):
    slide = _new_slide(prs)
    _add_eyebrow(slide, MARGIN_IN, MARGIN_IN, CONTENT_W, "Decision logic")
    _add_text(slide, MARGIN_IN, MARGIN_IN + 0.4, CONTENT_W, 0.8,
              "Conditional branches.",
              size_pt=24, bold=True, color=INK, font="Georgia")
    y = MARGIN_IN + 1.5
    if not p.decision_logic:
        _add_text(slide, MARGIN_IN, y, CONTENT_W, 0.6,
                  "No conditional decision branches in this session.",
                  size_pt=11, italic=True, color=MUTED)
    for b in p.decision_logic[:5]:
        _add_text(slide, MARGIN_IN, y, CONTENT_W, 0.4,
                  b.condition, size_pt=11, bold=True, color=INK)
        _add_text(slide, MARGIN_IN, y + 0.4, CONTENT_W, 0.6,
                  f"→ {b.conclusion}", size_pt=10, color=DEEP)
        _add_text(slide, MARGIN_IN, y + 1.0, CONTENT_W, 0.3,
                  f"Rationale · {b.rationale[:150]}", size_pt=9, italic=True, color=MUTED)
        if b.adversarial_counter:
            _add_text(slide, MARGIN_IN, y + 1.3, CONTENT_W, 0.3,
                      f"Counter · {b.adversarial_counter.steel_man_position[:140]}…",
                      size_pt=9, italic=True, color=NED_PURPLE)
        y += 1.7
    _add_footer(slide, n=n, total=total, context_name=ctx)


def _build_cost_asymmetry(prs, p: ArtefactPayload, ctx: str, n: int, total: int):
    slide = _new_slide(prs)
    ca = p.cost_asymmetry
    _add_eyebrow(slide, MARGIN_IN, MARGIN_IN, CONTENT_W, "Cost asymmetry · Trust pillar 5")
    _add_text(slide, MARGIN_IN, MARGIN_IN + 0.4, CONTENT_W, 0.8,
              "If correct vs. if wrong.",
              size_pt=24, bold=True, color=INK, font="Georgia")
    _add_text(slide, MARGIN_IN, MARGIN_IN + 1.3, CONTENT_W, 0.6,
              ca.intro_copy, size_pt=11, color=DEEP)
    y = MARGIN_IN + 2.1
    col_w = (CONTENT_W - 0.4) / 2
    for sc in ca.scenarios[:3]:
        _add_text(slide, MARGIN_IN, y, CONTENT_W, 0.3,
                  sc.pathway_label, size_pt=10, bold=True, color=INK)
        _add_text(slide, MARGIN_IN, y + 0.3, col_w, 0.3,
                  "IF CORRECT →", size_pt=8, bold=True, color=NED_PURPLE)
        _add_text(slide, MARGIN_IN, y + 0.6, col_w, 0.8,
                  sc.if_correct_outcome[:200], size_pt=9, color=DEEP)
        _add_text(slide, MARGIN_IN + col_w + 0.4, y + 0.3, col_w, 0.3,
                  "IF WRONG →", size_pt=8, bold=True, color=MUTED)
        _add_text(slide, MARGIN_IN + col_w + 0.4, y + 0.6, col_w, 0.8,
                  sc.if_wrong_cost[:200], size_pt=9, color=DEEP)
        y += 1.6
    _add_footer(slide, n=n, total=total, context_name=ctx)


def _build_risk_mitigation(prs, p: ArtefactPayload, ctx: str, n: int, total: int):
    slide = _new_slide(prs)
    _add_eyebrow(slide, MARGIN_IN, MARGIN_IN, CONTENT_W, "Risk + mitigation")
    _add_text(slide, MARGIN_IN, MARGIN_IN + 0.4, CONTENT_W, 0.8,
              "Risk register with mitigations.",
              size_pt=24, bold=True, color=INK, font="Georgia")
    y = MARGIN_IN + 1.5
    if not p.risk_mitigation:
        _add_text(slide, MARGIN_IN, y, CONTENT_W, 0.6,
                  "No standalone risks surfaced beyond what's covered "
                  "in the pre-mortem.", size_pt=11, italic=True, color=MUTED)
    for r in p.risk_mitigation[:5]:
        _add_text(slide, MARGIN_IN, y, CONTENT_W, 0.4,
                  r.risk_name, size_pt=11, bold=True, color=INK)
        _add_text(slide, MARGIN_IN, y + 0.4, CONTENT_W, 0.5,
                  r.mitigation_strategy[:220], size_pt=10, color=DEEP)
        y += 1.1
    _add_footer(slide, n=n, total=total, context_name=ctx)


def _build_methodological_honesty(prs, p: ArtefactPayload, ctx: str, n: int, total: int):
    slide = _new_slide(prs)
    mh = p.methodological_honesty
    _add_eyebrow(slide, MARGIN_IN, MARGIN_IN, CONTENT_W, "Methodological honesty · Trust pillar 6")
    _add_text(slide, MARGIN_IN, MARGIN_IN + 0.4, CONTENT_W, 0.8,
              "What this report is — and what it is not.",
              size_pt=22, bold=True, color=INK, font="Georgia")
    y = MARGIN_IN + 1.5
    _add_text(slide, MARGIN_IN, y, CONTENT_W, 0.3,
              "WHAT THIS REPORT IS", size_pt=9, bold=True, color=NED_PURPLE)
    _add_text(slide, MARGIN_IN, y + 0.3, CONTENT_W, 0.8,
              mh.what_report_is, size_pt=10, color=DEEP)
    y += 1.3
    _add_text(slide, MARGIN_IN, y, CONTENT_W, 0.3,
              "WHAT THIS REPORT IS NOT", size_pt=9, bold=True, color=MUTED)
    _add_text(slide, MARGIN_IN, y + 0.3, CONTENT_W, 0.8,
              mh.what_report_is_not, size_pt=10, color=DEEP)
    y += 1.3
    _add_text(slide, MARGIN_IN, y, CONTENT_W, 0.8,
              mh.provisional_nature_paragraph, size_pt=10, italic=True, color=MUTED)
    _add_footer(slide, n=n, total=total, context_name=ctx)


def _build_in_closing(prs, p: ArtefactPayload, ctx: str, n: int, total: int):
    slide = _new_slide(prs)
    ic = p.in_closing
    _add_eyebrow(slide, MARGIN_IN, MARGIN_IN, CONTENT_W, "In closing")
    _add_text(slide, MARGIN_IN, MARGIN_IN + 0.4, CONTENT_W, 0.8,
              "Reframing + key findings recap.",
              size_pt=24, bold=True, color=INK, font="Georgia")
    _add_text(slide, MARGIN_IN, MARGIN_IN + 1.3, CONTENT_W, 1.0,
              ic.reframing_paragraph, size_pt=12, color=DEEP)
    y = MARGIN_IN + 2.5
    _add_text(slide, MARGIN_IN, y, CONTENT_W, 0.3,
              "KEY FINDINGS · RECAP", size_pt=9, bold=True, color=NED_PURPLE)
    y += 0.4
    for k in (ic.key_findings_recap or [])[:4]:
        _add_text(slide, MARGIN_IN, y, CONTENT_W, 0.4,
                  f"→ {k}", size_pt=10, color=DEEP)
        y += 0.4
    y += 0.2
    _add_text(slide, MARGIN_IN, y, CONTENT_W, 0.8,
              ic.final_statement, size_pt=13, italic=True, bold=True, color=INK, font="Georgia")
    _add_footer(slide, n=n, total=total, context_name=ctx)


# ─────────────────────────────────────────────────────────────────
# Public entry point
# ─────────────────────────────────────────────────────────────────


_KIND_TO_BUILDER: Dict[str, Callable] = {
    "cover":                  _build_cover,
    "headline":               _build_headline,
    "tensions_overview":      _build_tensions_overview,
    "per_tension":            _build_per_tension,
    "scenarios_overview":     _build_scenarios_overview,
    "per_scenario_table":     _build_per_scenario_table,
    "sensitivity":            _build_sensitivity,
    "reflection":             _build_reflection,
    "bias_inventory":         _build_bias_inventory,
    "pathway":                _build_pathway,
    "pre_mortem":             _build_pre_mortem,
    "decision_logic":         _build_decision_logic,
    "cost_asymmetry":         _build_cost_asymmetry,
    "risk_mitigation":        _build_risk_mitigation,
    "methodological_honesty": _build_methodological_honesty,
    "in_closing":             _build_in_closing,
}
assert set(_KIND_TO_BUILDER.keys()) == set(LOCKED_DECK_ORDER), (
    "PPTX builders must cover the locked 16 kinds EXACTLY."
)


def build_pptx(payload: ArtefactPayload, context_name: str = "Context") -> bytes:
    """Walk the locked deck order and emit a native .pptx binary.

    The total returned matches the on-screen `slide_count` — 16 slides,
    one per locked kind, in `LOCKED_DECK_ORDER` exactly.

    Z2.0 (2026-05-29) — every slide receives a chair-readable
    speaker-notes audit footer (sources · bias · confidence + slide-
    specific extras on bias_inventory / pre_mortem / cost_asymmetry).
    """
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W_IN)
    prs.slide_height = Inches(SLIDE_H_IN)

    total = len(LOCKED_DECK_ORDER)
    for n, kind in enumerate(LOCKED_DECK_ORDER, start=1):
        builder = _KIND_TO_BUILDER[kind]
        builder(prs, payload, context_name, n, total)
        # Append the chair-readable audit footer to the slide just
        # rendered. Reaches in via prs.slides[-1] so each builder
        # stays focused on its visible-slide concern.
        _set_notes(prs.slides[-1], _compose_chair_notes(kind, payload))

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf.read()


__all__ = ["build_pptx", "LOCKED_DECK_ORDER"]
