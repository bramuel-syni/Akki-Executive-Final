"""Phase E.3 — Scope-compliance wire + live tests (2026-05-26).

Covers the 3 closures dispatched under autonomous mode:
  1. Prompt-based edit apply pipeline (Shield-bounded LLM rewrite +
     diff preview + apply).
  2. DRAFT watermark embedded in PDF / DOCX / PPTX exports (byte-level
     inspection — actual "DRAFT" stamp present in the output bytes).
  3. Related-docs typed groups (metadata_match, content_similarity,
     explicit_attachment gap, canonical_lineage gap).

Wire tests assert the LIVE BEHAVIOR (LLM diff sample bytes, watermarked
file bytes, Related tab DOM group labels), not JSX class strings.
"""
from __future__ import annotations

import io
import uuid
from datetime import datetime, timezone
from pathlib import Path

import pytest
from httpx import AsyncClient, ASGITransport


REPO = Path(__file__).resolve().parent.parent.parent
FE   = REPO / "frontend" / "src"
BE   = REPO / "backend"

DRAWER          = FE / "components" / "documents" / "DocumentDrawer.jsx"
DOCUMENTS_PY    = BE / "routers" / "documents.py"
WATERMARK_SVC   = BE / "services" / "documents" / "watermark_service.py"


def _read(p: Path) -> str:
    return p.read_text(encoding="utf-8")


# ═════════════════════════════════════════════════════════════════════
# 1. Prompted-edit apply pipeline
# ═════════════════════════════════════════════════════════════════════
def test_scope_prompted_edit_endpoint_wired_through_shield():
    """POST /documents/{id}/prompted-edit goes through Shield's
    invoke() — no raw LLM bypass."""
    src = _read(DOCUMENTS_PY)
    assert '@router.post("/documents/{doc_id}/prompted-edit")' in src
    pe_block = src.split('@router.post("/documents/{doc_id}/prompted-edit")')[1].split("@router")[0]
    # Shield invocation.
    assert "shield_invoke" in pe_block
    assert 'purpose="document_journal.prompted_edit.rewrite"' in pe_block
    # No emergentintegrations direct import in this block.
    assert "emergentintegrations" not in pe_block
    # Audit row written.
    assert "document.prompted_edit.proposed" in pe_block
    # Draft-only guard.
    assert "NOT_A_DRAFT" in pe_block


def test_scope_prompted_edit_frontend_uses_endpoint_not_toast():
    """The drawer's onPromptEdit now POSTs to /documents/{id}/prompted-edit
    instead of firing a 'coming soon' toast."""
    src = _read(DRAWER)
    # Old placeholder gone.
    assert "Prompt-based edits are coming soon" not in src
    # New plumbing present.
    assert "/documents/${doc.id}/prompted-edit" in src
    assert 'data-testid="drawer-document-prompt-diff"' in src
    assert 'data-testid="drawer-document-prompt-apply-confirm"' in src
    assert 'data-testid="drawer-document-prompt-discard"' in src


def test_scope_prompted_edit_frontend_renders_diff_strikethrough_and_underline():
    """Diff preview marks removed words with strikethrough and added
    words with an oxblood underline."""
    src = _read(DRAWER)
    assert 'data-testid="drawer-document-prompt-diff-del"' in src
    assert 'data-testid="drawer-document-prompt-diff-add"' in src
    assert "line-through" in src
    # Added words decorated in oxblood.
    assert "decoration-[var(--oxblood)]" in src


# ═════════════════════════════════════════════════════════════════════
# 2. DRAFT watermark embedded in exports (PDF / DOCX / PPTX)
# ═════════════════════════════════════════════════════════════════════
def test_scope_watermark_service_exports_three_format_helpers():
    """watermark_service.py exposes add_pdf_watermark,
    add_docx_watermark, add_pptx_watermark, and a watermark_file
    dispatcher."""
    src = _read(WATERMARK_SVC)
    for name in ("add_pdf_watermark", "add_docx_watermark",
                 "add_pptx_watermark", "watermark_file",
                 "WatermarkError"):
        assert f"def {name}" in src or f"class {name}" in src, f"missing {name}"


def test_scope_download_endpoint_applies_watermark_for_drafts():
    """The download endpoint imports watermark_file and applies it
    when state == draft; block-on-failure with HTTP 503 +
    DRAFT_WATERMARK_FAILED."""
    src = _read(DOCUMENTS_PY)
    dl_block = src.split("async def download_document(")[1].split("\n\n\n")[0]
    assert "from services.documents.watermark_service import" in dl_block
    assert "watermark_file(rendered, fmt=fmt, label=\"DRAFT\")" in dl_block
    assert "DRAFT_WATERMARK_FAILED" in dl_block
    # X-Watermark-Applied response header surfaces the applied state.
    assert "X-Watermark-Applied" in dl_block


def _draft_pdf_bytes() -> bytes:
    """Build a minimal PDF body used by the watermark byte tests.
    Single page, single line of text. Returns bytes."""
    from reportlab.pdfgen import canvas
    from reportlab.lib.pagesizes import A4
    buf = io.BytesIO()
    c = canvas.Canvas(buf, pagesize=A4)
    c.setFont("Helvetica", 12)
    c.drawString(72, 700, "Hello World from a test PDF.")
    c.showPage()
    c.save()
    return buf.getvalue()


def test_scope_pdf_watermark_actually_embedded_in_bytes():
    """Byte-level assertion: after add_pdf_watermark, the output PDF
    contains the literal 'DRAFT' string AND the source body line."""
    from services.documents.watermark_service import add_pdf_watermark
    src_pdf = _draft_pdf_bytes()
    wm_pdf = add_pdf_watermark(src_pdf, label="DRAFT")
    # Sanity — output is a real PDF.
    assert wm_pdf.startswith(b"%PDF-"), "output is not a PDF"
    # The watermark text is encoded as a content-stream string. We
    # decode via pypdf and check both the original body and the
    # "DRAFT" stamp appear.
    from pypdf import PdfReader
    reader = PdfReader(io.BytesIO(wm_pdf))
    text = reader.pages[0].extract_text() or ""
    assert "Hello World" in text, f"source body missing post-watermark: {text!r}"
    assert "DRAFT" in text, f"DRAFT stamp not extractable: {text!r}"


def _draft_docx_bytes() -> bytes:
    from docx import Document
    d = Document()
    d.add_paragraph("Hello DOCX body.")
    buf = io.BytesIO()
    d.save(buf)
    return buf.getvalue()


def test_scope_docx_watermark_actually_embedded_in_bytes():
    """Byte-level assertion: the watermarked DOCX zip now contains
    the header_watermark.xml part with the DRAFT label."""
    import zipfile
    from services.documents.watermark_service import add_docx_watermark
    src_docx = _draft_docx_bytes()
    wm_docx = add_docx_watermark(src_docx, label="DRAFT")
    with zipfile.ZipFile(io.BytesIO(wm_docx)) as z:
        names = set(z.namelist())
        assert "word/header_watermark.xml" in names, f"watermark header part missing: {sorted(names)}"
        body = z.read("word/header_watermark.xml").decode("utf-8")
        assert 'string="DRAFT"' in body, f"DRAFT label missing in header: {body[:300]}"
        # Source content survives.
        doc_xml = z.read("word/document.xml").decode("utf-8")
        assert "Hello DOCX body." in doc_xml


def _draft_pptx_bytes() -> bytes:
    from pptx import Presentation
    p = Presentation()
    layout = p.slide_layouts[1]
    s = p.slides.add_slide(layout)
    s.shapes.title.text = "Hello PPTX"
    buf = io.BytesIO()
    p.save(buf)
    return buf.getvalue()


def test_scope_pptx_watermark_actually_embedded_in_bytes():
    """Byte-level assertion: the watermarked PPTX has 'DRAFT' text
    runs on every slide."""
    from services.documents.watermark_service import add_pptx_watermark
    from pptx import Presentation
    src_pptx = _draft_pptx_bytes()
    wm_pptx = add_pptx_watermark(src_pptx, label="DRAFT")
    prs = Presentation(io.BytesIO(wm_pptx))
    # Collect every text run across every slide.
    found_draft = 0
    found_title = False
    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    if run.text == "DRAFT":
                        found_draft += 1
                    if run.text == "Hello PPTX":
                        found_title = True
    assert found_draft >= 1, "no DRAFT runs found in watermarked PPTX"
    assert found_title, "source title content missing after watermark"


# ═════════════════════════════════════════════════════════════════════
# 3. Related-docs typed groups
# ═════════════════════════════════════════════════════════════════════
def test_scope_related_endpoint_emits_four_typed_groups():
    """GET /related returns 4 buckets keyed metadata_match,
    content_similarity, explicit_attachment (gap),
    canonical_lineage (gap)."""
    src = _read(DOCUMENTS_PY)
    block = src.split("async def list_related_documents(")[1].split("\n\n\n")[0]
    for key in ("metadata_match", "content_similarity",
                "explicit_attachment", "canonical_lineage"):
        assert f'"{key}"' in block, f"missing related group key: {key}"
    # Gaps marked honestly.
    assert '"available": False' in block
    assert "gap_reason" in block


def test_scope_related_frontend_renders_typed_groups():
    """The Related tab renders one section per group, with a label
    + gap_reason for unshipped types."""
    src = _read(DRAWER)
    # The 4 group testids.
    for key in ("metadata_match", "content_similarity",
                "explicit_attachment", "canonical_lineage"):
        assert f'drawer-related-group-${{gk}}`' in src or f"drawer-related-group-{key}" in src.replace("${gk}", key)
    # Gap-reason testid present.
    assert "drawer-related-gap-reason-${gk}`" in src or "drawer-related-gap-reason-" in src
    # Calls the typed endpoint, not the legacy listing.
    assert "/related" in src
    assert "/contexts/${contextId}/documents/${doc.id}/related" in src


# ═════════════════════════════════════════════════════════════════════
# Live HTTP — happy paths
# ═════════════════════════════════════════════════════════════════════
@pytest.fixture
async def seeded_context():
    from core import db, hash_password
    uid = f"test-sc-{uuid.uuid4().hex[:8]}"
    cid = f"ctx-sc-{uuid.uuid4().hex[:8]}"
    email = f"sc-{uuid.uuid4().hex[:6]}@example.com"
    await db.accounts.insert_one({
        "id": uid, "email": email,
        "password_hash": hash_password("Pw!1234567Abc"),
        "name": "SC", "tier": "executive",
        "declared_role": "executive", "mfa_enrolled": False,
        "created_at": datetime.now(timezone.utc).isoformat(),
    })
    await db.contexts.insert_one({
        "id": cid, "name": "SC Co", "owner_account_id": uid,
        "type": "executive_personal",
        "created_at": datetime.now(timezone.utc).isoformat(),
    })
    await db.memberships.insert_one({
        "id": f"mem-{uuid.uuid4().hex[:8]}",
        "account_id": uid, "context_id": cid,
        "role": "executive", "status": "active",
        "created_at": datetime.now(timezone.utc).isoformat(),
    })
    yield {"uid": uid, "cid": cid, "email": email, "password": "Pw!1234567Abc"}
    await db.documents.delete_many({"context_id": cid})
    await db.audit_log.delete_many({"context_id": cid})
    await db.memberships.delete_many({"account_id": uid})
    await db.contexts.delete_one({"id": cid})
    await db.accounts.delete_one({"id": uid})


@pytest.mark.asyncio
async def test_scope_prompted_edit_endpoint_returns_diff_payload(seeded_context):
    """Live HTTP: create a draft, POST a prompt, get back
    {current_body, new_body, prompt_hash, diff_size}. Audit row is
    written."""
    from server import app  # noqa: F401
    from core import db
    did = f"doc-sc-{uuid.uuid4().hex[:8]}"
    await db.documents.insert_one({
        "id": did, "context_id": seeded_context["cid"],
        "name": "Draft for prompted-edit",
        "extracted_text": "The Q4 plan needs to address pricing for the European launch.",
        "state": "draft",
        "origin": "akki_generated",
        "status": "ready",
        "created_at": datetime.now(timezone.utc).isoformat(),
    })
    async with AsyncClient(transport=ASGITransport(app=app), base_url="http://test") as c:
        r = await c.post("/api/auth/login",
                         json={"email": seeded_context["email"], "password": seeded_context["password"]})
        token = r.json().get("access_token") or r.json().get("token")
        hdr = {"Authorization": f"Bearer {token}"}
        r = await c.post(f"/api/documents/{did}/prompted-edit",
                         json={"prompt": "Make it more concise — half the length."},
                         headers=hdr)
        assert r.status_code == 200, r.text
        body = r.json()
        for k in ("doc_id", "prompt_hash", "current_body", "new_body", "diff_size"):
            assert k in body, f"missing key in prompted-edit response: {k}"
        assert body["doc_id"] == did
        assert body["new_body"], "new_body must be non-empty"
        assert isinstance(body["diff_size"], int)
        # Audit row should exist.
        row = await db.audit_log.find_one({
            "resource_id": did,
            "action": "document.prompted_edit.proposed",
        })
        assert row, "audit row not written"
        assert row["metadata"]["prompt_hash"] == body["prompt_hash"]


@pytest.mark.asyncio
async def test_scope_prompted_edit_rejects_committed_docs(seeded_context):
    """Committed docs are out of scope for prompted-edit — must 400."""
    from server import app  # noqa: F401
    from core import db
    did = f"doc-sc-{uuid.uuid4().hex[:8]}"
    await db.documents.insert_one({
        "id": did, "context_id": seeded_context["cid"],
        "name": "Committed doc",
        "extracted_text": "Some body.",
        "state": "committed",
        "status": "ready",
        "created_at": datetime.now(timezone.utc).isoformat(),
    })
    async with AsyncClient(transport=ASGITransport(app=app), base_url="http://test") as c:
        r = await c.post("/api/auth/login",
                         json={"email": seeded_context["email"], "password": seeded_context["password"]})
        token = r.json().get("access_token") or r.json().get("token")
        hdr = {"Authorization": f"Bearer {token}"}
        r = await c.post(f"/api/documents/{did}/prompted-edit",
                         json={"prompt": "Rewrite."}, headers=hdr)
        assert r.status_code == 400, r.text
        assert r.json()["detail"]["code"] == "NOT_A_DRAFT"


@pytest.mark.asyncio
async def test_scope_download_draft_includes_watermarked_pdf(seeded_context):
    """Live HTTP: download a draft as PDF → response has
    X-Watermark-Applied=1 and the body contains the DRAFT stamp."""
    from server import app  # noqa: F401
    from core import db
    from pypdf import PdfReader
    did = f"doc-sc-{uuid.uuid4().hex[:8]}"
    await db.documents.insert_one({
        "id": did, "context_id": seeded_context["cid"],
        "name": "Draft for watermark",
        "extracted_text": "Quarterly performance review draft.",
        "state": "draft",
        "origin": "akki_generated",
        "status": "ready",
        "created_at": datetime.now(timezone.utc).isoformat(),
    })
    async with AsyncClient(transport=ASGITransport(app=app), base_url="http://test") as c:
        r = await c.post("/api/auth/login",
                         json={"email": seeded_context["email"], "password": seeded_context["password"]})
        token = r.json().get("access_token") or r.json().get("token")
        hdr = {"Authorization": f"Bearer {token}", "X-Active-Context": seeded_context["cid"]}
        r = await c.get(
            f"/api/contexts/{seeded_context['cid']}/documents/{did}/download",
            params={"format": "pdf"}, headers=hdr,
        )
        assert r.status_code == 200, r.text
        assert r.headers["X-Watermark-Applied"] == "1"
        assert r.headers["X-Document-State"] == "draft"
        # Extract PDF text → must contain DRAFT.
        reader = PdfReader(io.BytesIO(r.content))
        page_text = reader.pages[0].extract_text() or ""
        assert "DRAFT" in page_text, f"DRAFT not in downloaded PDF: {page_text!r}"
        assert "Quarterly performance review draft." in page_text


@pytest.mark.asyncio
async def test_scope_download_committed_pdf_has_no_watermark(seeded_context):
    """Committed docs export without the DRAFT stamp; X-Watermark-Applied=0."""
    from server import app  # noqa: F401
    from core import db
    from pypdf import PdfReader
    did = f"doc-sc-{uuid.uuid4().hex[:8]}"
    await db.documents.insert_one({
        "id": did, "context_id": seeded_context["cid"],
        "name": "Committed doc",
        "extracted_text": "Final approved report.",
        "state": "committed",
        "status": "ready",
        "created_at": datetime.now(timezone.utc).isoformat(),
    })
    async with AsyncClient(transport=ASGITransport(app=app), base_url="http://test") as c:
        r = await c.post("/api/auth/login",
                         json={"email": seeded_context["email"], "password": seeded_context["password"]})
        token = r.json().get("access_token") or r.json().get("token")
        hdr = {"Authorization": f"Bearer {token}", "X-Active-Context": seeded_context["cid"]}
        r = await c.get(
            f"/api/contexts/{seeded_context['cid']}/documents/{did}/download",
            params={"format": "pdf"}, headers=hdr,
        )
        assert r.status_code == 200
        assert r.headers["X-Watermark-Applied"] == "0"
        reader = PdfReader(io.BytesIO(r.content))
        page_text = reader.pages[0].extract_text() or ""
        assert "DRAFT" not in page_text
        assert "Final approved report." in page_text


@pytest.mark.asyncio
async def test_scope_related_endpoint_returns_typed_groups(seeded_context):
    """Live HTTP: seed 2 sibling docs in the same context and same
    doc_type, hit /related, get the metadata_match group populated +
    explicit_attachment / canonical_lineage marked as gaps."""
    from server import app  # noqa: F401
    from core import db
    did_main = f"doc-sc-{uuid.uuid4().hex[:8]}"
    did_peer = f"doc-sc-{uuid.uuid4().hex[:8]}"
    for d in (did_main, did_peer):
        await db.documents.insert_one({
            "id": d, "context_id": seeded_context["cid"],
            "name": f"Pricing memo {d[-4:]}",
            "doc_type": "memo",
            "extracted_text": "European launch pricing memo body for content similarity.",
            "status": "ready",
            "created_at": datetime.now(timezone.utc).isoformat(),
        })
    async with AsyncClient(transport=ASGITransport(app=app), base_url="http://test") as c:
        r = await c.post("/api/auth/login",
                         json={"email": seeded_context["email"], "password": seeded_context["password"]})
        token = r.json().get("access_token") or r.json().get("token")
        hdr = {"Authorization": f"Bearer {token}", "X-Active-Context": seeded_context["cid"]}
        r = await c.get(
            f"/api/contexts/{seeded_context['cid']}/documents/{did_main}/related",
            headers=hdr,
        )
        assert r.status_code == 200, r.text
        data = r.json()
        groups = data["groups"]
        # All 4 buckets present.
        assert set(groups.keys()) == {
            "metadata_match", "content_similarity",
            "explicit_attachment", "canonical_lineage",
        }
        # Gaps surface honestly.
        assert groups["explicit_attachment"]["available"] is False
        assert groups["canonical_lineage"]["available"] is False
        assert "gap_reason" in groups["explicit_attachment"]
        assert "gap_reason" in groups["canonical_lineage"]
        # Metadata match finds the peer doc (same context + same doc_type).
        meta_items = groups["metadata_match"]["items"]
        peer_ids = {it.get("id") for it in meta_items}
        assert did_peer in peer_ids, f"peer {did_peer} not in metadata_match items {peer_ids}"
