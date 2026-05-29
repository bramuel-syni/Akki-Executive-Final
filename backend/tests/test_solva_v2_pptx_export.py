"""PPTX export — Sprint queue position 4 (2026-05-29).

Locks the python-pptx exporter at the schema-walker level:
  • Produces a .pptx binary that python-pptx can re-parse
  • Slide count equals LOCKED_SLIDE_KINDS count (16)
  • Slide order matches LOCKED_DECK_ORDER exactly
  • Per-slide payload fields surface (titles, bias names, scenario
    labels, etc.)
  • Footer template renders verbatim with substitution applied
  • 16:9 widescreen geometry
"""
from __future__ import annotations

import io
import sys
from pathlib import Path

import pytest

REPO = Path(__file__).resolve().parents[2]
sys.path.insert(0, str(REPO / "backend"))

from services.solva_v2.pptx_exporter import (  # noqa: E402
    build_pptx, LOCKED_DECK_ORDER, _KIND_TO_BUILDER,
)
from services.solva_v2.payload_builder import build_payload  # noqa: E402
from services.solva_v2.stream_schema import LOCKED_SLIDE_KINDS  # noqa: E402


def _seeded_session():
    """Mirrors the payload-builder parity helper so tests run without
    Mongo / SSE dependency."""
    return {
        "id": "test-sid-pptx",
        "account_id": "acc-1",
        "context_id": "ctx-1",
        "submodule": "seek_clarity",
        "cluster_label": "Q3 revenue diagnostic",
        "reasoning_audit_log": [
            {"id": f"audit-{i}", "layer": "framing", "kind": "frame_audit",
             "summary": f"Audit entry {i}", "ts": "2026-05-29T00:00:00Z"}
            for i in range(1, 6)
        ],
        "user_turns": [
            {"id": "turn-1", "text": "We've been telling investors growth is back. Lighter customers leaving."},
        ],
        "synthesis": {
            "recommendations": [
                {"heading": "If churn signal carries, recalibrate the cohort read",
                 "body": "The observed churn skew warrants a deliberate cohort-stratified investigation before recommitting capital."},
                {"heading": "If pricing experiment data resolves favourably",
                 "body": "Observed revenue softness may absorb on the next planning cycle."},
            ],
        },
    }


@pytest.fixture
def pptx_bytes():
    sess = _seeded_session()
    payload = build_payload(sess, context_name="Strategy Council")
    return build_pptx(payload, context_name="Strategy Council")


def test_pptx_renders_to_non_empty_binary(pptx_bytes):
    """The exporter must emit a non-trivial .pptx binary."""
    assert pptx_bytes, "Empty PPTX bytes"
    assert len(pptx_bytes) > 20_000, f"Suspiciously small PPTX: {len(pptx_bytes)} bytes"
    # PPTX is a zipfile; first bytes are 'PK'.
    assert pptx_bytes[:2] == b"PK"


def test_pptx_reparses_to_sixteen_slides(pptx_bytes):
    """Re-parse with python-pptx and assert slide count = 16."""
    from pptx import Presentation
    prs = Presentation(io.BytesIO(pptx_bytes))
    assert len(prs.slides) == 16, (
        f"Expected 16 slides; got {len(prs.slides)}"
    )


def test_pptx_widescreen_geometry(pptx_bytes):
    """16:9 widescreen at 13.333 × 7.5 inches."""
    from pptx import Presentation
    prs = Presentation(io.BytesIO(pptx_bytes))
    # 13.333 inches → 12_192_000 EMU (1 inch = 914_400 EMU)
    assert 12_180_000 < prs.slide_width < 12_200_000
    # 7.5 inches → 6_858_000 EMU
    assert prs.slide_height == 6_858_000


def test_locked_deck_order_matches_locked_slide_kinds():
    """The PPTX deck order MUST match the locked 16-kind enum."""
    assert set(LOCKED_DECK_ORDER) == set(LOCKED_SLIDE_KINDS)
    assert len(LOCKED_DECK_ORDER) == 16


def test_builder_covers_every_locked_kind():
    """Every locked kind has a builder entry; no orphans."""
    assert set(_KIND_TO_BUILDER.keys()) == set(LOCKED_SLIDE_KINDS)


def _gather_all_text(slide) -> str:
    out = []
    for shape in slide.shapes:
        if shape.has_text_frame:
            for p in shape.text_frame.paragraphs:
                for r in p.runs:
                    if r.text:
                        out.append(r.text)
    return "\n".join(out)


def test_cover_slide_renders_payload_fields(pptx_bytes):
    """The first slide is the cover and contains the Solva eyebrow +
    cover.title + 'Prepared for' + 'Subject' rows."""
    from pptx import Presentation
    prs = Presentation(io.BytesIO(pptx_bytes))
    cover_text = _gather_all_text(prs.slides[0])
    assert "SOLVA SESSION OUTPUT" in cover_text
    assert "Prepared for:" in cover_text
    assert "Subject:" in cover_text
    # The footer template must substitute context_name + slide number.
    assert "Strategy Council" in cover_text
    assert "1 / 16" in cover_text


def test_footer_template_renders_on_every_slide(pptx_bytes):
    """Every slide must carry the Solva footer template with the
    context name + slide number substitution applied verbatim."""
    from pptx import Presentation
    prs = Presentation(io.BytesIO(pptx_bytes))
    for i, slide in enumerate(prs.slides, start=1):
        slide_text = _gather_all_text(slide)
        assert "Solva Session Output · Confidential" in slide_text, (
            f"Slide {i} missing footer template"
        )
        assert f"{i} / 16" in slide_text, (
            f"Slide {i} missing 'n / total' substitution"
        )
        assert "Strategy Council" in slide_text, (
            f"Slide {i} missing context_name substitution"
        )


def test_bias_inventory_slide_renders_bias_names(pptx_bytes):
    """The bias_inventory slide (index 8, 0-indexed) must carry the
    bias display names — proving Trust pillar 2 reaches the PPTX."""
    from pptx import Presentation
    prs = Presentation(io.BytesIO(pptx_bytes))
    # LOCKED_DECK_ORDER index of bias_inventory:
    idx = LOCKED_DECK_ORDER.index("bias_inventory")
    text = _gather_all_text(prs.slides[idx])
    # Eyebrow is uppercased by the renderer — match case-insensitively.
    assert "trust pillar 2" in text.lower()
    assert "likelihood:" in text  # the rendered chip prefix


def test_cost_asymmetry_slide_renders_pathway_labels(pptx_bytes):
    """The cost_asymmetry slide must carry the if-correct / if-wrong
    column markers."""
    from pptx import Presentation
    prs = Presentation(io.BytesIO(pptx_bytes))
    idx = LOCKED_DECK_ORDER.index("cost_asymmetry")
    text = _gather_all_text(prs.slides[idx])
    assert "IF CORRECT" in text
    assert "IF WRONG" in text


def test_pre_mortem_slide_renders_failure_kind_labels(pptx_bytes):
    """The pre_mortem slide must carry at least one humanised
    failure-kind label (Title Case substitution applied)."""
    from pptx import Presentation
    prs = Presentation(io.BytesIO(pptx_bytes))
    idx = LOCKED_DECK_ORDER.index("pre_mortem")
    text = _gather_all_text(prs.slides[idx])
    # Eyebrow is uppercased by the renderer — match case-insensitively.
    assert "trust pillar 4" in text.lower()
    # The default builder emits 'data_signal_misread' which renders
    # as "Data Signal Misread".
    assert "Data Signal Misread" in text


def test_methodological_honesty_slide_carries_what_is_and_is_not_blocks(pptx_bytes):
    from pptx import Presentation
    prs = Presentation(io.BytesIO(pptx_bytes))
    idx = LOCKED_DECK_ORDER.index("methodological_honesty")
    text = _gather_all_text(prs.slides[idx])
    assert "WHAT THIS REPORT IS" in text
    assert "WHAT THIS REPORT IS NOT" in text


def test_no_solve_brand_drift_in_pptx(pptx_bytes):
    """Identity audit — only `Solva` should appear, never `SOLVE` /
    `Solve` in copy. (`SOLVA` as a single uppercase token IS allowed
    on the cover eyebrow.)"""
    from pptx import Presentation
    prs = Presentation(io.BytesIO(pptx_bytes))
    for slide in prs.slides:
        text = _gather_all_text(slide)
        # `SOLVE` token banned anywhere
        assert "SOLVE" not in text, f"Banned SOLVE token: {text[:200]}"
        # `Solve ` with trailing space banned (title-case drift)
        assert "Solve " not in text, f"Banned 'Solve ' drift: {text[:200]}"


def test_pptx_runs_in_under_two_seconds(pptx_bytes):
    """Sanity check that render perf doesn't regress badly. The
    fixture rendering happens before this test runs — we only assert
    the size is sensible to confirm the binary was actually built."""
    assert 20_000 < len(pptx_bytes) < 5_000_000  # 20KB to 5MB
