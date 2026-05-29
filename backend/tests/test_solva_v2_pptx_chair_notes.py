"""Sprint Z2.0 — Chair-readable speaker notes lockdown test.

Locks the per-slide notes-pane copy so that:
  • Every slide carries the three baseline audit lines (Sourced …,
    Bias check: …, Confidence: …)
  • Bias inventory / pre-mortem / cost asymmetry slides also surface
    the slide-specific extras (named biases · failure-mode signals ·
    cost magnitudes)
  • Zero banned-vocabulary tokens land in the notes pane (mirrors the
    voice guide at `/app/docs/WEBSITE_BRIEF_V3.md`)
  • No emoji or marketing-style decorators leak through
"""
from __future__ import annotations

import io
import re
import sys
from pathlib import Path

import pytest

REPO = Path(__file__).resolve().parents[2]
sys.path.insert(0, str(REPO / "backend"))

from services.solva_v2.pptx_exporter import (  # noqa: E402
    build_pptx, LOCKED_DECK_ORDER,
)
from services.solva_v2.payload_builder import build_payload  # noqa: E402


def _seeded_session():
    return {
        "id": "test-sid-pptx-notes",
        "account_id": "acc-1",
        "context_id": "ctx-1",
        "submodule": "seek_clarity",
        "cluster_label": "Q3 revenue diagnostic",
        "reasoning_audit_log": [
            {"id": f"audit-{i}", "layer": "framing", "kind": "frame_audit",
             "summary": f"Audit entry {i}", "ts": "2026-05-29T00:00:00Z"}
            for i in range(1, 6)
        ],
        "user_turns": [
            {"id": "turn-1", "text": "We've been telling investors growth is back."},
        ],
        "synthesis": {
            "recommendations": [
                {"heading": "If churn signal carries, recalibrate the cohort read",
                 "body": "The observed churn skew warrants deliberate cohort-stratified investigation."},
            ],
        },
    }


@pytest.fixture(scope="module")
def parsed():
    from pptx import Presentation
    sess = _seeded_session()
    payload = build_payload(sess, context_name="Strategy Council")
    pptx_bytes = build_pptx(payload, context_name="Strategy Council")
    return Presentation(io.BytesIO(pptx_bytes))


def _notes_text(slide) -> str:
    """Return the verbatim concatenated notes pane text for a slide."""
    if not slide.has_notes_slide:
        return ""
    return slide.notes_slide.notes_text_frame.text


def test_every_slide_has_speaker_notes(parsed):
    """No silent slides — every one of the 16 slides carries the
    chair-readable audit footer."""
    missing = []
    for i, slide in enumerate(parsed.slides, start=1):
        if not _notes_text(slide).strip():
            missing.append(i)
    assert not missing, f"Slides without speaker notes: {missing}"


def test_three_baseline_lines_on_every_slide(parsed):
    """Every slide's notes carry the three audit lines."""
    for i, slide in enumerate(parsed.slides, start=1):
        text = _notes_text(slide)
        # Line 1 — sources + documents + grounding
        assert re.search(r"Sourced from \d+ input", text), (
            f"Slide {i} missing 'Sourced from N inputs' line · text={text[:200]}"
        )
        assert "document" in text.lower(), f"Slide {i} missing document count"
        assert "Evidence-grounding: passed" in text, (
            f"Slide {i} missing evidence-grounding line"
        )
        # Line 2 — bias check
        assert "Bias check:" in text, f"Slide {i} missing 'Bias check:' line"
        # Line 3 — confidence
        assert "Confidence:" in text, f"Slide {i} missing 'Confidence:' line"


def test_voice_exemplar_format_matches_brief(parsed):
    """The cover slide's notes MUST read in the Economist register
    locked by the voice brief (`/app/docs/WEBSITE_BRIEF_V3.md`)."""
    cover = parsed.slides[0]
    text = _notes_text(cover)
    # Spec exemplar from the brief:
    #   "Sourced from 4 inputs. 2 documents cited. Evidence-grounding: passed."
    # We assert the STRUCTURE not the literal numbers.
    line_re = re.compile(
        r"Sourced from \d+ inputs?\. "
        r"\d+ documents? cited\. "
        r"Evidence-grounding: passed\."
    )
    assert line_re.search(text), (
        f"Cover notes don't match locked exemplar format · text={text[:200]}"
    )


def test_bias_inventory_slide_surfaces_named_biases_in_notes(parsed):
    """Per the brief: 'For the bias-inventory slide specifically, also
    surface the 3 named biases in the notes so they survive even if
    the slide is hidden by a viewer.'"""
    idx = LOCKED_DECK_ORDER.index("bias_inventory")
    text = _notes_text(parsed.slides[idx])
    # The default payload_builder emits at least one bias.
    assert "This slide ·" in text, "Bias inventory notes missing slide-specific lines"
    # Likelihood pill from the chip surfaces in the notes.
    assert "likelihood" in text


def test_pre_mortem_slide_surfaces_triggering_signals_in_notes(parsed):
    """Per the brief: 'For the pre-mortem ... slides, surface the
    triggering signals ... in the notes so a chair scanning for
    "what should I worry about" can find them.'"""
    idx = LOCKED_DECK_ORDER.index("pre_mortem")
    text = _notes_text(parsed.slides[idx])
    assert "Watch for ·" in text
    assert "signals:" in text


def test_cost_asymmetry_slide_surfaces_magnitudes_in_notes(parsed):
    """Per the brief: '... surface the cost magnitudes in the notes.'"""
    idx = LOCKED_DECK_ORDER.index("cost_asymmetry")
    text = _notes_text(parsed.slides[idx])
    assert "Asymmetry ·" in text
    assert "magnitude" in text


# ─────────────────────────────────────────────────────────────────
# Voice-guide lockdown — banned vocabulary from
# /app/docs/WEBSITE_BRIEF_V3.md MUST NOT appear in notes
# ─────────────────────────────────────────────────────────────────


# Mirrors the 19-word ban list from the Founding Cohort Deck v3.
# Multi-word phrases are anchored. The Economist + senior-peer + restraint
# tests are validated by humans; this list is the machine-enforced
# floor.
_BANNED_PHRASES = (
    # generic vendor-speak
    "AI-powered", "AI-driven", "ai-powered", "ai-driven",
    "game-changer", "game changing",
    "synergy", "synergies",
    "leverage ", "leveraging",
    "supercharge", "supercharged",
    "seamless", "seamlessly",
    "revolutionary",
    "cutting-edge", "cutting edge",
    "disrupt", "disrupting", "disruptive",
    "frictionless",
    "unlock", "unlocking",
    "empower", "empowering",
    "all-under-one-roof", "all under one roof",
    # marketing decorators
    "✨", "🎯", "🚀", "🔥",
    # generic transform / insight / dashboard / solutions (banned in the
    # affirmative-only register)
    " insights ", " insights.",   # the noun used in vendor-speak
)


def test_no_banned_vocabulary_in_any_slide_notes(parsed):
    """Voice-brief lockdown: zero banned-vocabulary tokens across all
    16 slides' speaker notes."""
    offenders = []
    for i, slide in enumerate(parsed.slides, start=1):
        text = _notes_text(slide)
        for phrase in _BANNED_PHRASES:
            if phrase in text:
                offenders.append((i, phrase, text[:140]))
    assert not offenders, (
        "Banned vocabulary detected in speaker notes:\n  "
        + "\n  ".join(
            f"slide {i}: {p!r} in {t!r}"
            for (i, p, t) in offenders
        )
    )


def test_no_emoji_in_speaker_notes(parsed):
    """No emoji, no marketing decorators — chair-readable means
    plain-English plain-ASCII (plus the brand bullet middle-dot ·)."""
    emoji_re = re.compile(
        "["
        "\U0001F300-\U0001F9FF"   # general emoji blocks
        "\U0001F600-\U0001F64F"   # emoticons
        "\U0001F680-\U0001F6FF"   # transport
        "\u2600-\u27BF"           # misc symbols + dingbats
        "]"
    )
    offenders = []
    for i, slide in enumerate(parsed.slides, start=1):
        text = _notes_text(slide)
        m = emoji_re.search(text)
        if m:
            offenders.append((i, m.group(0), text[:140]))
    assert not offenders, (
        f"Emoji / dingbat detected in speaker notes: {offenders}"
    )


def test_notes_carry_solva_brand_only(parsed):
    """Solva-identity guard — `SOLVE` / `Solve ` drift banned in notes."""
    for i, slide in enumerate(parsed.slides, start=1):
        text = _notes_text(slide)
        assert "SOLVE" not in text, f"SOLVE drift in slide {i}: {text[:160]}"
        assert "Solve " not in text, f"'Solve ' drift in slide {i}: {text[:160]}"
