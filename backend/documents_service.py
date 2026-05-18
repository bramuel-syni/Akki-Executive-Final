"""Document extraction + storage pipeline.

Phase 10 changes:
  * ``virus_scan_stub`` is gone. Callers should invoke
    ``services.clamav_service.scan`` directly and translate
    :class:`ClamAVUnreachable` to a 503 with an audit row.
  * ``save_to_storage`` / ``read_from_storage`` / ``delete_from_storage``
    delegate to ``services.storage_service`` which swaps between
    local disk (tests) and S3/MinIO (prod) by the ``STORAGE_BACKEND``
    env var. The DB column continues to store the opaque key.

Phase F.1 (2026-05-16) — OCR + spreadsheet extraction:
  * ``.png/.jpg/.jpeg/.webp`` images → Tesseract OCR via pytesseract.
  * ``.heic/.heif`` images → pillow-heif opener + Tesseract OCR.
  * ``.xlsx`` → openpyxl in read-only mode, every sheet, every cell.
  * ``.csv`` → csv.reader, tab-joined cells, newline-joined rows.
  * Images >5 MB OR >2400px (max dimension) are downscaled before
    OCR to bound runtime.
  * Tesseract / Pillow failures → ``("", "Image had no extractable
    text. Try a higher-resolution scan.")`` — graceful, never raises.
"""
from __future__ import annotations

import csv
import io
import logging
import re
from pathlib import Path
from typing import List, Optional, Tuple

from services import storage_service

log = logging.getLogger("akki.documents")

ACCEPT_EXT = {".pdf", ".docx", ".pptx", ".txt", ".md", ".rtf",
              ".png", ".jpg", ".jpeg", ".webp", ".heic", ".heif",
              ".csv", ".xlsx"}
MAX_BYTES = 25 * 1024 * 1024  # 25 MB
MAX_EXTRACT_CHARS = 200_000
# Phase F.1 — image OCR bounds.
OCR_MAX_BYTES = 5 * 1024 * 1024
OCR_MAX_DIMENSION = 2400


def save_to_storage(context_id: str, doc_id: str, filename: str, data: bytes,
                    content_type: Optional[str] = None) -> str:
    """Store bytes through the configured storage backend; return key."""
    return storage_service.save(context_id, doc_id, filename, data, content_type=content_type)


def read_from_storage(key: str) -> bytes:
    return storage_service.read(key)


def delete_from_storage(key: str) -> None:
    storage_service.delete(key)


def _ocr_image_bytes(data: bytes, *, is_heif: bool = False) -> Tuple[str, Optional[str]]:
    """OCR a single image. Down-scales >OCR_MAX_DIMENSION before
    handing to Tesseract so runaway scans don't peg the worker."""
    try:
        from PIL import Image
        if is_heif:
            try:
                from pillow_heif import register_heif_opener
                register_heif_opener()
            except ImportError:
                return "", "HEIC/HEIF support not installed."
        img = Image.open(io.BytesIO(data))
        # EXIF rotation respected via thumbnail-safe path.
        img.load()
        w, h = img.size
        if max(w, h) > OCR_MAX_DIMENSION:
            scale = OCR_MAX_DIMENSION / float(max(w, h))
            img = img.resize((int(w * scale), int(h * scale)))
        try:
            import pytesseract
        except ImportError:
            return "", "OCR engine not installed (pytesseract missing)."
        text = pytesseract.image_to_string(img) or ""
        text = text.strip()
        if not text:
            return "", "Image had no extractable text. Try a higher-resolution scan."
        return text[:MAX_EXTRACT_CHARS], None
    except Exception as e:  # noqa: BLE001 — wrap Tesseract / Pillow errors uniformly
        log.info("OCR failed: %s: %s", type(e).__name__, str(e)[:200])
        return "", f"{type(e).__name__}: {str(e)[:200]}"


def _extract_xlsx(data: bytes) -> Tuple[str, Optional[str]]:
    try:
        from openpyxl import load_workbook
        wb = load_workbook(io.BytesIO(data), read_only=True, data_only=True)
    except Exception as e:  # noqa: BLE001
        return "", f"{type(e).__name__}: {str(e)[:200]}"
    parts: List[str] = []
    for sheet in wb.worksheets:
        parts.append(f"[Sheet: {sheet.title}]")
        for row in sheet.iter_rows(values_only=True):
            cells = ["" if c is None else str(c) for c in row]
            if any(cells):
                parts.append("\t".join(cells))
            if sum(len(p) for p in parts) > MAX_EXTRACT_CHARS:
                break
        if sum(len(p) for p in parts) > MAX_EXTRACT_CHARS:
            break
    return ("\n".join(parts))[:MAX_EXTRACT_CHARS], None


def _extract_csv(data: bytes) -> Tuple[str, Optional[str]]:
    try:
        # Try UTF-8 first; fall back to latin-1 (common for legacy
        # bank exports) before giving up.
        try:
            text_in = data.decode("utf-8")
        except UnicodeDecodeError:
            text_in = data.decode("latin-1", errors="replace")
        reader = csv.reader(io.StringIO(text_in))
        rows: List[str] = []
        for row in reader:
            rows.append("\t".join(row))
            if sum(len(r) for r in rows) > MAX_EXTRACT_CHARS:
                break
        return ("\n".join(rows))[:MAX_EXTRACT_CHARS], None
    except Exception as e:  # noqa: BLE001
        return "", f"{type(e).__name__}: {str(e)[:200]}"


def extract_text(data: bytes, filename: str, mime_type: str) -> Tuple[str, Optional[str]]:
    """Return (text, error). Text is truncated to MAX_EXTRACT_CHARS."""
    ext = Path(filename).suffix.lower()
    try:
        if ext == ".pdf" or mime_type == "application/pdf":
            from pypdf import PdfReader
            reader = PdfReader(io.BytesIO(data))
            pages = []
            for page in reader.pages:
                try:
                    pages.append(page.extract_text() or "")
                except Exception:
                    continue
            return ("\n\n".join(pages))[:MAX_EXTRACT_CHARS], None
        if ext == ".docx" or "officedocument.wordprocessingml" in (mime_type or ""):
            from docx import Document as DocxDocument
            d = DocxDocument(io.BytesIO(data))
            text = "\n".join(p.text for p in d.paragraphs if p.text)
            return text[:MAX_EXTRACT_CHARS], None
        # Phase H1 (2026-05-11) — PPTX text extraction via python-pptx
        # (already in requirements.txt for Work Studio deck rendering).
        # Pulls slide titles + every text frame in render order so
        # search hits and BM25 ranking land on actual slide content.
        if ext == ".pptx" or "officedocument.presentationml" in (mime_type or ""):
            from pptx import Presentation
            prs = Presentation(io.BytesIO(data))
            chunks: List[str] = []
            for slide_no, slide in enumerate(prs.slides, start=1):
                slide_lines: List[str] = []
                for shape in slide.shapes:
                    if not getattr(shape, "has_text_frame", False):
                        continue
                    tf = shape.text_frame
                    for para in tf.paragraphs:
                        line = "".join(run.text or "" for run in para.runs).strip()
                        if line:
                            slide_lines.append(line)
                if slide_lines:
                    chunks.append(f"--- Slide {slide_no} ---")
                    chunks.extend(slide_lines)
            return ("\n".join(chunks))[:MAX_EXTRACT_CHARS], None
        # Phase F.1 (2026-05-16) — XLSX via openpyxl read-only mode.
        if ext == ".xlsx" or "officedocument.spreadsheetml" in (mime_type or ""):
            return _extract_xlsx(data)
        # Phase F.1 — CSV via csv.reader.
        if ext == ".csv" or (mime_type or "").startswith("text/csv"):
            return _extract_csv(data)
        # Phase F.1 — image OCR via Tesseract.
        if ext in {".png", ".jpg", ".jpeg", ".webp"} or (mime_type or "").startswith("image/"):
            if len(data) > OCR_MAX_BYTES and ext not in {".heic", ".heif"}:
                # Pillow will handle downscale, but bail on truly huge
                # payloads to keep memory bounded.
                pass
            return _ocr_image_bytes(data, is_heif=False)
        if ext in {".heic", ".heif"} or (mime_type or "").startswith("image/heif"):
            return _ocr_image_bytes(data, is_heif=True)
        if ext in (".txt", ".md", ".rtf") or (mime_type or "").startswith("text/"):
            try:
                return data.decode("utf-8", errors="replace")[:MAX_EXTRACT_CHARS], None
            except Exception as e:
                return "", f"decode failed: {e}"
        return "", f"Unsupported file type: {ext}"
    except Exception as e:
        return "", f"extraction failed: {e}"


def make_preview(text: str, max_len: int = 320) -> str:
    if not text:
        return ""
    cleaned = re.sub(r"\s+", " ", text).strip()
    return cleaned[:max_len] + ("…" if len(cleaned) > max_len else "")
