"""
PPTX generator — programmatic python-pptx; no Jinja, no template files.
Mimics the MYDAWA Board Briefing signature:
  - 16:9 widescreen (13.33" × 7.5")
  - Persistent left sidebar on every content slide:
      <COMPANY>  ·  <DOCUMENT_TYPE>  ·  <PROGRAMME>
      with a slide-number badge (01, 02, ...)
  - Slide 1 (cover): title + subtitle + framework spine
  - Slide 2 (executive_brief one-pager): "If you read nothing else"
  - Content slides: ALL-CAPS section title, body, optional table or bullets
  - Closing slide: short recap + brand line
"""
from __future__ import annotations
from io import BytesIO
from typing import List, Optional

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from lxml import etree

from .brief import Brief, BriefSection, BriefTable, FIDELITY_HIGH, FIDELITY_LOW

# Brand palette ------------------------------------------------------------
INK    = RGBColor(0x0E, 0x0E, 0x0E)
OXBLD  = RGBColor(0x6E, 0x14, 0x18)
CREAM  = RGBColor(0xFA, 0xF6, 0xEF)
MUTED  = RGBColor(0x55, 0x55, 0x55)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
SIDEBAR_BG = RGBColor(0x0E, 0x0E, 0x0E)

BODY_FONT  = "Calibri"
HEAD_FONT  = "Calibri"   # match MYDAWA — sans throughout

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
SIDEBAR_W = Inches(0.85)
CONTENT_LEFT = Inches(1.10)
CONTENT_W = SLIDE_W - CONTENT_LEFT - Inches(0.6)


def _set_text(tf, text, *, font=BODY_FONT, size_pt=18, bold=False, italic=False,
              colour=INK, align=PP_ALIGN.LEFT, all_caps=False):
    tf.word_wrap = True
    tf.text = ""
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text.upper() if all_caps else text
    run.font.name = font
    run.font.size = Pt(size_pt)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = colour
    return p


def _add_text(slide, *, left, top, width, height, text,
              font=BODY_FONT, size_pt=18, bold=False, italic=False,
              colour=INK, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, all_caps=False):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.margin_left = Pt(0); tf.margin_right = Pt(0)
    tf.margin_top = Pt(0); tf.margin_bottom = Pt(0)
    tf.vertical_anchor = anchor
    _set_text(tf, text, font=font, size_pt=size_pt, bold=bold, italic=italic,
              colour=colour, align=align, all_caps=all_caps)
    return box


def _add_filled_rect(slide, left, top, width, height, colour):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = colour
    shape.line.fill.background()
    return shape


def _vertical_rotated_text(slide, *, left, top, width, height, text, colour,
                           size_pt=10, font=HEAD_FONT, bold=True):
    """Add a textbox with vertical text orientation by injecting bodyPr/@vert."""
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.margin_left = Pt(2); tf.margin_right = Pt(2)
    tf.margin_top = Pt(2); tf.margin_bottom = Pt(2)
    bodyPr = tf._txBody.bodyPr
    bodyPr.set("vert", "vert270")  # bottom-up vertical
    bodyPr.set("anchor", "ctr")
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size_pt)
    run.font.bold = bold
    run.font.color.rgb = colour
    # spacing for readability
    rPr = run._r.get_or_add_rPr()
    rPr.set("spc", "200")
    return box


def _add_sidebar(slide, brief: Brief, slide_number: int, *, fidelity: str):
    """Persistent left sidebar — only on content slides at high fidelity, or
    a thin slide-number marker at low fidelity."""
    if fidelity == FIDELITY_LOW:
        # thin oxblood number badge only
        _add_text(slide, left=Inches(0.35), top=Inches(0.30),
                  width=Inches(0.6), height=Inches(0.3),
                  text=f"{slide_number:02d}", font=HEAD_FONT, size_pt=10,
                  bold=True, colour=OXBLD)
        return

    # full dark sidebar
    _add_filled_rect(slide, left=Inches(0), top=Inches(0),
                     width=SIDEBAR_W, height=SLIDE_H, colour=SIDEBAR_BG)

    # rotated tagline
    sidebar_text_parts = [brief.company_label, brief.document_type]
    if brief.programme:
        sidebar_text_parts.append(brief.programme)
    tagline = "   ·   ".join(s for s in sidebar_text_parts if s)
    _vertical_rotated_text(slide,
                           left=Inches(0.10), top=Inches(0.6),
                           width=SIDEBAR_W - Inches(0.20), height=SLIDE_H - Inches(1.5),
                           text=tagline.upper(), colour=CREAM,
                           size_pt=10, bold=True)
    # slide-number badge
    _add_text(slide, left=Inches(0.10), top=SLIDE_H - Inches(0.7),
              width=SIDEBAR_W - Inches(0.20), height=Inches(0.4),
              text=f"{slide_number:02d}", font=HEAD_FONT, size_pt=14, bold=True,
              colour=CREAM, align=PP_ALIGN.CENTER)


def _add_table(slide, t: BriefTable, *, top: Emu, fidelity: str) -> Emu:
    """High-fidelity: real pptx table. Low-fidelity: bullets."""
    if fidelity == FIDELITY_LOW:
        # bullets fallback
        body = "  ·  ".join(t.headers) + "\n" + "\n".join(
            "  ·  ".join(r) for r in t.rows
        )
        box = _add_text(slide, left=CONTENT_LEFT, top=top,
                        width=CONTENT_W, height=Inches(2.5),
                        text=body, font=BODY_FONT, size_pt=12, colour=INK)
        return top + Inches(2.6)

    rows = 1 + len(t.rows)
    cols = len(t.headers)
    table_h = Inches(0.5) + Inches(0.4) * len(t.rows)
    table_shape = slide.shapes.add_table(rows, cols, CONTENT_LEFT, top,
                                         CONTENT_W, table_h)
    tbl = table_shape.table
    # column widths: equal
    col_w = int(CONTENT_W / cols)
    for c in range(cols):
        tbl.columns[c].width = col_w

    # header
    for c, h in enumerate(t.headers):
        cell = tbl.cell(0, c)
        cell.fill.solid()
        cell.fill.fore_color.rgb = SIDEBAR_BG
        cell.text_frame.text = ""
        p = cell.text_frame.paragraphs[0]
        run = p.add_run()
        run.text = h.upper()
        run.font.name = HEAD_FONT
        run.font.size = Pt(10)
        run.font.bold = True
        run.font.color.rgb = WHITE
    # body
    for r_idx, row in enumerate(t.rows):
        for c_idx, val in enumerate(row[:cols]):
            cell = tbl.cell(r_idx + 1, c_idx)
            cell.text_frame.text = ""
            p = cell.text_frame.paragraphs[0]
            run = p.add_run()
            run.text = str(val)
            run.font.name = BODY_FONT
            run.font.size = Pt(10)
            run.font.color.rgb = INK
    return top + table_h + Inches(0.2)


def _add_section_slide(prs, brief: Brief, sec: BriefSection, slide_number: int):
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    _add_sidebar(slide, brief, slide_number, fidelity=brief.fidelity)

    # Kicker
    kicker = (sec.kicker or sec.title).upper()
    _add_text(slide, left=CONTENT_LEFT, top=Inches(0.45),
              width=CONTENT_W, height=Inches(0.4),
              text=kicker, font=HEAD_FONT, size_pt=11, bold=True,
              colour=OXBLD, all_caps=True)
    # Title — ALL CAPS structurally per MYDAWA
    _add_text(slide, left=CONTENT_LEFT, top=Inches(0.85),
              width=CONTENT_W, height=Inches(1.0),
              text=sec.title, font=HEAD_FONT, size_pt=28, bold=True,
              colour=INK, all_caps=True)
    # Body
    cur = Inches(2.0)
    if sec.body_paragraphs:
        joined = "\n\n".join(sec.body_paragraphs)
        # fit body to remaining vertical space
        _add_text(slide, left=CONTENT_LEFT, top=cur,
                  width=CONTENT_W, height=Inches(2.5),
                  text=joined, font=BODY_FONT, size_pt=14, colour=INK)
        cur = cur + Inches(2.6)
    if sec.bullets:
        bullets = "\n".join(f"•  {b}" for b in sec.bullets)
        _add_text(slide, left=CONTENT_LEFT, top=cur,
                  width=CONTENT_W, height=Inches(2.0),
                  text=bullets, font=BODY_FONT, size_pt=14, colour=INK)
        cur = cur + Inches(2.0)
    for t in sec.tables:
        cur = _add_table(slide, t, top=cur, fidelity=brief.fidelity)


def _add_cover(prs, brief: Brief):
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    if brief.fidelity == FIDELITY_HIGH:
        _add_filled_rect(slide, Inches(0), Inches(0), SLIDE_W, SLIDE_H, CREAM)
        # accent bar
        _add_filled_rect(slide, Inches(1.0), Inches(0.9),
                         Inches(0.6), Inches(0.08), OXBLD)
    # Document type kicker
    _add_text(slide, left=Inches(1.0), top=Inches(1.05),
              width=Inches(11), height=Inches(0.4),
              text=brief.document_type, font=HEAD_FONT, size_pt=14, bold=True,
              colour=OXBLD, all_caps=True)
    # Title
    _add_text(slide, left=Inches(1.0), top=Inches(1.6),
              width=Inches(11), height=Inches(2.4),
              text=brief.title, font=HEAD_FONT, size_pt=44, bold=True, colour=INK)
    # Subtitle
    if brief.subtitle:
        _add_text(slide, left=Inches(1.0), top=Inches(4.1),
                  width=Inches(11), height=Inches(0.6),
                  text=brief.subtitle, font=HEAD_FONT, size_pt=18,
                  italic=True, colour=MUTED)
    # Framework spine
    if brief.framework_spine:
        _add_text(slide, left=Inches(1.0), top=Inches(5.0),
                  width=Inches(11), height=Inches(0.5),
                  text=brief.framework_spine, font=HEAD_FONT, size_pt=14,
                  bold=True, colour=OXBLD, all_caps=True)
    # Host / version / date row
    meta = []
    if brief.host_org_line: meta.append(brief.host_org_line)
    if brief.version: meta.append(brief.version)
    if brief.date_text: meta.append(brief.date_text)
    if meta:
        _add_text(slide, left=Inches(1.0), top=Inches(6.5),
                  width=Inches(11), height=Inches(0.4),
                  text="  ·  ".join(meta), font=HEAD_FONT, size_pt=12,
                  colour=MUTED)


def _add_exec_one_pager(prs, brief: Brief):
    """The "if you read nothing else" device — distinct from cover, used
    for executive_brief depth as the single content slide."""
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    _add_sidebar(slide, brief, slide_number=2, fidelity=brief.fidelity)
    _add_text(slide, left=CONTENT_LEFT, top=Inches(0.45),
              width=CONTENT_W, height=Inches(0.4),
              text="IF YOU READ NOTHING ELSE", font=HEAD_FONT, size_pt=11,
              bold=True, colour=OXBLD, all_caps=True)
    _add_text(slide, left=CONTENT_LEFT, top=Inches(0.85),
              width=CONTENT_W, height=Inches(1.4),
              text=brief.title, font=HEAD_FONT, size_pt=28, bold=True,
              colour=INK)
    if brief.cover_lead_paragraph:
        _add_text(slide, left=CONTENT_LEFT, top=Inches(2.4),
                  width=CONTENT_W, height=Inches(2.0),
                  text=brief.cover_lead_paragraph, font=BODY_FONT, size_pt=15,
                  colour=INK)
    # Top 3 recommendations
    recs = brief.sections[0].bullets if brief.sections else []
    if recs:
        _add_text(slide, left=CONTENT_LEFT, top=Inches(4.5),
                  width=CONTENT_W, height=Inches(0.4),
                  text="THE CALL", font=HEAD_FONT, size_pt=11, bold=True,
                  colour=OXBLD, all_caps=True)
        bullets = "\n".join(f"•  {r}" for r in recs[:3])
        _add_text(slide, left=CONTENT_LEFT, top=Inches(4.95),
                  width=CONTENT_W, height=Inches(2.1),
                  text=bullets, font=BODY_FONT, size_pt=14, colour=INK)


def _add_closing(prs, brief: Brief, slide_number: int):
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    _add_sidebar(slide, brief, slide_number, fidelity=brief.fidelity)
    _add_text(slide, left=CONTENT_LEFT, top=Inches(0.45),
              width=CONTENT_W, height=Inches(0.4),
              text="RECAP", font=HEAD_FONT, size_pt=11, bold=True,
              colour=OXBLD, all_caps=True)
    if brief.closing_recap:
        _add_text(slide, left=CONTENT_LEFT, top=Inches(2.5),
                  width=CONTENT_W, height=Inches(2.5),
                  text=brief.closing_recap, font=HEAD_FONT, size_pt=22,
                  bold=False, colour=INK)
    if brief.closing_brand_line:
        _add_text(slide, left=CONTENT_LEFT, top=Inches(6.5),
                  width=CONTENT_W, height=Inches(0.4),
                  text=brief.closing_brand_line, font=HEAD_FONT, size_pt=12,
                  bold=True, colour=MUTED, italic=True, all_caps=True)


def render_pptx(brief: Brief) -> bytes:
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    _add_cover(prs, brief)
    if brief.depth == "executive_brief":
        # 1-3 slides total: cover, exec one-pager, closing
        _add_exec_one_pager(prs, brief)
        _add_closing(prs, brief, slide_number=3)
    else:
        # board_summary: 8-12 slides ; deep_dive: 18+ slides
        slide_no = 2
        for sec in brief.sections:
            _add_section_slide(prs, brief, sec, slide_number=slide_no)
            slide_no += 1
        _add_closing(prs, brief, slide_number=slide_no)

    buf = BytesIO()
    prs.save(buf)
    return buf.getvalue()
