"""Phase C.2 — Work Studio export router.

Endpoints:
  POST /api/contexts/{cid}/work-studio/export/{kind}    — start
  GET  /api/contexts/{cid}/work-studio/exports/{eid}    — poll status
  GET  /api/contexts/{cid}/work-studio/exports/{eid}/download  — pin-token + bytes

Storage: `/app/backend/uploads/work_studio_exports/<eid>.<ext>` (local
backend in dev — same path the rest of the app uses).

LLM content stage: two-pass canonical via `services.two_pass` machinery.
Pass 1 (silent reasoning) and Pass 2 (strict JSON content_dict) both
recorded in `chat_audit_log` per the B.2 contract, with the new payload
keys `export_kind` and `export_artefact_id`.

Thin-input refusal (deterministic): when the user-supplied description /
objective / scope fails the same `detect_thin_input` heuristic used by
the chat path, the export is marked `failed` with reason `thin_input`
and the verbatim memo refusal text is surfaced in the API response so
the UI can render it.
"""
from __future__ import annotations

import asyncio
import io
import json
import logging
import os
import re
import secrets
import time
import uuid
import hashlib
from datetime import datetime, timezone, timedelta
from pathlib import Path
from typing import Any, Dict, List, Optional, Literal, Tuple

from fastapi import APIRouter, BackgroundTasks, Depends, File, Form, HTTPException, Request, UploadFile
from fastapi.responses import Response
from pydantic import BaseModel, Field

from core import (
    db, now, iso, write_audit, require_context_membership, get_current_account,
)
from services import two_pass as _tp
from services import work_studio_export as _ex

logger = logging.getLogger("akki.work_studio_export.router")

router = APIRouter(prefix="/api")

# Local storage root
_EXPORTS_ROOT = Path("/app/backend/uploads/work_studio_exports")
_EXPORTS_ROOT.mkdir(parents=True, exist_ok=True)

# Download token TTL.
_DOWNLOAD_TTL_SECONDS = 15 * 60

# =============================================================================
# Thin-input shape detector — independent of the regex-based chat
# detector in services.two_pass (which scans for decision/strategy
# verbs). The memo's deterministic refusal MUST fire before any LLM
# call; that means we cannot rely on the model to recognise
# "x"/"x"/"x"-style inputs as thin. Conservative thresholds:
#   • Any field with stripped len < 5      → thin.
#   • A field lacking ANY 4+-letter word   → thin.
# Used by both Export (3 fields) and Enhance (1 instructions field).
# =============================================================================
def _is_thin_text_fields(fields: Dict[str, str]) -> Optional[Dict[str, Any]]:
    short = [k for k, v in fields.items() if len((v or "").strip()) < 5]
    real_word = re.compile(r"[A-Za-z]{4,}")
    no_real_word = [k for k, v in fields.items() if not real_word.search(v or "")]
    if not short and not no_real_word:
        return None
    return {
        "trigger": "thin_signature",
        "fields_too_short": short,
        "fields_lacking_4letter_word": no_real_word,
        "field_lens": {k: len((v or "").strip()) for k, v in fields.items()},
    }


def _is_thin_export_shape(body: "ExportRequestIn") -> Optional[Dict[str, Any]]:
    return _is_thin_text_fields({
        "description": body.description or "",
        "objective":   body.objective or "",
        "scope":       body.scope or "",
    })


def _is_thin_enhance_shape(instructions: str) -> Optional[Dict[str, Any]]:
    return _is_thin_text_fields({"instructions": instructions or ""})


# =============================================================================
# Request / response shapes
# =============================================================================
class ExportRequestIn(BaseModel):
    description: str = Field(min_length=1, max_length=2000)
    objective:   str = Field(min_length=1, max_length=2000)
    scope:       str = Field(min_length=1, max_length=2000)
    output_format: Literal["docx", "pptx", "pdf", "auto"] = "auto"


class ExportStatusOut(BaseModel):
    export_id: str
    status: str
    kind: str
    output_format: str
    file_name: Optional[str] = None
    sensitivity_band: Optional[str] = None
    error: Optional[str] = None
    refusal_text: Optional[str] = None
    created_at: str
    completed_at: Optional[str] = None
    download_token: Optional[str] = None
    chat_audit_id: Optional[str] = None
    continue_chat_id: Optional[str] = None
    continue_doc_id: Optional[str] = None
    source: Optional[str] = None


# C.3 — Enhance accepts file_upload OR source_artefact_id (existing
# completed export), an instructions string, and an output_format. We
# parse these from multipart/form fields rather than JSON so the file
# upload can ride the same request.
#
# Chunk 3 (2026-05-13, WS-R06) — `minutes` added as a registered
# enhance kind. Minutes are functionally a Report-style narrative
# artefact, so the renderer treats them as Report (docx-only output).
# Accepted upload extensions extended to match (`.docx`, `.pdf`,
# `.txt` — plain text is common for minutes drafts pasted from notes).
# QA-2026-05-16-043 (2026-05-18): `committee_pack` added as a
# first-class enhance kind. The QA spec lists "Minutes / Deck /
# Report / Committee Pack" as the four enhance outputs. Backend-
# wise a committee pack uses the same docx-report render path as
# a Report — the renderer doesn't need a new code path, only the
# kind guard.
_ENHANCE_KINDS = ("deck", "report", "minutes", "committee_pack")
_ENHANCE_MAX_BYTES = 25 * 1024 * 1024  # 25 MB ceiling, mirrors chat_attach
_ENHANCE_ACCEPTED_EXT_BY_KIND = {
    "deck":    {".pptx", ".pdf"},
    "report":  {".docx", ".pdf"},
    "minutes": {".docx", ".pdf", ".txt"},
    # QA-2026-05-16-043 — committee pack ships as DOCX/PDF input.
    "committee_pack": {".docx", ".pdf"},
}


# =============================================================================
# Output-format selection
# =============================================================================
# QA-2026-05-16-043 — committee_pack mirrors report shape (docx-out).
_AUTO_FORMAT = {"brief": "docx", "deck": "pptx", "report": "docx", "minutes": "docx", "committee_pack": "docx"}
_VALID_KINDS = ("brief", "deck", "report", "minutes", "committee_pack")


def _resolve_format(kind: str, requested: str) -> str:
    if requested == "auto":
        return _AUTO_FORMAT[kind]
    if kind == "deck" and requested == "pdf":
        # Soft fork — see services/work_studio_export.render_deck_pdf.
        # We force pptx so the user always gets a file. The audit log
        # records the substitution.
        return "pptx"
    return requested


# =============================================================================
# Hash-chain helper for chat_audit_log (mirrors routers/chat.py)
# =============================================================================
async def _append_chat_audit(
    *, account_id: str, chat_id: str, action: str,
    payload: Dict[str, Any], request: Optional[Request],
) -> Dict[str, Any]:
    """Append a row to chat_audit_log with the same canonical-JSON
    hashing layout used by routers/chat.py:_append_audit. We don't
    import the helper directly to avoid coupling, but the hash shape
    must match exactly so reviewers can grep-verify both paths."""
    # Find the previous hash for this account.
    prev = await db.chat_audit_log.find_one(
        {"account_id": account_id},
        sort=[("at", -1)], projection={"_id": 0, "row_hash": 1},
    )
    prev_hash = (prev or {}).get("row_hash", "GENESIS-AKKI-CHAT-AUDIT-2026")
    rid = str(uuid.uuid4())
    at = iso(now())
    ip = ""
    ua_sha = ""
    if request is not None:
        try:
            ip = request.client.host if request.client else ""
            ua = request.headers.get("user-agent", "")
            ua_sha = hashlib.sha256(ua.encode()).hexdigest()[:16] if ua else ""
        except Exception:
            pass
    canonical = {
        "prev": prev_hash, "id": rid, "at": at,
        "account_id": account_id, "chat_id": chat_id,
        "action": action, "payload": payload, "ip": ip, "ua_sha": ua_sha,
    }
    row_hash = hashlib.sha256(
        json.dumps(canonical, sort_keys=True, separators=(",", ":")).encode()
    ).hexdigest()
    row = {
        "id": rid, "at": at, "account_id": account_id, "chat_id": chat_id,
        "action": action, "payload": payload, "ip": ip, "ua_sha": ua_sha,
        "prev_hash": prev_hash, "row_hash": row_hash,
    }
    await db.chat_audit_log.insert_one(row)
    return row


# =============================================================================
# Deterministic thin-input refusal — shared by Export and Enhance
# =============================================================================
async def _emit_thin_refusal(
    *, export_id: str, account_id: str, context_id: str,
    kind: str, source: str = "export",
    detection: Dict[str, Any],
) -> None:
    """Persist a deterministic thin-input refusal for an export OR
    enhance row. Emits the verbatim memo refusal text (substring guarantees
    "candidate framings" + "don't have enough to weight them honestly"),
    appends a `chat.refused` row to the hash chain, and records a
    `work_studio.{export|enhance}.failed` audit row."""
    refusal_text = _tp.THIN_INPUT_REFUSAL_TEMPLATE.format(
        evidence_phrase=_tp.THIN_INPUT_FALLBACK_EVIDENCE
    )
    audit_row = await _append_chat_audit(
        account_id=account_id, chat_id=f"{source}-{export_id}",
        action="chat.refused", payload={
            "refusal_reason": "thin_input",
            "turn_class": "strategic_deliverable",
            "evidence_phrase": _tp.THIN_INPUT_FALLBACK_EVIDENCE,
            "evidence_source": "fallback_static",
            "detection": detection,
            "channel": source,
            "deterministic": True,
            "export_kind": (
                kind if source == "export" else f"enhance_{kind}"
            ),
            "export_artefact_id": export_id,
        }, request=None,
    )
    await db.work_studio_exports.update_one(
        {"id": export_id},
        {"$set": {
            "status": "failed",
            "error": "thin_input",
            "refusal_reason": "thin_input",
            "refusal_text": refusal_text,
            "chat_audit_id": audit_row["id"],
            "completed_at": iso(now()),
        }},
    )
    try:
        await write_audit(
            account_id=account_id, context_id=context_id,
            action=f"work_studio.{source}.failed", target_id=export_id,
            metadata={"export_id": export_id, "kind": kind,
                      "reason": "thin_input", "source": source},
        )
    except Exception:
        pass


# =============================================================================
# BM25 grounding (mirrors the chat path, scoped to context documents)
# =============================================================================
async def _grounding_for_context(context_id: str, query: str, top_k: int = 8) -> List[Dict[str, Any]]:
    """Return top-K paragraphs from the context's documents for the
    given query. Mirrors routers/chat._retrieve_grounding_paragraphs
    but limits to docs only (no scope filter)."""
    cursor = db.documents.find(
        {"context_id": context_id, "status": {"$ne": "deleted"}},
        {"_id": 0, "id": 1, "name": 1, "extracted_text": 1, "paragraphs": 1},
    )
    docs = await cursor.to_list(200)
    if not docs:
        return []
    try:
        from bm25 import score_bm25
    except Exception:
        return []
    corpus: List[Dict[str, Any]] = []
    for d in docs:
        paras = d.get("paragraphs") or []
        if paras:
            for p in paras:
                txt = p.get("text") or ""
                if txt.strip():
                    corpus.append({
                        "text": txt, "doc_id": d["id"],
                        "doc_name": d.get("name") or "Document",
                        "anchor": p.get("anchor_id"),
                    })
        elif (d.get("extracted_text") or "").strip():
            corpus.append({
                "text": d["extracted_text"][:3000], "doc_id": d["id"],
                "doc_name": d.get("name") or "Document", "anchor": None,
            })
    if not corpus:
        return []
    # bm25.score_bm25(query, chunks, *, k=12) returns top-k
    # [(score, chunk_dict), ...] already sorted desc. Each chunk dict
    # in `corpus` carries `text` (required by score_bm25) plus the
    # doc_id / doc_name / anchor we pass through to _format_grounding.
    ranked = score_bm25(query, corpus, k=top_k)
    return [c for s, c in ranked if s > 0]


def _format_grounding(paragraphs: List[Dict[str, Any]]) -> str:
    if not paragraphs:
        return ""
    rail = ["[GROUNDING]"]
    for p in paragraphs:
        head = f"  - {p['doc_name']}"
        if p.get("anchor"):
            head += f" (¶ {p['anchor']})"
        head += f": {p['text'][:600]}"
        rail.append(head)
    rail.append("[/GROUNDING]")
    return "\n".join(rail)


def _flat_user_input(body: ExportRequestIn) -> str:
    return (
        f"DESCRIPTION: {body.description.strip()}\n\n"
        f"OBJECTIVE: {body.objective.strip()}\n\n"
        f"SCOPE: {body.scope.strip()}"
    )


# =============================================================================
# LLM stage — two-pass canonical
# =============================================================================
def _strip_json_fence(s: str) -> str:
    s = (s or "").strip()
    # Strip ```json ... ``` fences if present.
    m = re.search(r"```(?:json)?\s*(.*?)```", s, re.DOTALL | re.IGNORECASE)
    if m:
        return m.group(1).strip()
    return s


def _try_parse_json(raw: str) -> Optional[Dict[str, Any]]:
    s = _strip_json_fence(raw)
    # Find the outermost JSON object.
    if not s.startswith("{"):
        first_brace = s.find("{")
        if first_brace == -1:
            return None
        s = s[first_brace:]
    last_brace = s.rfind("}")
    if last_brace != -1:
        s = s[: last_brace + 1]
    try:
        return json.loads(s)
    except Exception:
        return None


async def _set_export_phase(export_id: str, phase_index: int, phase_label: str) -> None:
    """Phase L.a (2026-05-27) — write a phase marker on the export row.

    Best-effort: failures must NEVER block the worker. The streaming
    observer endpoint polls `phase_index` at 500ms cadence and emits
    SSE events when it advances; missing markers degrade UX (the
    progress log freezes between phases) but do not corrupt the
    underlying export."""
    try:
        await db.work_studio_exports.update_one(
            {"id": export_id},
            {"$set": {
                "phase_index": phase_index,
                "phase_label": phase_label,
                "phase_at":    iso(now()),
            }},
        )
    except Exception:
        logger.debug("phase marker write failed (export_id=%s, phase=%d)", export_id, phase_index)


async def _run_two_pass_for_export(
    *, kind: str, body: ExportRequestIn, ctx_doc: Dict[str, Any],
    grounding_block: str, citations_manifest: List[Dict[str, Any]],
    on_pass_boundary=None,
) -> Tuple[str, Dict[str, Any], int, int, Dict[str, Any]]:
    """Run Pass 1 (silent reasoning) + Pass 2 (strict JSON) over the
    universal LLM proxy. Returns (pass_1_text, content_dict,
    pass_1_ms, pass_2_ms, llm_meta).  `llm_meta` carries
    {"llm_pass1": {"provider","fallback"}, "llm_pass2": {...}} so the
    caller can persist it on the work_studio_exports row. Raises
    RuntimeError / ChatError on irrecoverable failure.

    Phase B.3 hardening — Pass 1 + Pass 2 used to call `LlmChat` directly
    against the universal LLM proxy. When the proxy 502'd on Claude every
    Work Studio export+enhance failed with `llm_error:ChatError`. Now
    each pass goes through `services.llm_streaming.collect_llm_text`
    which tries the direct Anthropic SDK first (when ANTHROPIC_API_KEY
    is set) and falls back to the proxy on any direct-call 5xx/network
    error. Audit signal lives on the row via `llm_meta`.
    """
    from services.llm_streaming import collect_llm_text, provider_for_model
    # Phase B (2026-05-13): `ChatError` was previously imported from
    # `emergentintegrations.llm.chat` to flag LLM-failure rows. We now
    # use a local marker class (`_WorkStudioLLMError`) so the CI guard
    # `test_no_direct_llm_calls_outside_shield` finds zero SDK imports
    # outside Shield. Behaviour is unchanged: outer `except Exception`
    # blocks still catch and persist `llm_error:_WorkStudioLLMError`.
    class _WorkStudioLLMError(RuntimeError):
        """Marker for Pass 1 / Pass 2 LLM failures. Preserves the prior
        `ChatError` ergonomics without importing the SDK exception."""

    user_text = _flat_user_input(body)

    # STUDIO sprint (2026-05-12) — W-22 fix: capture partial pass content
    # + per-pass provider metadata in a single dict reference. On any
    # raise, attach it to the exception as `.partial` so the caller can
    # persist what we DID succeed at before the failure. Without this,
    # debugging a "validation error on Pass 2" requires re-running the
    # whole thing.
    partial: Dict[str, Any] = {
        "llm_pass1": None,
        "llm_pass2": None,
        "pass1_text": None,
        "pass2_text_head": None,
    }

    def _attach(exc: Exception) -> Exception:
        try:
            setattr(exc, "partial", partial)
        except Exception:
            pass
        return exc

    # ── Pass 1 — reasoning rail (silent)
    pass_1_system = (
        "You are AKKI's reasoning rail for a Work Studio export. "
        "Apply the four-layer reasoning architecture to the user's "
        "task: candidate framings, triangulation against the grounding "
        "block, probability weighting (with confidence intervals), "
        "and reflection (three questions: what would change my mind? "
        "what's the explanation in six months if I got this wrong? "
        "what am I disappointed by?). Do NOT produce the deliverable. "
        "Do NOT emit JSON. Just the reasoning. Operating preferences "
        "(banned words, no glazing, lead with substance) apply."
    )
    p1_t0 = time.monotonic()
    p1_input = (grounding_block + "\n\n" if grounding_block else "") + user_text
    p1_model = "claude-sonnet-4-5-20250929"
    p1_text, p1_provider, p1_fallback, p1_err = await collect_llm_text(
        provider=provider_for_model(p1_model),
        model_id=p1_model,
        system_msg=pass_1_system,
        user_text=p1_input,
        session_id=f"akki-export-p1-{uuid.uuid4().hex[:8]}",
    )
    if p1_err:
        # Re-raise as ChatError so the caller's existing
        # `except Exception` handler tags the row as `llm_error:ChatError`
        # — preserves the existing exception ergonomics.
        partial["llm_pass1"] = {"provider": p1_provider, "fallback": bool(p1_fallback), "error": p1_err[:300]}
        raise _attach(_WorkStudioLLMError(p1_err))
    pass_1_text = p1_text
    pass_1_ms = int((time.monotonic() - p1_t0) * 1000)
    partial["llm_pass1"] = {"provider": p1_provider, "fallback": bool(p1_fallback)}
    partial["pass1_text"] = pass_1_text[:8000] if pass_1_text else None

    # Phase L.a (2026-05-27) — Pass 1 complete; emit "Composing." phase marker.
    if on_pass_boundary is not None:
        try:
            await on_pass_boundary("pass2_start")
        except Exception:
            pass

    # ── Pass 2 — JSON content_dict
    schema_doc = {
        "brief": (
            'Required keys: title, subtitle, classification ("Public"|'
            '"Internal"|"Confidential"|"Restricted"), period, '
            'generated_for, executive_summary (80-180 words), sections '
            '(list of {heading, paragraphs (list of strings, 1-3 items),'
            ' pullquote (optional), cites (list of 1-based indices into '
            'citations[])}), citations (list of {doc_id, doc_name, '
            'paragraph_anchor (optional)}). Do NOT invent doc_ids — use '
            'only entries from the citations manifest given below.'
        ),
        "deck": (
            'Required keys: title, subtitle, classification, period, '
            'generated_for, executive_summary (one or two short '
            'sentences), conclusion (one sentence), sections (4-6 '
            'items, each {heading, bullets (list of 3-5 prose lines), '
            'callout (optional), cites}), citations.'
        ),
        "report": (
            'Required keys: title, subtitle, classification, period, '
            'generated_for, executive_summary, sections (4-8 items, '
            'each {heading, subheading (optional), paragraphs (list of '
            '2-5 strings), pullquote (optional), cites}), '
            'recommendations (list of 1-5 short imperatives), '
            'citations.'
        ),
        # Chunk 3 (2026-05-13, WS-R06) — Minutes shape. Identical
        # structural keys to a Report (the same DOCX renderer is used)
        # but tighter section sizing (a typical minutes doc has many
        # short sections, not 4-8 long ones). Recommendations are
        # renamed `actions` semantically but we keep the key name as
        # `recommendations` so the renderer schema is unchanged.
        "minutes": (
            'Required keys: title (e.g. "Board Minutes — <date>"), '
            'subtitle (attendees in one line), classification, period, '
            'generated_for, executive_summary (one short paragraph of '
            'context and overall outcome), sections (one per agenda '
            'item, each {heading, subheading (optional, e.g. owner), '
            'paragraphs (list of 1-3 strings capturing the discussion '
            'and the decision), pullquote (optional, the verbatim '
            'decision), cites}), recommendations (list of 0-8 next '
            'actions, each phrased as "<owner> to <action> by <when>"), '
            'citations.'
        ),
        # QA-2026-05-16-043 (2026-05-18) — Committee Pack shape. Same
        # underlying schema as Report (the renderer is the same) but
        # the prompt steers the model towards a committee-meeting
        # deliverable: cover, attendees / quorum line in subtitle,
        # 4-6 substantive sections per agenda item, decision-oriented
        # pullquotes, and a clean action list. Recommendations key is
        # reused so the renderer stays untouched.
        "committee_pack": (
            'Required keys: title (e.g. "<Committee> Pack — <date>"), '
            'subtitle (attendees / quorum line), classification, period, '
            'generated_for, executive_summary (one short paragraph of '
            'context and the committee\'s overall recommendation), '
            'sections (4-6 items per agenda topic, each {heading, '
            'subheading (optional, e.g. owner / decision-status), '
            'paragraphs (list of 2-4 strings capturing the substance '
            'and the recommendation), pullquote (optional, the verbatim '
            'committee decision), cites}), recommendations (list of '
            '0-8 actions for the board, each phrased as "<owner> to '
            '<action> by <when>"), citations.'
        ),
    }[kind]

    cite_manifest_text = "\n".join(
        f"  [{i+1}] doc_id={c['doc_id']}  doc_name={c['doc_name']}"
        for i, c in enumerate(citations_manifest)
    ) or "  (no documents tethered — use citations: [] in your output)"

    pass_2_system = (
        "You are AKKI's deliverable composer. Produce a SINGLE JSON "
        "object — no markdown fences, no preamble, no commentary, just "
        "valid JSON parseable by json.loads. The schema is:\n"
        f"{schema_doc}\n\n"
        "Operating preferences (apply on every string you emit):\n"
        " - No glazing. Lead with substance.\n"
        " - Banned words: " + ", ".join(_tp.BANNED_WORDS) + ".\n"
        " - The Economist test, senior peer test, restraint test must pass.\n"
        " - If you make a claim you cannot trace to the grounding "
        "block or the citations manifest below, do not include it.\n"
        " - Citations manifest (use ONLY these; cites[] are 1-based "
        "indices into THIS list):\n"
        f"{cite_manifest_text}\n\n"
        "The Pass 1 reasoning has already been produced and is in "
        "[PASS_1_REASONING] below. Honour the selected framing fully."
    )
    p2_t0 = time.monotonic()
    p2_input = (
        f"[PASS_1_REASONING]\n{pass_1_text}\n[/PASS_1_REASONING]\n\n"
        + (grounding_block + "\n\n" if grounding_block else "")
        + user_text
        + "\n\nProduce the JSON now. Set classification to 'Internal' "
        "unless the description mentions sensitive material that warrants "
        "Confidential or Restricted. Set generated_for to the context name."
    )
    p2_model = "claude-sonnet-4-5-20250929"
    p2_text, p2_provider, p2_fallback, p2_err = await collect_llm_text(
        provider=provider_for_model(p2_model),
        model_id=p2_model,
        system_msg=pass_2_system,
        user_text=p2_input,
        session_id=f"akki-export-p2-{uuid.uuid4().hex[:8]}",
    )
    if p2_err:
        partial["llm_pass2"] = {"provider": p2_provider, "fallback": bool(p2_fallback), "error": p2_err[:300]}
        raise _attach(_WorkStudioLLMError(p2_err))
    pass_2_text = p2_text
    pass_2_ms = int((time.monotonic() - p2_t0) * 1000)
    partial["llm_pass2"] = {"provider": p2_provider, "fallback": bool(p2_fallback)}
    partial["pass2_text_head"] = pass_2_text[:2000] if pass_2_text else None

    parsed = _try_parse_json(pass_2_text)
    if parsed is None:
        raise _attach(RuntimeError(
            f"Pass 2 LLM did not return parseable JSON (kind={kind}). "
            f"Head: {pass_2_text[:200]!r}"
        ))

    # Force the canonical context name + period overrides if absent.
    parsed.setdefault("generated_for", ctx_doc.get("name") or "—")
    if not parsed.get("citations"):
        parsed["citations"] = citations_manifest

    # B.3 audit metadata — provider used + whether the proxy fallback fired.
    llm_meta: Dict[str, Any] = {
        "llm_pass1": {"provider": p1_provider, "fallback": bool(p1_fallback)},
        "llm_pass2": {"provider": p2_provider, "fallback": bool(p2_fallback)},
    }
    return pass_1_text, parsed, pass_1_ms, pass_2_ms, llm_meta


# =============================================================================
# Export worker (BackgroundTasks)
# =============================================================================
async def _run_export(
    *, export_id: str, account_id: str, context_id: str, kind: str,
    output_format: str, body: ExportRequestIn,
):
    """Execute the LLM stage + render + persist + audit. Runs in
    BackgroundTasks. Updates db.work_studio_exports as it progresses.
    """
    # Phase L.a (2026-05-27) — phase 0 ("Reading the cycle inputs.") fires
    # at entry. Subsequent phases are stamped at each boundary.
    await _set_export_phase(export_id, 0, "Reading the cycle inputs.")
    ctx_doc = await db.contexts.find_one({"id": context_id}, {"_id": 0})
    if not ctx_doc:
        await db.work_studio_exports.update_one(
            {"id": export_id},
            {"$set": {"status": "failed", "error": "context_missing",
                      "completed_at": iso(now())}},
        )
        return
    ctx_meta = {
        "context_name":     ctx_doc.get("name") or "—",
        "classification":   "Internal",
        "period":           "—",
        "generated_at_human": _ex.now_human(),
    }

    # ── PRE-GROUNDING deterministic thin-input refusal (Phase C.2 follow-up).
    # The export form collects three explicit fields; if any is too short
    # or lacks a real word, refuse here BEFORE any LLM call so the
    # refusal is resilient to upstream proxy errors. Emits the verbatim
    # memo refusal text — both required substrings are guaranteed to
    # appear in the persisted `refusal_text` and the failure error
    # column (`error="thin_input:<head>"`).
    shape_thin = _is_thin_export_shape(body)
    if shape_thin is not None:
        await _emit_thin_refusal(
            export_id=export_id, account_id=account_id, context_id=context_id,
            kind=kind, detection={**shape_thin, "stage": "pre_grounding"},
        )
        return

    # ── Synisense-shielded grounding context for the LLM stage. We
    # rely on the existing chat path's grounding paragraphs so the
    # documents are already de-identified at upload-time.
    await _set_export_phase(export_id, 1, "Checking the grounding contract.")
    user_input = _flat_user_input(body)
    grounding = await _grounding_for_context(context_id, user_input, top_k=8)
    grounding_block = _format_grounding(grounding)

    # Citation manifest = unique docs in grounding (id + name).
    seen: Dict[str, str] = {}
    for g in grounding:
        seen.setdefault(g["doc_id"], g["doc_name"])
    citations_manifest = [
        {"doc_id": k, "doc_name": v, "paragraph_anchor": None}
        for k, v in seen.items()
    ]

    # ── Thin-input deterministic refusal — regex pattern fallback for
    # decision/strategy verbs that the shape check (above) does not
    # catch (e.g. "should I", "what should we do"). Still pre-LLM.
    thin = _tp.detect_thin_input(
        turn_class="strategic_deliverable",
        text=user_input,
        attached_document_ids=[d["doc_id"] for d in citations_manifest],
        prior_substantive_turns=0,
    )
    if thin is not None:
        await _emit_thin_refusal(
            export_id=export_id, account_id=account_id, context_id=context_id,
            kind=kind, detection={**thin, "stage": "pre_llm_regex"},
        )
        return

    # ── Two-pass canonical LLM stage
    await _set_export_phase(export_id, 2, "Drafting the outline.")
    try:
        async def _on_pass_boundary(marker: str):
            if marker == "pass2_start":
                await _set_export_phase(export_id, 3, "Composing.")
        pass_1, content_dict, p1_ms, p2_ms, llm_meta = await _run_two_pass_for_export(
            kind=kind, body=body, ctx_doc=ctx_doc,
            grounding_block=grounding_block,
            citations_manifest=citations_manifest,
            on_pass_boundary=_on_pass_boundary,
        )
    except Exception as e:
        logger.exception("Export LLM stage failed")
        # STUDIO sprint (W-22): persist whatever partial pass metadata
        # + text we captured before the failure, so a later debug pass
        # can see exactly which pass failed and why.
        # Chunk 3 (2026-05-13) — also capture the exception message so
        # the user / ops can see WHICH key / arg / etc. blew up. Was:
        # `llm_error:KeyError` (opaque). Now: `llm_error:KeyError: 'sections'`.
        partial = getattr(e, "partial", None) or {}
        err_class = type(e).__name__
        err_msg = str(e).replace("\n", " ").strip()[:300]
        err_field = f"llm_error:{err_class}" + (f": {err_msg}" if err_msg else "")
        await db.work_studio_exports.update_one(
            {"id": export_id},
            {"$set": {
                "status": "failed",
                "error": err_field,
                "completed_at": iso(now()),
                "llm_pass1": partial.get("llm_pass1"),
                "llm_pass2": partial.get("llm_pass2"),
                "llm_pass1_text": partial.get("pass1_text"),
                "llm_pass2_text_head": partial.get("pass2_text_head"),
            }},
        )
        try:
            await write_audit(
                account_id=account_id, context_id=context_id,
                action="work_studio.export.failed", target_id=export_id,
                metadata={"export_id": export_id, "kind": kind,
                          "reason": "llm_error"},
            )
        except Exception:
            pass
        return

    # ── Pre-render banned-word scan on the JSON content text. We don't
    # block on this — we record any hit on the audit row.
    pre_text = _ex.scrape_content_text(content_dict)
    pre_hit = _ex.scan_for_banned_words(pre_text)

    # ── Render
    await _set_export_phase(export_id, 4, "Rendering the artefact.")
    try:
        ctx_meta_full = {
            **ctx_meta,
            "context_name": ctx_doc.get("name") or "—",
        }
        if output_format == "docx" and kind == "brief":
            data, sha, fname = _ex.render_brief_docx(content_dict, ctx_meta_full)
        elif output_format == "docx" and kind == "report":
            data, sha, fname = _ex.render_report_docx(content_dict, ctx_meta_full)
        elif output_format == "pptx" and kind == "deck":
            data, sha, fname = _ex.render_deck_pptx(content_dict, ctx_meta_full)
        elif output_format == "pdf" and kind == "brief":
            data, sha, fname = _ex.render_brief_pdf(content_dict, ctx_meta_full)
        elif output_format == "pdf" and kind == "report":
            data, sha, fname = _ex.render_report_pdf(content_dict, ctx_meta_full)
        elif output_format == "pdf" and kind == "deck":
            # Soft fork — should have been resolved at request time.
            raise _ex.ContentValidationError(
                "PDF output is not supported for decks in this pod."
            )
        else:
            raise _ex.ContentValidationError(
                f"No renderer for kind={kind} format={output_format}."
            )
    except _ex.ContentValidationError as ve:
        logger.warning("Render rejected content: %s", ve)
        await db.work_studio_exports.update_one(
            {"id": export_id},
            {"$set": {"status": "failed",
                      "error": f"validation:{str(ve)[:200]}",
                      "completed_at": iso(now())}},
        )
        try:
            await write_audit(
                account_id=account_id, context_id=context_id,
                action="work_studio.export.failed", target_id=export_id,
                metadata={"export_id": export_id, "kind": kind,
                          "reason": "validation"},
            )
        except Exception:
            pass
        return
    except Exception as e:
        logger.exception("Render failed")
        await db.work_studio_exports.update_one(
            {"id": export_id},
            {"$set": {"status": "failed", "error": f"render:{type(e).__name__}",
                      "completed_at": iso(now())}},
        )
        return

    # ── Post-render banned-word scan (text scraped from the file).
    file_text = _ex.scrape_text(data, output_format) or pre_text
    post_hit = _ex.scan_for_banned_words(file_text)

    # ── Sensitivity scoring (deterministic — reuse Studio scorer).
    await _set_export_phase(export_id, 5, "Validating.")
    sensitivity_band = "INTERNAL"
    sensitivity_score = 0
    sensitivity_reasons: List[str] = []
    try:
        from studio_sensitivity import score_text
        sc = score_text(file_text)
        sensitivity_band = sc.get("band") or sensitivity_band
        sensitivity_score = sc.get("score") or 0
        sensitivity_reasons = sc.get("reasons") or []
    except Exception as e:  # noqa: BLE001
        logger.warning("studio_sensitivity unavailable: %s", e.__class__.__name__)

    # ── Persist file to storage_service path.
    fpath = _EXPORTS_ROOT / f"{export_id}.{output_format}"
    fpath.write_bytes(data)
    # ── Two-pass canonical chat_audit_log row (the headline B.2 audit row).
    msg_audit = await _append_chat_audit(
        account_id=account_id, chat_id=f"export-{export_id}",
        action="message.received", payload={
            "user_message_id": None,
            "reply_id": None,
            "model_id": "claude-sonnet-4-5", "mode": "export",
            "latency_ms": p1_ms + p2_ms,
            "char_len_reply": len(file_text),
            "channel": "export",
            "turn_class": "strategic_deliverable",
            "classifier_source": "forced",
            "classifier_latency_ms": 0,
            "four_check_surfaced": {"label": None, "ran": True},
            "refusal_reason": None,
            "two_pass": {
                "pass_1": pass_1[:24000],
                "pass_2": json.dumps(content_dict)[:24000],
                "pass_1_visible": False,
                "pass_1_present": True,
                "pass_1_ms": p1_ms,
                "pass_2_ms": p2_ms,
            },
            "voice_violation":
                ({"banned_word": post_hit, "retry_outcome": "ship_with_violation",
                  "before_text": "(post-render scrape)", "after_text": None}
                 if post_hit else None),
            "export_kind": kind,
            "export_artefact_id": export_id,
            "pre_render_banned_word": pre_hit,
            "sensitivity_band": sensitivity_band,
            "sensitivity_score": sensitivity_score,
        }, request=None,
    )

    if post_hit:
        try:
            await write_audit(
                account_id=account_id, context_id=context_id,
                action="work_studio.export.voice_violation",
                target_id=export_id,
                metadata={"export_id": export_id, "kind": kind,
                          "banned_word": post_hit},
            )
        except Exception:
            pass

    completed_at = iso(now())
    # Phase L.a (2026-05-27) — phase 6 ("Almost there.") fires just before the
    # final updates land + the export row flips to status=complete.
    await _set_export_phase(export_id, 6, "Almost there.")
    # ── C.3 — mint Continue-in-chat (chat row + artefact doc row)
    cont_chat_id, cont_doc_id = await _create_continue_chat(
        account_id=account_id, context_id=context_id, kind=kind,
        source="export", export_id=export_id,
        file_name=fname, file_path=str(fpath), output_format=output_format,
        extracted_text=file_text, sensitivity_band=sensitivity_band,
    )
    await db.work_studio_exports.update_one(
        {"id": export_id},
        {"$set": {
            "status": "complete",
            "file_name": fname,
            "file_path": str(fpath),
            "sha256": sha,
            "byte_len": len(data),
            "sensitivity_band": sensitivity_band,
            "sensitivity_score": sensitivity_score,
            "sensitivity_reasons": sensitivity_reasons,
            "completed_at": completed_at,
            "chat_audit_id": msg_audit["id"],
            "pass_1_ms": p1_ms, "pass_2_ms": p2_ms,
            # Phase B.3 — provider-used + fallback-fired audit signal,
            # parallel to what briefings.py persists. Either pass shows
            # `proxy_buffered` if the direct path 502'd and we fell back.
            "llm_pass1": llm_meta.get("llm_pass1"),
            "llm_pass2": llm_meta.get("llm_pass2"),
            "voice_violation": post_hit,
            "continue_chat_id": cont_chat_id,
            "continue_doc_id": cont_doc_id,
        }},
    )

    # Phase R.3 (2026-05-27) — emit work_studio.export.completed feature event.
    try:
        from services.cohort.feature_events import (
            emit_feature_event, WORK_STUDIO_EXPORT_COMPLETED,
        )
        acct_doc = await db.accounts.find_one(
            {"id": account_id}, {"_id": 0, "cohort_tag": 1},
        ) or {}
        await emit_feature_event(
            event_type=WORK_STUDIO_EXPORT_COMPLETED,
            account_id=account_id,
            cohort_tag=acct_doc.get("cohort_tag"),
            payload={
                "export_id": export_id, "context_id": context_id,
                "kind": kind, "output_format": output_format,
                "sensitivity_band": sensitivity_band,
            },
        )
    except Exception:
        pass
    try:
        await write_audit(
            account_id=account_id, context_id=context_id,
            action="work_studio.export.completed", target_id=export_id,
            metadata={
                "export_id": export_id, "kind": kind,
                "output_format": output_format,
                "sensitivity_band": sensitivity_band,
                "sha256": sha,
                "continue_chat_id": cont_chat_id,
                "continue_doc_id": cont_doc_id,
            },
        )
    except Exception:
        pass


# =============================================================================
# C.3 — Continue-in-chat helper
# After a successful export OR enhance, mint a chat tethered to the
# artefact: insert a row in db.chats (so existing chat router serves
# it as-is) and a row in db.documents (so the chat composer can show
# the artefact as an attachment chip via the standard ?attach=<doc_id>
# URL parameter on the Chat page).
# =============================================================================
async def _create_continue_chat(
    *, account_id: str, context_id: str, kind: str, source: str,
    export_id: str, file_name: str, file_path: str, output_format: str,
    extracted_text: str, sensitivity_band: str,
) -> Tuple[str, str]:
    chat_id = str(uuid.uuid4())
    doc_id = str(uuid.uuid4())
    created_at = iso(now())

    # Default model = the same Sonnet 4.5 the export used; user can
    # change the model in the Chat composer if they prefer.
    chat_row = {
        "id": chat_id,
        "account_id": account_id,
        "title": f"Continue · {kind.capitalize()} · {file_name[:60]}",
        "model_id": "claude-sonnet-4-5-20250929",
        "shielding_policy": "auto",
        "context_id": context_id,
        "status": "active",
        "message_count": 0,
        "last_message_preview": "",
        "last_message_at": None,
        "created_at": created_at,
        "updated_at": created_at,
        # Phase C.3 markers — non-hot-path, purely informational.
        "continue_source": source,        # "export" | "enhance"
        "continue_artefact_id": export_id,
    }
    await db.chats.insert_one(chat_row)

    # Best-effort sensitivity for the document chip.
    band_lower = (sensitivity_band or "INTERNAL").lower()
    doc_row = {
        "id": doc_id,
        "context_id": context_id,
        "name": file_name,
        "original_filename": file_name,
        "mime_type": {
            "docx": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
            "pptx": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
            "pdf":  "application/pdf",
        }.get(output_format, "application/octet-stream"),
        "size_bytes": (Path(file_path).stat().st_size if Path(file_path).exists() else 0),
        "storage_key": f"work_studio_exports/{export_id}.{output_format}",
        "status": "extracted",
        "extracted_text": extracted_text or "",
        "extracted_chars": len(extracted_text or ""),
        "preview": (extracted_text or "")[:320],
        "data_trust": "trusted",
        "doc_type": "work_studio_artefact",
        "uploaded_by": account_id,
        "source_channel": "work_studio_export",
        "chat_id": chat_id,
        "synisense_version": 0,  # the artefact is composed on already-shielded grounding
        "body_redacted": None,
        "sensitivity_band": band_lower,
        "sensitivity_score": 0,
        "sensitivity_label": (sensitivity_band or "INTERNAL").upper(),
        "created_at": created_at,
        "updated_at": created_at,
        # Phase C.3 marker so the journal can group "Continue in chat" docs.
        "work_studio_export_id": export_id,
    }
    await db.documents.insert_one(doc_row)

    # Hash-chained chat.created row so the audit chain is preserved.
    try:
        await _append_chat_audit(
            account_id=account_id, chat_id=chat_id,
            action="chat.created", payload={
                "model_id": "claude-sonnet-4-5-20250929",
                "shielding_policy": "auto",
                "context_id": context_id,
                "continue_source": source,
                "continue_artefact_id": export_id,
            }, request=None,
        )
    except Exception:  # noqa: BLE001
        logger.warning("continue_chat audit-append failed (non-fatal)")
    return chat_id, doc_id


# =============================================================================
# C.3 — Enhance source resolution + text extraction
# =============================================================================
def _file_ext(name: str) -> str:
    ext = (Path(name or "").suffix or "").lower()
    return ext


async def _resolve_enhance_source(
    *, context_id: str, account_id: str, kind: str,
    file: Optional[UploadFile], source_artefact_id: Optional[str],
) -> Tuple[bytes, str, str]:
    """Return (raw_bytes, source_filename, source_label).
    Validates extension matches `kind`. Raises HTTPException on errors."""
    if file is not None:
        data = await file.read()
        if not data:
            raise HTTPException(status_code=400, detail="Empty file upload.")
        if len(data) > _ENHANCE_MAX_BYTES:
            raise HTTPException(status_code=413, detail="File exceeds 25 MB ceiling.")
        ext = _file_ext(file.filename or "")
        accepted = _ENHANCE_ACCEPTED_EXT_BY_KIND[kind]
        if ext not in accepted:
            raise HTTPException(
                status_code=400,
                detail=f"For kind={kind}, accepted extensions are {sorted(accepted)}; got {ext or '(none)'}.",
            )
        # ClamAV scan with the same dev escape hatch as chat_attach.
        try:
            from services import clamav_service
            scan_res = await clamav_service.scan(data, file.filename, user_id=account_id)
            if not scan_res.clean:
                raise HTTPException(
                    status_code=400,
                    detail=f"Virus scan failed: {scan_res.signature or 'unclean'}.",
                )
        except HTTPException:
            raise
        except Exception:  # ClamAVUnreachable or anything else
            allow_unsafe = os.getenv("ALLOW_UNSAFE_UPLOADS", "false").lower() == "true"
            if not allow_unsafe:
                raise HTTPException(status_code=503, detail="Virus scanner unavailable.")
            logger.warning("clamav unreachable; proceeding under ALLOW_UNSAFE_UPLOADS")
        return data, (file.filename or f"upload{ext}"), "upload"

    if source_artefact_id:
        prev = await db.work_studio_exports.find_one(
            {"id": source_artefact_id, "context_id": context_id, "status": "complete"},
            {"_id": 0},
        )
        if not prev:
            raise HTTPException(status_code=404, detail="Source artefact not found.")
        if prev.get("account_id") != account_id:
            raise HTTPException(status_code=403, detail="Source artefact not yours.")
        fp = Path(prev.get("file_path") or "")
        if not fp.exists():
            raise HTTPException(status_code=404, detail="Source artefact file missing on disk.")
        return fp.read_bytes(), (prev.get("file_name") or fp.name), "source_artefact"

    raise HTTPException(
        status_code=400,
        detail="Provide either `file` (multipart) or `source_artefact_id` (form field).",
    )


def _extract_text_from_bytes(data: bytes, filename: str) -> str:
    """Best-effort text extraction. Reuses backend.documents_service for
    .docx/.pdf/.txt; handles .pptx natively (python-pptx) since the
    docs service does not."""
    ext = (Path(filename or "").suffix or "").lower()
    # Native .pptx path — extract every text frame across slides.
    if ext == ".pptx":
        try:
            from pptx import Presentation
            prs = Presentation(io.BytesIO(data))
            chunks: List[str] = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if not getattr(shape, "has_text_frame", False):
                        continue
                    tf = shape.text_frame
                    for para in tf.paragraphs:
                        line = "".join(run.text for run in para.runs).strip()
                        if line:
                            chunks.append(line)
            return ("\n".join(chunks))[:200000]
        except Exception as exc:  # noqa: BLE001
            logger.warning("pptx text-extract failed: %s", exc)
            return ""
    # Everything else — DOCX / PDF / TXT — defer to documents_service.
    try:
        from documents_service import extract_text
        text, err = extract_text(data, filename, "")
        if err:
            logger.warning("enhance text-extract returned error: %s", err)
        return text or ""
    except Exception as exc:  # noqa: BLE001
        logger.warning("enhance text-extract failed: %s", exc)
        return ""


# =============================================================================
# C.3 — Enhance worker
# =============================================================================
async def _run_enhance(
    *, export_id: str, account_id: str, context_id: str, kind: str,
    output_format: str, instructions: str, source_data: bytes,
    source_filename: str, source_label: str,
) -> None:
    """Run the canonical two-pass on (extracted text + instructions),
    render the enhanced output via the existing renderer, persist to
    work_studio_exports with source='enhance', and mint a continue chat."""
    ctx_doc = await db.contexts.find_one({"id": context_id}, {"_id": 0})
    if not ctx_doc:
        await db.work_studio_exports.update_one(
            {"id": export_id},
            {"$set": {"status": "failed", "error": "context_missing",
                      "completed_at": iso(now())}},
        )
        return

    # ── Extract text from the source bytes
    source_text = _extract_text_from_bytes(source_data, source_filename)
    if not source_text:
        # We cannot enhance what we cannot read.
        await db.work_studio_exports.update_one(
            {"id": export_id},
            {"$set": {"status": "failed",
                      "error": "source_extract_empty",
                      "completed_at": iso(now())}},
        )
        try:
            await write_audit(
                account_id=account_id, context_id=context_id,
                action="work_studio.enhance.failed", target_id=export_id,
                metadata={"export_id": export_id, "kind": kind,
                          "reason": "source_extract_empty",
                          "source_filename": source_filename},
            )
        except Exception:
            pass
        return

    # ── Synisense the source text (surface=enhance)
    body_text = source_text
    try:
        from services.synisense import shield_payload_async as _syn_shield
        shielded, _meta = await _syn_shield(
            source_text, surface="enhance", context_id=context_id, mode="redact",
        )
        body_text = shielded or source_text
    except Exception as exc:  # noqa: BLE001
        logger.warning("synisense shielding skipped on enhance: %s", exc)

    # ── Synthesise an ExportRequestIn from instructions + source
    SOURCE_CAP = 16000
    body_synth = ExportRequestIn(
        description=f"Enhance my {kind}: {instructions[:1500]}",
        objective=(
            "Apply the user's instructions to the source artefact text below. "
            "Preserve any factual claims that already cite a [n] reference; "
            "rewrite, drop, or reorder sections only as the instructions direct."
        ),
        scope=(
            f"SOURCE ARTEFACT (kind={kind}, filename={source_filename!r}, "
            f"truncated to {SOURCE_CAP} chars):\n\n"
            + body_text[:SOURCE_CAP]
        ),
        output_format=output_format,  # already resolved; not "auto"
    )

    # ── Two-pass canonical LLM stage. Grounding-block is the source
    # text itself (we treat it as the only [1] citation). Citation
    # manifest = the source artefact.
    grounding_block = (
        "[1] " + (body_text[:8000].replace("\n", " ").strip() or "(empty)")
    )
    citations_manifest = [
        {"doc_id": "source-artefact", "doc_name": source_filename,
         "paragraph_anchor": None}
    ]

    try:
        pass_1, content_dict, p1_ms, p2_ms, llm_meta = await _run_two_pass_for_export(
            kind=kind, body=body_synth, ctx_doc=ctx_doc,
            grounding_block=grounding_block,
            citations_manifest=citations_manifest,
        )
    except Exception as e:
        logger.exception("Enhance LLM stage failed")
        # Chunk 3 (2026-05-13) — capture the actual exception message,
        # not just the class name. Pre-Chunk 3 the row stored
        # `llm_error:KeyError` which gave the user (and ops) no clue
        # which key was missing. Now stores
        # `llm_error:KeyError: 'content_pass2'` so the failure point
        # is identifiable without backend log access.
        err_class = type(e).__name__
        err_msg = str(e).replace("\n", " ").strip()[:300]
        err_field = f"llm_error:{err_class}" + (f": {err_msg}" if err_msg else "")
        await db.work_studio_exports.update_one(
            {"id": export_id},
            {"$set": {"status": "failed", "error": err_field,
                      "completed_at": iso(now())}},
        )
        try:
            await write_audit(
                account_id=account_id, context_id=context_id,
                action="work_studio.enhance.failed", target_id=export_id,
                metadata={"export_id": export_id, "kind": kind,
                          "reason": "llm_error"},
            )
        except Exception:
            pass
        return

    # Defensive citations override — the LLM sometimes returns a
    # well-formed `citations: [{}]` (truthy but with empty doc_id),
    # which slips past `_run_two_pass_for_export`'s `if not citations:`
    # guard and lands the validator at n_cit=0. Force the source
    # artefact citation when the effective list is unusable.
    clean_cites = [
        c for c in (content_dict.get("citations") or [])
        if isinstance(c, dict) and (c.get("doc_id") or "").strip()
    ]
    if not clean_cites:
        content_dict["citations"] = citations_manifest

    # ── Pre-render banned-word scan
    pre_text = _ex.scrape_content_text(content_dict)
    pre_hit = _ex.scan_for_banned_words(pre_text)

    # ── Render
    try:
        ctx_meta_full = {
            "context_name":      ctx_doc.get("name") or "—",
            "classification":    "Internal",
            "period":            "—",
            "generated_at_human": _ex.now_human(),
        }
        if output_format == "docx" and kind == "report":
            data, sha, fname = _ex.render_report_docx(content_dict, ctx_meta_full)
        elif output_format == "docx" and kind == "minutes":
            # Chunk 3 (WS-R06) — Minutes use the same DOCX renderer as
            # Report. Minutes are a narrative artefact identical in
            # structure (headings + body paragraphs) to a Report, so
            # the same renderer path applies cleanly. A bespoke
            # Minutes renderer can branch off this if PO later wants
            # a different visual treatment.
            data, sha, fname = _ex.render_report_docx(content_dict, ctx_meta_full)
        elif output_format == "docx" and kind == "committee_pack":
            # QA-2026-05-16-043 (2026-05-18) — Committee Pack uses the
            # same DOCX renderer as Report. Same narrative shape, same
            # output format. A bespoke renderer can branch off this
            # once the PO finalises committee-specific styling.
            data, sha, fname = _ex.render_report_docx(content_dict, ctx_meta_full)
        elif output_format == "pptx" and kind == "deck":
            data, sha, fname = _ex.render_deck_pptx(content_dict, ctx_meta_full)
        elif output_format == "pdf" and kind == "report":
            data, sha, fname = _ex.render_report_pdf(content_dict, ctx_meta_full)
        elif output_format == "pdf" and kind == "minutes":
            data, sha, fname = _ex.render_report_pdf(content_dict, ctx_meta_full)
        elif output_format == "pdf" and kind == "committee_pack":
            data, sha, fname = _ex.render_report_pdf(content_dict, ctx_meta_full)
        else:
            raise _ex.ContentValidationError(
                f"Enhance has no renderer for kind={kind} format={output_format}."
            )
    except _ex.ContentValidationError as ve:
        logger.warning("Enhance render rejected content: %s", ve)
        await db.work_studio_exports.update_one(
            {"id": export_id},
            {"$set": {"status": "failed",
                      "error": f"validation:{str(ve)[:200]}",
                      "completed_at": iso(now())}},
        )
        return
    except Exception as e:
        logger.exception("Enhance render failed")
        await db.work_studio_exports.update_one(
            {"id": export_id},
            {"$set": {"status": "failed", "error": f"render:{type(e).__name__}",
                      "completed_at": iso(now())}},
        )
        return

    # ── Post-render banned-word scan + sensitivity
    file_text = _ex.scrape_text(data, output_format) or pre_text
    post_hit = _ex.scan_for_banned_words(file_text)
    sensitivity_band = "INTERNAL"
    sensitivity_score = 0
    sensitivity_reasons: List[str] = []
    try:
        from studio_sensitivity import score_text
        sc = score_text(file_text)
        sensitivity_band = sc.get("band") or sensitivity_band
        sensitivity_score = sc.get("score") or 0
        sensitivity_reasons = sc.get("reasons") or []
    except Exception as e:  # noqa: BLE001
        logger.warning("studio_sensitivity unavailable: %s", e.__class__.__name__)

    # ── Persist
    fpath = _EXPORTS_ROOT / f"{export_id}.{output_format}"
    fpath.write_bytes(data)

    # ── Hash-chained chat_audit_log message.received row.
    msg_audit = await _append_chat_audit(
        account_id=account_id, chat_id=f"enhance-{export_id}",
        action="message.received", payload={
            "user_message_id": None,
            "reply_id": None,
            "model_id": "claude-sonnet-4-5", "mode": "enhance",
            "latency_ms": p1_ms + p2_ms,
            "char_len_reply": len(file_text),
            "channel": "enhance",
            "turn_class": "strategic_deliverable",
            "classifier_source": "forced",
            "classifier_latency_ms": 0,
            "four_check_surfaced": {"label": None, "ran": True},
            "refusal_reason": None,
            "two_pass": {
                "pass_1": pass_1[:24000],
                "pass_2": json.dumps(content_dict)[:24000],
                "pass_1_visible": False,
                "pass_1_present": True,
                "pass_1_ms": p1_ms,
                "pass_2_ms": p2_ms,
            },
            "voice_violation":
                ({"banned_word": post_hit, "retry_outcome": "ship_with_violation",
                  "before_text": "(post-render scrape)", "after_text": None}
                 if post_hit else None),
            "export_kind": f"enhance_{kind}",
            "export_artefact_id": export_id,
            "pre_render_banned_word": pre_hit,
            "sensitivity_band": sensitivity_band,
            "sensitivity_score": sensitivity_score,
            "source_filename": source_filename,
            "source_label": source_label,
        }, request=None,
    )

    # ── Continue-in-chat
    cont_chat_id, cont_doc_id = await _create_continue_chat(
        account_id=account_id, context_id=context_id, kind=kind,
        source="enhance", export_id=export_id,
        file_name=fname, file_path=str(fpath), output_format=output_format,
        extracted_text=file_text, sensitivity_band=sensitivity_band,
    )

    completed_at = iso(now())
    await db.work_studio_exports.update_one(
        {"id": export_id},
        {"$set": {
            "status": "complete",
            "file_name": fname,
            "file_path": str(fpath),
            "sha256": sha,
            "byte_len": len(data),
            "sensitivity_band": sensitivity_band,
            "sensitivity_score": sensitivity_score,
            "sensitivity_reasons": sensitivity_reasons,
            "completed_at": completed_at,
            "chat_audit_id": msg_audit["id"],
            "pass_1_ms": p1_ms, "pass_2_ms": p2_ms,
            # Phase B.3 — provider/fallback signal, parallel to export.
            "llm_pass1": llm_meta.get("llm_pass1"),
            "llm_pass2": llm_meta.get("llm_pass2"),
            "voice_violation": post_hit,
            "continue_chat_id": cont_chat_id,
            "continue_doc_id": cont_doc_id,
        }},
    )
    try:
        await write_audit(
            account_id=account_id, context_id=context_id,
            action="work_studio.enhance.completed", target_id=export_id,
            metadata={
                "export_id": export_id, "kind": kind,
                "output_format": output_format,
                "sensitivity_band": sensitivity_band,
                "sha256": sha,
                "continue_chat_id": cont_chat_id,
                "continue_doc_id": cont_doc_id,
            },
        )
    except Exception:
        pass


# =============================================================================
# Endpoints
# =============================================================================
@router.post("/contexts/{context_id}/work-studio/export/{kind}")
async def start_export(
    context_id: str,
    kind: str,
    body: ExportRequestIn,
    background: BackgroundTasks,
    request: Request,
    ctx: Dict[str, Any] = Depends(require_context_membership()),
):
    """Phase C.2 — start a Work Studio export. Returns 202 immediately
    with the export_id. The actual LLM + render run in the background."""
    if kind not in _VALID_KINDS:
        raise HTTPException(status_code=400, detail=f"Unknown kind. Allowed: {', '.join(_VALID_KINDS)}.")
    output_format = _resolve_format(kind, body.output_format)

    export_id = str(uuid.uuid4())
    created_at = iso(now())
    row = {
        "id": export_id,
        "context_id": context_id,
        "account_id": ctx["account"]["id"],
        "kind": kind,
        "output_format": output_format,
        "status": "running",
        "description_chars": len(body.description),
        "objective_chars": len(body.objective),
        "scope_chars": len(body.scope),
        "created_at": created_at,
        "completed_at": None,
        "file_name": None,
        "file_path": None,
        "sha256": None,
        "sensitivity_band": None,
        "error": None,
        "refusal_text": None,
        "chat_audit_id": None,
        # Track A Phase 5 (2026-06-04) — additive fields for the
        # fig-53 card row 2 (sources count · contributors count ·
        # Akki Generated badge). source_count for a fresh export is
        # 0 (sources are not bound until a compilation chains in;
        # the /compilations/{id}/start path overrides this on insert).
        "source_count": 0,
        "contributor_count": 1,
        "akki_generated": True,
    }
    await db.work_studio_exports.insert_one(row)

    try:
        await write_audit(
            account_id=ctx["account"]["id"], context_id=context_id,
            action="work_studio.export.requested", target_id=export_id,
            metadata={"export_id": export_id, "kind": kind,
                      "output_format": output_format},
        )
    except Exception:
        pass

    # Schedule the worker. The wrapper detaches from the request scope.
    #
    # Chunk 3 (2026-05-13) — same fix as the enhance-runner above.
    # The pre-Chunk-3 code wrote the literal `"worker_crash"` here too,
    # hiding the real exception. Use a typed-error message so the UI
    # has something to show beyond an opaque token.
    async def _runner():
        try:
            await _run_export(
                export_id=export_id, account_id=ctx["account"]["id"],
                context_id=context_id, kind=kind,
                output_format=output_format, body=body,
            )
        except Exception as exc:
            logger.exception("export worker crashed (export_id=%s, kind=%s)", export_id, kind)
            err_class = type(exc).__name__
            err_msg = str(exc).replace("\n", " ").strip()[:300] or "(no message)"
            await db.work_studio_exports.update_one(
                {"id": export_id},
                {"$set": {
                    "status": "failed",
                    "error": f"{err_class}: {err_msg}",
                    "completed_at": iso(now()),
                }},
            )
    background.add_task(_runner)

    return {
        "export_id": export_id,
        "status": "running",
        "kind": kind,
        "output_format": output_format,
        "created_at": created_at,
    }


@router.get("/work_studio/artefacts/{kind}/{artefact_id}/synisense-breakdown")
async def get_studio_synisense_breakdown(
    kind: str,
    artefact_id: str,
    current: Dict[str, Any] = Depends(get_current_account),
) -> Dict[str, Any]:
    """STUDIO sprint (2026-05-12) — per-artefact Synisense breakdown.

    Returns the total identifier count + three-layer breakdown grouped
    by Work Studio surface (briefing / deck / report / brief / minutes /
    enhance). Powers the inline badge on the artefact detail drawer and
    the redaction-count line stamped on exports.

    Filter is on `artefact_id` (post-sprint sessions) with a fallback to
    account_id + surface family when no artefact_id is present (legacy
    artefacts created before the threading)."""
    if kind not in {"briefing", "deck", "report", "brief", "minutes", "committee_pack"}:
        raise HTTPException(status_code=400, detail=f"Unknown artefact kind: {kind}")
    # Map the artefact kind → the set of surface keys we aggregate over.
    # Every kind also includes its corresponding `enhance` surface so
    # an "Enhance pass" counts toward the same artefact.
    surface_family = {
        "briefing": ["briefing", "enhance"],
        "deck":     ["deck", "enhance"],
        "report":   ["report", "enhance"],
        "brief":    ["brief", "enhance"],
        "minutes":  ["minutes", "enhance"],
        # QA-2026-05-16-043 — committee_pack is a first-class artefact
        # kind alongside deck/report/minutes.
        "committee_pack": ["committee_pack", "enhance"],
    }[kind]

    pipeline = [
        {
            "$match": {
                "artefact_id": artefact_id,
                "surface": {"$in": surface_family},
            }
        },
        {
            "$group": {
                "_id": "$surface",
                "runs": {"$sum": 1},
                "identifiers": {"$sum": {"$size": {"$ifNull": ["$spans", []]}}},
                "layers": {"$push": "$stats.layer_won"},
            }
        },
    ]
    rows = await db.synisense_runs.aggregate(pipeline).to_list(length=32)

    per_surface: List[Dict[str, Any]] = []
    total_identifiers = 0
    total_calls = 0
    for row in rows:
        layers = row.get("layers") or []
        breakdown = {"regex": 0, "presidio": 0, "llm": 0}
        for lw in layers:
            if not lw:
                continue
            if lw == "regex":
                breakdown["regex"] += 1
            elif lw == "presidio":
                breakdown["presidio"] += 1
            elif lw in ("llm", "llm_fallback"):
                breakdown["llm"] += 1
        ident = int(row.get("identifiers") or 0)
        calls = int(row.get("runs") or 0)
        total_identifiers += ident
        total_calls += calls
        per_surface.append({
            "surface": row.get("_id"),
            "identifiers_count": ident,
            "model_calls": calls,
            "layers": breakdown,
        })
    per_surface.sort(key=lambda r: r["surface"])
    storyline = (
        f"This artefact passed through three layers of redaction. "
        f"{total_identifiers} identifier{'s were' if total_identifiers != 1 else ' was'} "
        f"masked deterministically across {total_calls} model call{'s' if total_calls != 1 else ''}. "
        f"Nothing left your tenant."
    )
    return {
        "artefact_id": artefact_id,
        "artefact_kind": kind,
        "per_surface": per_surface,
        "total_identifiers_count": total_identifiers,
        "model_calls_total": total_calls,
        "storyline": storyline,
    }


@router.get("/contexts/{context_id}/work-studio/exports/{export_id}/stream")
async def stream_export_progress(
    context_id: str,
    export_id: str,
    request: Request,
    ctx: Dict[str, Any] = Depends(require_context_membership()),
):
    """Phase L.a (2026-05-27) — SSE observer for an active export.

    The `POST /export/{kind}` endpoint kicks off a BackgroundTasks
    worker (`_run_export`) and returns 202 immediately. This endpoint
    OBSERVES that worker by polling `db.work_studio_exports.{phase_index}`
    at 500ms cadence and emits SSE `phase` events when the marker
    advances. When the row flips to `status=complete`, the final
    `complete` event carries the same payload the legacy
    `GET /exports/{eid}` endpoint returns. On `status=failed`, an
    `error` event is emitted. The frontend `StreamingLogScene` reads
    these events via `useStreamingProgress`.

    Auth: same membership check as the polling endpoint. Same
    cross-account narrowing.

    Why "observer" pattern (not synchronous SSE):
      The legacy POST→202+poll architecture is preserved so existing
      clients (Cycle.jsx + ExportModal polling code paths) keep
      working unchanged. The stream endpoint is purely additive.
    """
    row = await db.work_studio_exports.find_one(
        {"id": export_id, "context_id": context_id}, {"_id": 0},
    )
    if not row:
        raise HTTPException(status_code=404, detail="Export not found.")
    if row["account_id"] != ctx["account"]["id"]:
        raise HTTPException(status_code=403, detail="Not the export owner.")

    from services.streaming.sse import sse_headers
    from services.streaming.progress import PhaseEmitter
    from fastapi.responses import StreamingResponse

    async def _gen():
        emitter = PhaseEmitter(request, surface="work-studio-compile")
        last_phase_emitted = -1
        max_iterations = 600  # 600 * 500ms = 5 minutes hard cap
        async with emitter.start() as e:
            yield e.script_event()

            # Reset internal index; we drive advance() to match the
            # worker's marker.
            iter_count = 0
            while iter_count < max_iterations:
                iter_count += 1
                if e.cancelled or await emitter._stream.is_disconnected():
                    return

                cur = await db.work_studio_exports.find_one(
                    {"id": export_id, "context_id": context_id},
                    {"_id": 0, "phase_index": 1, "phase_label": 1, "status": 1,
                     "error": 1, "refusal_text": 1, "file_name": 1,
                     "sensitivity_band": 1, "completed_at": 1,
                     "kind": 1, "output_format": 1, "created_at": 1,
                     "chat_audit_id": 1, "continue_chat_id": 1,
                     "continue_doc_id": 1, "source": 1},
                )
                if not cur:
                    yield await emitter.error("Export row disappeared.", code="not_found")
                    return

                worker_phase = cur.get("phase_index", -1)
                if isinstance(worker_phase, int) and worker_phase > last_phase_emitted:
                    # Catch up phase by phase so the frontend renders
                    # each intermediate line (in case the worker
                    # jumped 2 phases inside one 500ms window).
                    while last_phase_emitted < worker_phase:
                        chunk = await emitter.advance()
                        if chunk is None:
                            return
                        yield chunk
                        last_phase_emitted += 1

                status_val = cur.get("status")
                if status_val == "complete":
                    # Mint download token (same path as get_export_status).
                    token = secrets.token_urlsafe(32)
                    await db.work_studio_export_tokens.insert_one({
                        "token": token, "export_id": export_id,
                        "expires_at": (
                            datetime.now(timezone.utc) + timedelta(seconds=_DOWNLOAD_TTL_SECONDS)
                        ).isoformat(),
                        "used": False,
                        "issued_at": iso(now()),
                        "account_id": ctx["account"]["id"],
                    })
                    yield await emitter.complete({
                        "export_id": export_id,
                        "status": "complete",
                        "kind": cur.get("kind"),
                        "output_format": cur.get("output_format"),
                        "file_name": cur.get("file_name"),
                        "sensitivity_band": cur.get("sensitivity_band"),
                        "completed_at": cur.get("completed_at"),
                        "download_token": token,
                        "chat_audit_id": cur.get("chat_audit_id"),
                        "continue_chat_id": cur.get("continue_chat_id"),
                        "continue_doc_id": cur.get("continue_doc_id"),
                        "source": cur.get("source") or "export",
                    })
                    return
                if status_val == "failed":
                    yield await emitter.error(
                        cur.get("error") or "Composition did not complete.",
                        code="export_failed",
                    )
                    return

                # Heartbeat every ~15s (30 iterations * 500ms) to keep
                # the connection warm during long LLM passes.
                if iter_count % 30 == 0:
                    hb = await emitter.heartbeat()
                    if hb is None:
                        return
                    yield hb

                await asyncio.sleep(0.5)

            # 5-minute observer cap — escalate as a timeout. The worker
            # keeps running; the user can poll the legacy GET endpoint.
            yield await emitter.error(
                "Composition is taking longer than expected. Try the polling fallback.",
                code="stream_timeout",
            )

    return StreamingResponse(_gen(), headers=sse_headers())


@router.get("/contexts/{context_id}/work-studio/exports/{export_id}")
async def get_export_status(
    context_id: str,
    export_id: str,
    ctx: Dict[str, Any] = Depends(require_context_membership()),
):
    row = await db.work_studio_exports.find_one(
        {"id": export_id, "context_id": context_id}, {"_id": 0},
    )
    if not row:
        raise HTTPException(status_code=404, detail="Export not found.")
    if row["account_id"] != ctx["account"]["id"]:
        # Membership covers context-level access; we still narrow to
        # the requesting account so cross-account peeks are blocked.
        raise HTTPException(status_code=403, detail="Not the export owner.")

    download_token: Optional[str] = None
    if row["status"] == "complete":
        # Mint a fresh single-use token.
        token = secrets.token_urlsafe(32)
        await db.work_studio_export_tokens.insert_one({
            "token": token, "export_id": export_id,
            "expires_at": (datetime.now(timezone.utc) + timedelta(seconds=_DOWNLOAD_TTL_SECONDS)).isoformat(),
            "used": False,
            "issued_at": iso(now()),
            "account_id": ctx["account"]["id"],
        })
        download_token = token

    return {
        "export_id": row["id"],
        "status": row["status"],
        "kind": row["kind"],
        "output_format": row["output_format"],
        "file_name": row.get("file_name"),
        "sensitivity_band": row.get("sensitivity_band"),
        "error": row.get("error"),
        "refusal_text": row.get("refusal_text"),
        "created_at": row["created_at"],
        "completed_at": row.get("completed_at"),
        "download_token": download_token,
        "chat_audit_id": row.get("chat_audit_id"),
        "continue_chat_id": row.get("continue_chat_id"),
        "continue_doc_id": row.get("continue_doc_id"),
        "source": row.get("source") or "export",
    }


@router.get("/contexts/{context_id}/work-studio/exports/{export_id}/download")
async def download_export(
    context_id: str,
    export_id: str,
    token: str,
    inline: bool = False,
    ctx: Dict[str, Any] = Depends(require_context_membership()),
):
    row = await db.work_studio_exports.find_one(
        {"id": export_id, "context_id": context_id}, {"_id": 0},
    )
    if not row:
        raise HTTPException(status_code=404, detail="Export not found.")
    if row["account_id"] != ctx["account"]["id"]:
        raise HTTPException(status_code=403, detail="Not the export owner.")
    if row["status"] != "complete":
        raise HTTPException(status_code=404, detail="Export not ready.")

    tok = await db.work_studio_export_tokens.find_one(
        {"token": token, "export_id": export_id, "used": False},
        {"_id": 0},
    )
    if not tok:
        raise HTTPException(status_code=403, detail="Invalid download token.")
    if tok.get("account_id") != ctx["account"]["id"]:
        raise HTTPException(status_code=403, detail="Token does not match account.")
    try:
        exp = datetime.fromisoformat((tok.get("expires_at") or "").replace("Z", "+00:00"))
        if exp < datetime.now(timezone.utc):
            raise HTTPException(status_code=403, detail="Token expired.")
    except HTTPException:
        raise
    except Exception:
        raise HTTPException(status_code=403, detail="Bad token shape.")

    # Mark used (best-effort single-use).
    await db.work_studio_export_tokens.update_one(
        {"token": token},
        {"$set": {"used": True, "used_at": iso(now())}},
    )

    fpath = Path(row["file_path"]) if row.get("file_path") else None
    if not fpath or not fpath.exists():
        raise HTTPException(status_code=404, detail="File missing on disk.")
    data = fpath.read_bytes()

    fmt = row["output_format"]
    media_types = {
        "docx": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        "pptx": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
        "pdf":  "application/pdf",
    }
    # Track A Phase 5 (2026-06-04, Tightening 4) — `?inline=true` makes
    # the response renderable inside an `<iframe>` for the Document
    # Review Side Drawer's PDF surface. The default (False) keeps the
    # legacy `attachment` disposition for the Download button. The
    # media-type stays `application/pdf` regardless of disposition.
    disposition = "inline" if inline else "attachment"
    return Response(
        content=data,
        media_type=media_types.get(fmt, "application/octet-stream"),
        headers={
            "Content-Disposition": f'{disposition}; filename="{row["file_name"]}"',
            "X-AKKI-Sensitivity-Band": row.get("sensitivity_band") or "INTERNAL",
            "Cache-Control": "private, no-store",
        },
    )



# =============================================================================
# C.3 — POST /api/contexts/{cid}/work-studio/enhance/{kind}
# Multipart form-data body: instructions (Form, required),
# output_format (Form, default "auto"), file (UploadFile, optional),
# source_artefact_id (Form, optional). One of file/source_artefact_id
# must be provided.
# =============================================================================
@router.post("/contexts/{context_id}/work-studio/enhance/{kind}")
async def start_enhance(
    context_id: str,
    kind: str,
    background: BackgroundTasks,
    request: Request,
    instructions: str = Form(...),
    output_format: str = Form("auto"),
    source_artefact_id: Optional[str] = Form(None),
    file: Optional[UploadFile] = File(None),
    ctx: Dict[str, Any] = Depends(require_context_membership()),
):
    """Phase C.3 — start a Work Studio enhance run. Same persistence
    shape as export (rows in db.work_studio_exports with source='enhance').
    """
    if kind not in _ENHANCE_KINDS:
        raise HTTPException(
            status_code=400,
            detail=f"Unknown enhance kind. Allowed: {', '.join(_ENHANCE_KINDS)}.",
        )
    if not instructions or not instructions.strip():
        raise HTTPException(status_code=400, detail="instructions is required.")
    fmt = _resolve_format(kind, output_format if output_format in ("docx","pptx","pdf","auto") else "auto")

    # Resolve source bytes BEFORE persisting the row so we surface
    # ClamAV / 4xx errors as HTTP responses (not as failed-row state).
    source_bytes, source_filename, source_label = await _resolve_enhance_source(
        context_id=context_id, account_id=ctx["account"]["id"], kind=kind,
        file=file, source_artefact_id=source_artefact_id,
    )

    export_id = str(uuid.uuid4())
    created_at = iso(now())
    row = {
        "id": export_id,
        "context_id": context_id,
        "account_id": ctx["account"]["id"],
        "kind": kind,
        "output_format": fmt,
        "status": "running",
        "source": "enhance",
        "instructions_chars": len(instructions),
        "source_label": source_label,
        "source_filename": source_filename,
        "source_artefact_id": source_artefact_id,
        "created_at": created_at,
        "completed_at": None,
        "file_name": None,
        "file_path": None,
        "sha256": None,
        "sensitivity_band": None,
        "error": None,
        "refusal_text": None,
        "chat_audit_id": None,
        # Track A Phase 5 (2026-06-04) — fig-53 card row 2 fields.
        # Enhance always has exactly one source (the file or the
        # source_artefact_id) and is solo-authored.
        "source_count": 1,
        "contributor_count": 1,
        "akki_generated": True,
    }
    await db.work_studio_exports.insert_one(row)

    try:
        await write_audit(
            account_id=ctx["account"]["id"], context_id=context_id,
            action="work_studio.enhance.requested", target_id=export_id,
            metadata={"export_id": export_id, "kind": kind,
                      "output_format": fmt, "source_label": source_label,
                      "source_filename": source_filename},
        )
    except Exception:
        pass

    # Append a hash-chained `enhance.requested` row to chat_audit_log.
    try:
        await _append_chat_audit(
            account_id=ctx["account"]["id"],
            chat_id=f"enhance-{export_id}", action="enhance.requested",
            payload={
                "export_kind": f"enhance_{kind}",
                "export_artefact_id": export_id,
                "source_label": source_label,
                "source_filename": source_filename,
                "instructions_preview": instructions[:200],
                "output_format": fmt,
                "channel": "enhance",
                "deterministic": True,
            }, request=None,
        )
    except Exception:
        logger.warning("enhance.requested audit-append failed (non-fatal)")

    # ── PRE-LLM deterministic thin-input refusal on the instructions
    # (closes the C.2 carry-over: refusal must fire BEFORE any LLM call
    # so it is resilient to upstream proxy errors).
    thin_shape = _is_thin_enhance_shape(instructions)
    if thin_shape is not None:
        await _emit_thin_refusal(
            export_id=export_id, account_id=ctx["account"]["id"],
            context_id=context_id, kind=kind, source="enhance",
            detection={**thin_shape, "stage": "pre_llm_enhance"},
        )
        return {
            "export_id": export_id,
            "status": "failed",
            "kind": kind,
            "output_format": fmt,
            "source": "enhance",
            "created_at": created_at,
            "error": "thin_input",
        }

    # Schedule the worker.
    #
    # Chunk 3 (2026-05-13, WS-R06/R12/R15) — pre-Chunk-3 this catch-all
    # wrote the literal string `"worker_crash"` as the row's error,
    # hiding the actual exception type and message from the user (and
    # from anyone reading the row without backend log access). The real
    # exception was logged but never persisted. Now the error string
    # carries the actual exception type + message so the UI can show
    # something actionable ("Refine failed: TimeoutError — LLM provider
    # returned 503 after 60s") instead of an opaque "worker_crash".
    # logger.exception still captures the full traceback for ops.
    async def _runner():
        try:
            await _run_enhance(
                export_id=export_id, account_id=ctx["account"]["id"],
                context_id=context_id, kind=kind, output_format=fmt,
                instructions=instructions, source_data=source_bytes,
                source_filename=source_filename, source_label=source_label,
            )
        except Exception as exc:
            logger.exception("enhance worker crashed (export_id=%s, kind=%s)", export_id, kind)
            # Sanitise: cap at 300 chars to avoid bloating the DB row;
            # strip newlines so the message renders cleanly inline; keep
            # the exception class name as the leading token so the UI
            # can decide whether to offer Retry or not.
            err_class = type(exc).__name__
            err_msg = str(exc).replace("\n", " ").strip()[:300] or "(no message)"
            await db.work_studio_exports.update_one(
                {"id": export_id},
                {"$set": {
                    "status": "failed",
                    "error": f"{err_class}: {err_msg}",
                    "completed_at": iso(now()),
                }},
            )
    background.add_task(_runner)

    return {
        "export_id": export_id,
        "status": "running",
        "kind": kind,
        "output_format": fmt,
        "source": "enhance",
        "created_at": created_at,
    }


# ─────────────────────────────────────────────────────────────────────
# Track A Phase 5 (2026-06-04) — W5 Drafting Drawer backing endpoint
#
# Powers the [Save] CTA on the Drafting Side Drawer (W5 + W6 blank
# routes). Persists a `work_studio_exports` row with lifecycle_state=
# "draft", source="draft", akki_generated=true. Subsequent saves with
# the same `draft_session_id` update the existing row in place; this
# avoids the race where two near-simultaneous autosave debounces fire
# on first-open and create two rows (Tightening 5).
#
# Note: this endpoint does NOT render a file — the draft body lives in
# `structured_content`. The Render endpoint at work_studio_render.py
# already handles structured_content → HTML conversion for the drawer.
# File materialisation (DOCX/PPTX/PDF) happens on Enhance → at which
# point the row flips to lifecycle_state="in_review" with source=
# "enhance" via the existing enhance pipeline.
# ─────────────────────────────────────────────────────────────────────
class _SaveDraftIn(BaseModel):
    draft_session_id: str = Field(..., min_length=8, max_length=64)
    title: str = Field(default="Untitled draft", max_length=300)
    structured_content: Dict[str, Any] = Field(default_factory=dict)
    kind: Literal["report", "deck", "minutes", "committee_pack", "briefing", "board_pack"] = "report"
    document_id: Optional[str] = None  # link back to the source document (W5 path)


@router.post("/contexts/{context_id}/work-studio/documents/save-draft")
async def save_draft(
    context_id: str,
    body: _SaveDraftIn,
    ctx: Dict[str, Any] = Depends(require_context_membership()),
):
    kind = body.kind
    # Format follows the kind: report/minutes/committee_pack/briefing →
    # docx; deck/board_pack → pptx.
    fmt = "pptx" if kind in ("deck", "board_pack") else "docx"

    # Idempotency (Tightening 5) — keyed on
    # `(context_id, account_id, draft_session_id)`. Two near-simultaneous
    # POSTs with the same draft_session_id collapse to ONE row.
    existing = await db.work_studio_exports.find_one(
        {
            "context_id": context_id,
            "account_id": ctx["account"]["id"],
            "draft_session_id": body.draft_session_id,
        },
        {"_id": 0},
    )
    now_iso = iso(now())
    if existing:
        await db.work_studio_exports.update_one(
            {"id": existing["id"]},
            {"$set": {
                "document_title":     body.title.strip() or "Untitled draft",
                "structured_content": body.structured_content,
                "updated_at":         now_iso,
            }},
        )
        refreshed = await db.work_studio_exports.find_one({"id": existing["id"]}, {"_id": 0})
        return {
            "export_id":          refreshed["id"],
            "draft_session_id":   body.draft_session_id,
            "lifecycle_state":    refreshed.get("lifecycle_state", "draft"),
            "updated_at":         now_iso,
            "idempotent_update":  True,
        }

    export_id = str(uuid.uuid4())
    row = {
        "id": export_id,
        "context_id": context_id,
        "account_id": ctx["account"]["id"],
        "kind": kind,
        "output_format": fmt,
        "status": "complete",          # draft is its own terminal state
        "lifecycle_state": "draft",
        "source": "draft",
        "draft_session_id": body.draft_session_id,
        "source_document_id": body.document_id,
        "document_title": body.title.strip() or "Untitled draft",
        "structured_content": body.structured_content,
        "file_name": None,
        "file_path": None,
        "sha256": None,
        "sensitivity_band": None,
        "error": None,
        "refusal_text": None,
        "chat_audit_id": None,
        "created_at": now_iso,
        "updated_at": now_iso,
        "completed_at": now_iso,
        # Phase 5 additive fields.
        "source_count": 1 if body.document_id else 0,
        "contributor_count": 1,
        "akki_generated": True,
    }
    await db.work_studio_exports.insert_one(row)
    try:
        await write_audit(
            account_id=ctx["account"]["id"], context_id=context_id,
            action="work_studio.draft.saved", target_id=export_id,
            metadata={"draft_session_id": body.draft_session_id, "kind": kind},
        )
    except Exception:  # noqa: BLE001
        # Swallow contract: audit-log failures must not block the save.
        pass

    return {
        "export_id":         export_id,
        "draft_session_id":  body.draft_session_id,
        "lifecycle_state":   "draft",
        "updated_at":        now_iso,
        "idempotent_update": False,
    }
