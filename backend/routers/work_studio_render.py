"""T4.1 — Work Studio Compiled Document on-the-fly render (G6 ratified).

Endpoint:
  GET /api/contexts/{cid}/work-studio/documents/{aid}/render?format={docx|pdf|pptx}

Renders the artefact's `structured_content` to the requested format
on demand and streams the binary back with the right Content-Type
header. This is the surface that backs the three G6 download buttons
on the Compiled Document toolbar (spec §4.C → W3 + §6 → G6 ratified).

Notes
-----
* The render is server-produced and synchronous — no background worker,
  no token round-trip. The artefact's already-derived structured_content
  carries no PII that hasn't already passed through Shield during the
  original compile (Phase C.2 + Phase C.3), so this endpoint does not
  re-Shield: it just serialises.
* If `structured_content` is missing or empty, returns 409 so the
  frontend can surface a clear "this artefact hasn't been compiled yet"
  state instead of returning an empty / malformed file.
* DOCX uses `python-docx`. PPTX uses `python-pptx`. PDF uses ReportLab
  (already installed; lighter than weasyprint for a basic letterhead-
  style brief).
"""
from __future__ import annotations

import io
import logging
from typing import Any, Dict, List, Optional

from fastapi import APIRouter, Depends, HTTPException
from fastapi.responses import Response

from core import (
    db, write_audit, require_context_membership,
)

logger = logging.getLogger("akki.work_studio_render")
router = APIRouter(prefix="/api")


_MEDIA_TYPES = {
    "docx": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
    "pptx": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
    "pdf":  "application/pdf",
}


# ── Sectional helpers ─────────────────────────────────────────────────
def _safe_sections(structured_content: Optional[Dict[str, Any]]) -> List[Dict[str, Any]]:
    if not structured_content or not isinstance(structured_content, dict):
        return []
    sections = structured_content.get("sections") or []
    if not isinstance(sections, list):
        return []
    cleaned = []
    for s in sections:
        if not isinstance(s, dict):
            continue
        heading = s.get("heading") or ""
        paragraphs = s.get("paragraphs") or []
        if not isinstance(paragraphs, list):
            paragraphs = []
        paragraphs = [p for p in paragraphs if isinstance(p, str)]
        cleaned.append({"heading": str(heading), "paragraphs": paragraphs})
    return cleaned


# ── DOCX renderer ─────────────────────────────────────────────────────
def _render_docx(title: str, sections: List[Dict[str, Any]]) -> bytes:
    from docx import Document
    from docx.shared import Pt

    doc = Document()
    style = doc.styles["Normal"]
    style.font.name = "Calibri"
    style.font.size = Pt(11)
    doc.add_heading(title or "Document", level=0)
    for s in sections:
        if s.get("heading"):
            doc.add_heading(s["heading"], level=1)
        for p in s.get("paragraphs") or []:
            doc.add_paragraph(p)
    buf = io.BytesIO()
    doc.save(buf)
    return buf.getvalue()


# ── PPTX renderer ─────────────────────────────────────────────────────
def _render_pptx(title: str, sections: List[Dict[str, Any]]) -> bytes:
    from pptx import Presentation
    from pptx.util import Inches, Pt

    prs = Presentation()
    # Title slide
    title_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_layout)
    slide.shapes.title.text = title or "Document"
    subtitle = slide.placeholders[1] if len(slide.placeholders) > 1 else None
    if subtitle:
        subtitle.text = "Akki Work Studio"

    # One slide per section
    bullet_layout = prs.slide_layouts[1]
    for s in sections:
        slide = prs.slides.add_slide(bullet_layout)
        slide.shapes.title.text = s.get("heading") or " "
        body_ph = slide.placeholders[1] if len(slide.placeholders) > 1 else None
        if body_ph and body_ph.has_text_frame:
            tf = body_ph.text_frame
            tf.clear()
            paragraphs = s.get("paragraphs") or []
            if not paragraphs:
                # Keep the placeholder non-empty to avoid orphan layout.
                tf.text = " "
            else:
                tf.text = paragraphs[0]
                for p in paragraphs[1:]:
                    new_p = tf.add_paragraph()
                    new_p.text = p
                    new_p.font.size = Pt(14)

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


# ── PDF renderer ──────────────────────────────────────────────────────
def _render_pdf(title: str, sections: List[Dict[str, Any]]) -> bytes:
    from reportlab.lib.pagesizes import A4
    from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
    from reportlab.lib.units import cm
    from reportlab.platypus import (
        SimpleDocTemplate, Paragraph, Spacer, PageBreak,
    )

    buf = io.BytesIO()
    doc = SimpleDocTemplate(
        buf, pagesize=A4,
        topMargin=2.0 * cm, bottomMargin=2.0 * cm,
        leftMargin=2.0 * cm, rightMargin=2.0 * cm,
        title=title or "Document",
    )
    styles = getSampleStyleSheet()
    title_style = ParagraphStyle(
        "AkkiTitle",
        parent=styles["Title"],
        spaceAfter=18,
    )
    heading_style = ParagraphStyle(
        "AkkiH1",
        parent=styles["Heading1"],
        spaceBefore=14, spaceAfter=8,
    )
    body_style = ParagraphStyle(
        "AkkiBody",
        parent=styles["BodyText"],
        leading=14,
        spaceAfter=8,
    )
    flow: List[Any] = []
    flow.append(Paragraph(title or "Document", title_style))
    for s in sections:
        if s.get("heading"):
            flow.append(Paragraph(s["heading"], heading_style))
        for p in s.get("paragraphs") or []:
            # ReportLab interprets `<` etc. as XML — escape the basics.
            safe = (
                p.replace("&", "&amp;")
                 .replace("<", "&lt;")
                 .replace(">", "&gt;")
            )
            flow.append(Paragraph(safe, body_style))
        flow.append(Spacer(1, 0.4 * cm))
    if not flow:
        flow.append(Paragraph(" ", body_style))
    doc.build(flow)
    return buf.getvalue()


# ── Endpoint ──────────────────────────────────────────────────────────
@router.get(
    "/contexts/{context_id}/work-studio/documents/{artefact_id}/render"
)
async def render_compiled_document(
    context_id: str,
    artefact_id: str,
    format: str,
    ctx: Dict[str, Any] = Depends(require_context_membership()),
):
    """Server-produced render in DOCX / PDF / PPTX (G6 ratified).

    Authorisation: the caller must be a member of `context_id`; the
    require_context_membership dependency enforces that. The artefact
    must belong to the same context.

    422 paths:
      - `format` not in {docx, pdf, pptx}

    404 paths:
      - artefact not found in this context

    409 paths:
      - artefact has no `structured_content` yet (e.g. fresh Draft
        created via T3.1 D5 that hasn't been compiled)
    """
    fmt = (format or "").strip().lower()
    if fmt not in _MEDIA_TYPES:
        raise HTTPException(
            status_code=422,
            detail=(
                f"Unsupported format `{format}`. "
                f"G6 ratified accepts: docx, pdf, pptx."
            ),
        )

    row = await db.work_studio_exports.find_one(
        {"id": artefact_id, "context_id": context_id},
        {"_id": 0},
    )
    if not row:
        raise HTTPException(status_code=404, detail="Artefact not found.")

    sections = _safe_sections(row.get("structured_content"))
    if not sections:
        raise HTTPException(
            status_code=409,
            detail=(
                "This artefact has no compiled content yet. "
                "Open it in Work Studio and run a compile or revision first."
            ),
        )
    title = row.get("title") or "Document"

    try:
        if fmt == "docx":
            data = _render_docx(title, sections)
            file_name = f"{_slug(title)}.docx"
        elif fmt == "pptx":
            data = _render_pptx(title, sections)
            file_name = f"{_slug(title)}.pptx"
        else:
            data = _render_pdf(title, sections)
            file_name = f"{_slug(title)}.pdf"
    except Exception:
        logger.exception("render_compiled_document: %s render failed", fmt)
        raise HTTPException(
            status_code=500,
            detail=f"Render failed for format `{fmt}`.",
        )

    try:
        await write_audit(
            context_id=context_id,
            account_id=ctx["account"]["id"],
            action="work_studio.compiled_document.rendered",
            resource_type=f"work_studio_artefact.{row.get('kind') or 'unknown'}",
            resource_id=artefact_id,
            metadata={"format": fmt, "size_bytes": len(data)},
        )
    except Exception:
        logger.exception("render_compiled_document: audit write failed (non-fatal)")

    return Response(
        content=data,
        media_type=_MEDIA_TYPES[fmt],
        headers={
            "Content-Disposition": f'attachment; filename="{file_name}"',
            "X-AKKI-Sensitivity-Band": row.get("sensitivity_band") or "INTERNAL",
            "Cache-Control": "private, no-store",
        },
    )


# ── helpers ───────────────────────────────────────────────────────────
def _slug(s: str) -> str:
    if not s:
        return "document"
    out = []
    for ch in s.strip():
        if ch.isalnum():
            out.append(ch)
        elif ch in (" ", "-", "_"):
            out.append("_")
    slug = "".join(out).strip("_")[:80]
    return slug or "document"
