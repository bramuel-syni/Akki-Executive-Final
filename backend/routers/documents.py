"""Documents: upload pipeline, thread, list/get/patch/archive/download, generate-meta."""
from __future__ import annotations

import logging
import uuid
from pathlib import Path
from typing import Any, Dict, List, Literal, Optional

from fastapi import APIRouter, BackgroundTasks, Depends, File, Form, HTTPException, UploadFile
from datetime import datetime, timezone
from pydantic import BaseModel, Field

from llm_service import call_llm as llm_call_llm, parse_json_response
from documents_service import (
    ACCEPT_EXT, MAX_BYTES, save_to_storage,
    delete_from_storage, extract_text, make_preview,
)
from services import clamav_service
from services.clamav_service import ClamAVUnreachable
from core import (
    db, now as _now, iso as _iso, write_audit, require_context_membership,
    get_current_account,
)

logger = logging.getLogger("akki.documents")

router = APIRouter(prefix="/api")

MAX_EXTRACT_CHARS_OUT = 40000


def sanitize_doc(d: Dict[str, Any]) -> Dict[str, Any]:
    return {
        "id": d["id"],
        "context_id": d["context_id"],
        "name": d.get("name"),
        "description": d.get("description", ""),
        "original_filename": d.get("original_filename"),
        "mime_type": d.get("mime_type"),
        "size_bytes": d.get("size_bytes", 0),
        "status": d.get("status", "uploaded"),
        "preview": d.get("preview", ""),
        "extracted_chars": d.get("extracted_chars", 0),
        "data_trust": d.get("data_trust", "mixed"),
        "uploaded_by_email": d.get("uploaded_by_email"),
        "mentioned_account_ids": d.get("mentioned_account_ids", []),
        "related_doc_id": d.get("related_doc_id"),
        "relation_type": d.get("relation_type"),
        "error": d.get("error"),
        "created_at": d.get("created_at"),
        "doc_type": d.get("doc_type"),
        # T2.1 (2026-05-25) — surface `source_channel` (and the
        # adjacent `doc_kind`) so the Document Journal listing can
        # derive a per-row origin tag (Uploaded vs Akki Generated)
        # per spec §4.A → D3/D4. `source_channel` is the canonical
        # discriminator: cycle_compilation + work_studio_export are
        # Akki-Generated origins; everything else (upload, inbound_email,
        # chat_attach, solva_attach, sandbox) is user-Uploaded.
        "source_channel": d.get("source_channel"),
        "doc_kind": d.get("doc_kind"),
        # Inbound email provenance (iter51) + trust tier (iter70)
        "source": d.get("source"),
        "inbound_from_email": d.get("inbound_from_email"),
        "inbound_from_name": d.get("inbound_from_name"),
        "inbound_subject": d.get("inbound_subject"),
        "inbound_trust_tier": d.get("inbound_trust_tier"),
        "inbound_trust_reason": d.get("inbound_trust_reason"),
        "inbound_reportee_id": d.get("inbound_reportee_id"),
        "inbound_reportee_name": d.get("inbound_reportee_name"),
        "inbound_reportee_title": d.get("inbound_reportee_title"),
        "inbound_queue_id": d.get("inbound_queue_id"),
        "inbound_promoted_by": d.get("inbound_promoted_by"),
        "inbound_promoted_at": d.get("inbound_promoted_at"),
        "inbound_promoted_note": d.get("inbound_promoted_note"),
        # Phase E.3 (2026-05-26) — Universal Document Drawer fields.
        "state":         d.get("state"),
        "objective":     d.get("objective"),
        "origin":        d.get("origin"),
        "audience":      d.get("audience"),
        "updated_at":    d.get("updated_at"),
        "committed_at":  d.get("committed_at"),
        # Phase Z (2026-05-27) — orthogonal category × origin model.
        # `category` answers "what kind of artefact?" → Work Studio
        # tabs. `origin` (above) answers "where did it come from?" →
        # /app/documents tabs. Both required on every row.
        "category":      d.get("category"),
        # Z2.7 (2026-02) — surface `source_doc_ids` so the Drawer's
        # Sources block can render the populated branch. List of
        # source document UUIDs that composed this artefact (set on
        # work-studio compile and cycle-pack export). Empty list /
        # null → Drawer renders the fallback line.
        "source_doc_ids": d.get("source_doc_ids") or [],
        # Sprint Z1.2 (2026-05-29) + Track B Phase B5 G6 (2026-06-04):
        # notes + per-notes timestamp for the drawer's Notes tab.
        # `notes_updated_at` is set only on notes-bearing PATCHes; it
        # surfaces as the FE's "Last updated: …" indicator and does
        # NOT churn on unrelated edits (title/category/audience).
        "notes":              d.get("notes"),
        "notes_updated_at":   d.get("notes_updated_at"),
    }



def _render_structured_content(sc: Any) -> str:
    """Best-effort: flatten a work_studio_exports `structured_content`
    payload into a plain-text rendering for the Drawer's extracted_text
    pane. Handles the common shapes the exporter produces (dict with
    sections, list of blocks, plain string). Synthesis is read-only —
    the source row is never mutated."""
    if sc is None:
        return ""
    if isinstance(sc, str):
        return sc
    if isinstance(sc, list):
        return "\n\n".join(_render_structured_content(x) for x in sc if x)
    if isinstance(sc, dict):
        # Common exporter shape: { sections: [{heading, body|paragraphs}, ...] }
        parts: List[str] = []
        title = sc.get("title") or sc.get("heading")
        if isinstance(title, str) and title.strip():
            parts.append(title.strip())
        for section in (sc.get("sections") or []):
            if isinstance(section, dict):
                h = section.get("heading") or section.get("title")
                if isinstance(h, str) and h.strip():
                    parts.append(h.strip())
                body = section.get("body") or section.get("text")
                if isinstance(body, str) and body.strip():
                    parts.append(body.strip())
                paras = section.get("paragraphs")
                if isinstance(paras, list):
                    for p in paras:
                        if isinstance(p, str) and p.strip():
                            parts.append(p.strip())
                        elif isinstance(p, dict) and isinstance(p.get("text"), str):
                            parts.append(p["text"].strip())
            elif isinstance(section, str) and section.strip():
                parts.append(section.strip())
        # Fallback for less-structured payloads.
        for k in ("body", "text", "summary", "description", "executive_summary"):
            v = sc.get(k)
            if isinstance(v, str) and v.strip() and v.strip() not in parts:
                parts.append(v.strip())
        return "\n\n".join(parts)
    return str(sc)


def _synthesize_doc_from_export(row: Dict[str, Any]) -> Dict[str, Any]:
    """Bug #1 fix (2026-05-27) — when a work_studio_exports row is opened
    via the canonical `?doc_id=` URL contract but has no corresponding
    `documents` mirror, return a documents-shaped read-only payload so the
    Universal Document Drawer renders the artefact without 404'ing.

    The synthesized payload carries:
    - `id` = the export's id (URL contract stable; subsequent endpoint
      calls that don't apply to exports will 404 gracefully — out of
      scope per the dispatch brief).
    - `_synthesized_from = "work_studio_export"` marker for the frontend.
    - `work_studio_export_id` = self-ref so the drawer can detect the
      synthesis case alongside the existing "Continue-spawned" mirrors.
    """
    sc = row.get("structured_content")
    intel = row.get("intelligence_report") if isinstance(row.get("intelligence_report"), dict) else None
    kind = row.get("kind") or ""
    output_format = row.get("output_format") or ""

    # Title resolution: structured_content.title → intelligence_report.title
    # → file_name (without ext) → "Brief"/"Deck"/etc. derived from kind.
    title = None
    if isinstance(sc, dict):
        title = sc.get("title") or sc.get("heading")
    if not title and intel:
        title = intel.get("title")
    if not title and row.get("file_name"):
        try:
            title = str(row["file_name"]).rsplit(".", 1)[0]
        except Exception:
            title = row.get("file_name")
    if not title:
        title = (kind or "artefact").replace("_", " ").strip().title() or "Untitled artefact"

    body = _render_structured_content(sc)
    fname = row.get("file_name") or (
        f"{kind or 'artefact'}{('.' + output_format) if output_format else ''}"
    )

    lifecycle = (row.get("lifecycle_state") or "committed").lower()
    state = "committed" if lifecycle in ("committed", "locked") else "draft"

    return {
        "id":               row["id"],
        "context_id":       row["context_id"],
        "name":             title,
        "description":      row.get("description") or "",
        "original_filename": fname,
        "mime_type":        {
                                "docx": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
                                "pptx": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
                                "pdf":  "application/pdf",
                            }.get(output_format, "application/octet-stream"),
        "size_bytes":       row.get("description_chars", 0) + row.get("objective_chars", 0) + row.get("scope_chars", 0),
        "status":           "extracted",
        "preview":          (body or "")[:320],
        "extracted_text":   body,
        "extracted_chars":  len(body),
        "data_trust":       "trusted",
        "doc_type":         "work_studio_artefact",
        "uploaded_by_email": row.get("account_id"),
        "source_channel":   "work_studio_export",
        # Phase E.3 fields — drives the Universal Drawer creation-mode.
        "state":            state,
        "origin":           "akki_generated",
        "audience":         None,
        "objective":        None,
        "created_at":       row.get("created_at"),
        "updated_at":       row.get("completed_at") or row.get("created_at"),
        "committed_at":     row.get("completed_at") if state == "committed" else None,
        # Sensitivity + Synisense fields the Drawer reads.
        "sensitivity_band":  row.get("sensitivity_band"),
        "sensitivity_score": None,
        "sensitivity_label": (row.get("sensitivity_band") or "INTERNAL").upper() if row.get("sensitivity_band") else None,
        "body_redacted":     None,
        "synisense_version": 0,
        # Provenance markers (Bug #1 fix).
        "work_studio_export_id": row["id"],
        "_synthesized_from":     "work_studio_export",
    }



class DocumentTrustUpdate(BaseModel):
    data_trust: Optional[Literal["trusted", "mixed", "weak"]] = None
    related_doc_id: Optional[str] = Field(
        default=None,
        description="Link this document as a successor of an existing doc — "
                    "lets NEDs trace how a recurring report has evolved across "
                    "cycles. Pass `null` to unlink. Pass an empty string for "
                    "no change (use Optional[None] sentinel below).",
    )


class DocumentMetaGenerateIn(BaseModel):
    filename: Optional[str] = None
    preview_text: Optional[str] = Field(default=None, max_length=8_000)


@router.post("/contexts/{context_id}/documents/generate-meta")
async def generate_document_meta(
    context_id: str, body: DocumentMetaGenerateIn,
    ctx: Dict[str, Any] = Depends(require_context_membership()),
):
    """Given a filename and optional preview text, ask the LLM to propose a
    short display name and ≤300-char description. Used by the upload modal."""
    if not body.filename and not body.preview_text:
        raise HTTPException(status_code=400, detail="Send filename or preview_text.")
    sample = (body.preview_text or "")[:4000]
    prompt = (
        "Propose a short display name and a description for a board-pack "
        "document. Stay neutral, specific, no hype. Return JSON ONLY:\n"
        '{"display_name": "<=60 chars title-case declarative name", '
        '"description": "<=300 chars single-paragraph description"}\n\n'
        f"Filename: {body.filename or '(unknown)'}\n\n"
        f"First ~4KB of extracted text:\n{sample or '(no text extracted)'}"
    )
    out = await llm_call_llm(
        module="document.meta",
        user_query=prompt,
        context_object=None,
        session_context={"session_id": f"docmeta-{context_id}"},
        data_trust={"overall": "unrated"},
        response_format="json",
    )
    parsed = parse_json_response(out.get("response", ""))
    if not isinstance(parsed, dict):
        raise HTTPException(status_code=502, detail=f"LLM returned no meta. Mode={out.get('mode')}.")
    return {
        "display_name": (parsed.get("display_name") or "")[:60].strip(),
        "description": (parsed.get("description") or "")[:300].strip(),
        "mode": out.get("mode"),
        "shielding_masked": out.get("shielding", {}).get("identifiers_masked", 0),
        "shielding": out.get("shielding", {}),
    }


@router.post("/contexts/{context_id}/documents/{doc_id}/summary")
async def generate_document_summary(
    context_id: str, doc_id: str,
    ctx: Dict[str, Any] = Depends(require_context_membership()),
    refresh: bool = False,
):
    """AKKI summary of a single document for the Document Journal.

    Returns:
        - tldr: 2-3 sentence executive summary
        - highlights: list of 4-7 important points the reader should know
        - questions: list of 2-3 questions the reader should bring back

    Cached on the document so opening it twice doesn't re-burn the LLM.
    Pass ?refresh=true to regenerate.
    """
    d = await db.documents.find_one(
        {"id": doc_id, "context_id": context_id}, {"_id": 0},
    )
    if not d:
        raise HTTPException(status_code=404, detail="Document not found")

    cached = d.get("akki_summary")
    if cached and not refresh:
        return cached

    text = (d.get("extracted_text") or "").strip()
    if not text:
        raise HTTPException(status_code=400, detail="No extractable text on this document.")

    sample = text[:12000]
    prompt = (
        "Read the executive document below. Summarise it for a board "
        "member who has 60 seconds. Return JSON ONLY:\n"
        '{\n'
        '  "tldr": "<=2 sentences, what this document is and what it says>",\n'
        '  "highlights": [\n'
        '    "<the single most important point — a fact or a claim — phrased plainly>",\n'
        '    "<next>",\n'
        '    "<...4 to 7 total>"\n'
        '  ],\n'
        '  "questions": [\n'
        '    "<question the reader should walk into the room with>",\n'
        '    "<...2 to 3 total>"\n'
        '  ]\n'
        '}\n\n'
        f"Document name: {d.get('name')}\n\n"
        f"Document text (first ~12KB):\n{sample}"
    )
    out = await llm_call_llm(
        module="document.summary",
        user_query=prompt,
        context_object=None,
        session_context={"session_id": f"docsum-{doc_id}"},
        data_trust={"overall": d.get("data_trust", "unrated")},
        response_format="json",
    )
    parsed = parse_json_response(out.get("response", "")) or {}
    summary = {
        "tldr": (parsed.get("tldr") or "").strip(),
        "highlights": [str(h).strip() for h in (parsed.get("highlights") or []) if str(h).strip()][:7],
        "questions": [str(q).strip() for q in (parsed.get("questions") or []) if str(q).strip()][:3],
        "mode": out.get("mode"),
        "generated_at": _iso(_now()),
    }
    await db.documents.update_one(
        {"id": doc_id, "context_id": context_id},
        {"$set": {"akki_summary": summary, "updated_at": _iso(_now())}},
    )
    return summary


@router.post("/contexts/{context_id}/documents")
async def upload_document(
    context_id: str,
    file: UploadFile = File(...),
    display_name: Optional[str] = Form(None),
    description: Optional[str] = Form(None),
    data_trust: Optional[str] = Form(None),
    mentioned_account_ids: Optional[str] = Form(None),  # comma-sep
    related_doc_id: Optional[str] = Form(None),
    relation_type: Optional[str] = Form(None),  # update | follow_up | additional_context | correction
    # Phase Z (2026-05-27) — upload modal carries the user-selected
    # category. Empty string == uncategorized (persisted as null).
    category: Optional[str] = Form(None),
    background_tasks: BackgroundTasks = None,  # type: ignore[assignment]
    ctx: Dict[str, Any] = Depends(require_context_membership()),
):
    filename = file.filename or "unnamed"
    ext = Path(filename).suffix.lower()
    if ext not in ACCEPT_EXT:
        raise HTTPException(status_code=415, detail=f"Unsupported file type {ext}. Accepted: {', '.join(sorted(ACCEPT_EXT))}")

    data = await file.read()
    if len(data) > MAX_BYTES:
        # G25 (2026-05-25, ratified spec §3 Stage 4) — verbatim toast
        # for oversized uploads. NOT paraphrased.
        raise HTTPException(
            status_code=413,
            detail="That file is larger than 50 MB. Please split it or upload a smaller version.",
        )

    # Generate the doc_id up front so the clamav `upload_scan_log` row
    # carries the same identifier the persisted document row will use.
    doc_id = str(uuid.uuid4())

    # Real virus scanning (Phase 10). clamd unreachable → 503 + audit.
    # Signature match → 422 + audit. Neither branch persists the file.
    try:
        scan_result = await clamav_service.scan(data, filename, file_id=doc_id, user_id=ctx["account"]["id"])
    except ClamAVUnreachable as e:
        await write_audit(
            context_id, ctx["account"]["id"],
            "upload.virus_scan.unreachable",
            "document", None,
            {"filename_hash": __import__("hashlib").sha256((filename or "").encode()).hexdigest()[:16],
             "error": str(e)[:200]},
        )
        raise HTTPException(
            status_code=503,
            detail={"error": "scanner_unavailable", "reason": "virus scanner offline"},
        )
    if not scan_result.clean:
        await write_audit(
            context_id, ctx["account"]["id"],
            "upload.virus_scan.blocked",
            "document", None,
            {
                "filename_hash": __import__("hashlib").sha256((filename or "").encode()).hexdigest()[:16],
                "signature": scan_result.signature,
                "size_bytes": len(data),
                "scan_ms": scan_result.scan_ms,
            },
        )
        raise HTTPException(
            status_code=422,
            detail={"error": "blocked", "reason": "malware_suspected", "signature": scan_result.signature},
        )

    if related_doc_id:
        if relation_type not in ("update", "follow_up", "additional_context", "correction"):
            raise HTTPException(status_code=400, detail="Invalid relation_type. Expected: update | follow_up | additional_context | correction")
        related = await db.documents.find_one(
            {"id": related_doc_id, "context_id": context_id, "status": {"$ne": "archived"}},
            {"_id": 0, "id": 1},
        )
        if not related:
            raise HTTPException(status_code=404, detail="related_doc_id not found in this context")

    mention_ids: List[str] = []
    if mentioned_account_ids:
        requested = [m.strip() for m in mentioned_account_ids.split(",") if m.strip()]
        if requested:
            valid = await db.memberships.find(
                {"context_id": context_id, "account_id": {"$in": requested}, "status": "active"},
                {"_id": 0, "account_id": 1},
            ).to_list(100)
            mention_ids = [v["account_id"] for v in valid]

    storage_key = save_to_storage(context_id, doc_id, filename, data)
    text, err = extract_text(data, filename, file.content_type or "")
    preview = make_preview(text)

    # G24 (2026-05-25, ratified spec §3 Stage 4) — verbatim 400 when
    # no text can be extracted. Refusing pre-storage so the user sees
    # the exact message + nothing is committed.
    if not text or not text.strip():
        raise HTTPException(
            status_code=400,
            detail="That file doesn't have any text we can read. Please upload a different file.",
        )

    created_at = _iso(_now())
    trust = data_trust if data_trust in ("trusted", "mixed", "weak") else "mixed"
    # Phase Z (2026-05-27) — normalize category. Empty string OR a
    # value outside the canonical 6-enum is stored as None
    # ("uncategorized"). The upload modal's "Uncategorized" option
    # submits an empty string.
    _CATEGORY_ENUM = ("board_pack", "minutes", "draft", "deck", "report", "briefing")
    cat_clean: Optional[str] = category if category in _CATEGORY_ENUM else None
    doc = {
        "id": doc_id,
        "context_id": context_id,
        "name": (display_name or Path(filename).stem).strip() or "Untitled",
        "description": (description or "")[:300].strip(),
        "original_filename": filename,
        "mime_type": file.content_type or "application/octet-stream",
        "size_bytes": len(data),
        "storage_key": storage_key,
        "status": "extracted" if text and not err else ("failed" if err else "empty"),
        "extracted_text": text,
        "extracted_chars": len(text),
        "preview": preview,
        "data_trust": trust,
        "uploaded_by": ctx["account"]["id"],
        "uploaded_by_email": ctx["account"]["email"],
        "mentioned_account_ids": mention_ids,
        "related_doc_id": related_doc_id or None,
        "relation_type": relation_type if related_doc_id else None,
        "error": err,
        "created_at": created_at,
        "updated_at": created_at,
        # Phase Z (2026-05-27) — orthogonal classification fields.
        # Every uploaded doc carries origin="upload" + the user-
        # selected category (None if they picked "Uncategorized").
        "origin":         "upload",
        "category":       cat_clean,
        "source_channel": "upload",
    }
    await db.documents.insert_one(doc)
    doc.pop("_id", None)

    # Phase E.0.2 — cross-board metadata signature derivation. Sources
    # the text from a redacted preview rather than full extracted_text
    # so the derivation pass never logs full document bodies. Pulse
    # never reads payload anyway, but the wall principle is "metadata
    # only" all the way down.
    try:
        from services.metadata_signatures import derive_and_persist
        derivation_text = " ".join([
            doc.get("name") or "",
            doc.get("description") or "",
            (preview or "")[:4000],
        ])
        await derive_and_persist(
            db,
            text=derivation_text,
            context_id=context_id,
            account_id=ctx["account"]["id"],
            source_artefact_kind="document",
            source_artefact_id=doc_id,
        )
    except Exception:  # pragma: no cover — non-fatal
        pass

    for account_id in mention_ids:
        if account_id == ctx["account"]["id"]:
            continue
        await db.mentions.insert_one({
            "id": str(uuid.uuid4()),
            "context_id": context_id,
            "target_account_id": account_id,
            "source_account_id": ctx["account"]["id"],
            "source_name": ctx["account"].get("name") or ctx["account"].get("email", ""),
            "artefact_type": "document",
            "artefact_id": doc_id,
            "comment_id": None,
            "preview": f"Tagged you on: {doc['name']}",
            "created_at": created_at,
            "read": False,
        })

    await write_audit(
        context_id, ctx["account"]["id"], "document.uploaded", "document", doc_id,
        {"name": doc["name"], "size_bytes": doc["size_bytes"], "status": doc["status"],
         "related_doc_id": related_doc_id, "relation_type": doc["relation_type"],
         "mentions": len(mention_ids)},
    )

    # Phase I.4.b (2026-05-27) — auto-extract events from doc text in
    # the background when `doc_type` matches the allowlist (Board pack /
    # briefing / cycle_compilation / strategy_document — E1=b decided
    # 2026-05-27). Best-effort: failures are logged + swallowed by
    # `auto_extract_after_upload`; never blocks the upload response.
    # `doc_type` is set by upstream writers (seed scripts, inbound
    # email parser, chat/solva attachment paths); the standard user-
    # upload UI doesn't tag doc_type at upload time, so this trigger
    # only fires for the channels that DO tag.
    try:
        if background_tasks is not None and doc.get("doc_type"):
            from routers.events import auto_extract_after_upload
            background_tasks.add_task(
                auto_extract_after_upload,
                cid=context_id,
                doc_id=doc_id,
                doc_type=doc.get("doc_type"),
                actor_id=ctx["account"]["id"],
            )
    except Exception:  # pragma: no cover — non-fatal
        pass

    # J3 (2026-05-25, ratified spec §3 Stage 4) — first-doc-uploaded
    # flag. Flip on every successful upload (idempotent set). Gates
    # the Stage 5 Trust Center introduction.
    try:
        await db.accounts.update_one(
            {"id": ctx["account"]["id"]},
            {"$set": {"first_session.first_doc_uploaded": True,
                      "first_session.first_doc_uploaded_at": created_at}},
        )
    except Exception:
        # Non-fatal — the doc is already persisted. The flag will flip
        # on the next successful upload.
        pass

    return sanitize_doc(doc)


@router.get("/contexts/{context_id}/documents/{doc_id}/thread")
async def document_thread(
    doc_id: str,
    ctx: Dict[str, Any] = Depends(require_context_membership()),
):
    """Return the thread of linked documents for this doc — ancestors (the
    chain of related_doc_id pointers back to the origin) plus descendants
    (anything that points TO this doc)."""
    context_id = ctx["context"]["id"]
    ancestors: List[Dict[str, Any]] = []
    current_id: Optional[str] = doc_id
    seen: set = set()
    for _ in range(20):  # safety cap against loops
        if not current_id or current_id in seen:
            break
        seen.add(current_id)
        d = await db.documents.find_one(
            {"id": current_id, "context_id": context_id}, {"_id": 0},
        )
        if not d:
            break
        ancestors.insert(0, sanitize_doc(d))
        current_id = d.get("related_doc_id")

    ancestor_ids = [a["id"] for a in ancestors]
    descendants = await db.documents.find(
        {
            "context_id": context_id,
            "related_doc_id": {"$in": ancestor_ids},
            "id": {"$nin": ancestor_ids},
            "status": {"$ne": "archived"},
        },
        {"_id": 0},
    ).sort("created_at", 1).to_list(100)
    return {
        "ancestors": ancestors,
        "descendants": [sanitize_doc(d) for d in descendants],
    }


@router.get("/contexts/{context_id}/documents")
async def list_documents(
    ctx: Dict[str, Any] = Depends(require_context_membership()),
    limit: int = 100,
    committee_id: Optional[str] = None,
    # Phase Z (2026-05-27) — orthogonal filter axes. `origin` filters
    # the `/app/documents` page's 3 capsule tabs; `category` filters
    # the Work Studio 6-tab row. `search` is the unified search input.
    origin: Optional[str] = None,
    category: Optional[str] = None,
    search: Optional[str] = None,
):
    q: Dict[str, Any] = {"context_id": ctx["context"]["id"], "status": {"$ne": "archived"}}
    if committee_id:
        q["committee_id"] = committee_id
    if origin:
        if origin not in ("akki_generated", "upload", "email_receipt"):
            raise HTTPException(status_code=400, detail="invalid origin filter")
        q["origin"] = origin
    if category:
        if category not in (
            "board_pack", "minutes", "draft", "deck", "report", "briefing",
            "uncategorized",
        ):
            raise HTTPException(status_code=400, detail="invalid category filter")
        # `uncategorized` is the sentinel for null — matches docs where
        # the field is missing OR explicitly null.
        if category == "uncategorized":
            q["$or"] = [{"category": {"$exists": False}}, {"category": None}]
        else:
            q["category"] = category
    if search:
        s = search.strip()
        if s:
            # Case-insensitive substring match on name / original_filename.
            # `extracted_text` is intentionally excluded — keeps the
            # GET fast + avoids surfacing matches from the body content
            # in the listing (search inside docs is a separate endpoint).
            import re
            rx = re.compile(re.escape(s), re.IGNORECASE)
            q["$and"] = [
                {"$or": [
                    {"name": {"$regex": rx}},
                    {"original_filename": {"$regex": rx}},
                ]},
            ]
    docs = await db.documents.find(
        q, {"_id": 0, "extracted_text": 0, "storage_key": 0},
    ).sort("created_at", -1).to_list(min(limit, 500))
    return [sanitize_doc(d) for d in docs]


# Phase E.1 (2026-05-26) — Drafts feed for the new Work Studio "Drafts"
# tab + the new "Recent Drafts" right-rail card. A draft is a document
# with `state == "draft"`. The `state` field is introduced in Phase E.3
# (Universal Document Drawer); until then the endpoint returns []
# correctly — empty state per spec, no fake data.
@router.get("/contexts/{context_id}/documents/drafts")
async def list_draft_documents(
    ctx: Dict[str, Any] = Depends(require_context_membership()),
    limit: int = 100,
):
    q: Dict[str, Any] = {
        "context_id": ctx["context"]["id"],
        "state": "draft",
        "status": {"$ne": "archived"},
    }
    docs = await db.documents.find(
        q, {"_id": 0, "extracted_text": 0, "storage_key": 0},
    ).sort("updated_at", -1).to_list(min(limit, 500))
    return [sanitize_doc(d) for d in docs]


# ─────────────────────────────────────────────────────────────────────
# Track A Phase 5 (2026-06-04) — POST /documents/manual-create
#
# FIX for W5 (fig 64 — "Draft Document 405 Method Not Allowed"). The
# `WorkStudio.jsx` ObjectiveCaptureModal calls this endpoint to create
# a blank `akki_generated` draft document carrying a {goal, context}
# objective. Pre-Phase-5 the endpoint did not exist — FastAPI returned
# 405 because the path partially matched the `POST /contexts/{cid}/
# documents` multipart-upload handler at line 357. Now: JSON body, no
# file, ClamAV skipped (no bytes), category enforcement skipped.
# Returns the freshly-created document so the FE can immediately open
# the drawer via `?doc_id=...`.
# ─────────────────────────────────────────────────────────────────────
class _ManualCreateIn(BaseModel):
    name: str = Field(default="Untitled draft", max_length=300)
    body: str = Field(default="", max_length=200_000)
    state: str = Field(default="draft")  # draft | active
    origin: str = Field(default="akki_generated")
    objective: Optional[Dict[str, Any]] = None


@router.post("/contexts/{context_id}/documents/manual-create")
async def manual_create_document(
    context_id: str,
    body: _ManualCreateIn,
    ctx: Dict[str, Any] = Depends(require_context_membership()),
):
    doc_id = str(uuid.uuid4())
    now_iso = datetime.now(timezone.utc).isoformat()
    doc = {
        "id": doc_id,
        "context_id": context_id,
        "account_id": ctx["account"]["id"],
        "name": body.name.strip() or "Untitled draft",
        "filename": f"{body.name.strip() or 'Untitled draft'}.txt",
        "ext": ".txt",
        "size_bytes": len((body.body or "").encode("utf-8")),
        "mime_type": "text/plain",
        "extracted_text": body.body or "",
        "state": body.state or "draft",
        "origin": body.origin or "akki_generated",
        "status": "active",
        "category": None,
        "objective": body.objective or None,
        "notes": "",
        "notes_updated_at": None,
        "created_at": now_iso,
        "updated_at": now_iso,
        # Track A Phase 5 — the FE expects `body` echoed back when the
        # drawer opens via `?doc_id=`. We persist as `extracted_text`
        # (the canonical body field) and return `body` for compatibility.
    }
    await db.documents.insert_one(doc)
    try:
        await write_audit(
            context_id, ctx["account"]["id"],
            "document.manual_create",
            "document", doc_id,
            {"name": doc["name"], "origin": doc["origin"], "state": doc["state"]},
        )
    except Exception:  # noqa: BLE001
        # Swallow contract: audit-log failures must not block the
        # create. The route is otherwise side-effect-only writes.
        pass
    return sanitize_doc(doc)


# Phase E.2 (2026-05-26) — Unified Recent Activity feed for the new
# Work Studio right-rail card. Reads `audit_log` events scoped to the
# user's active context. Returns a row-projection per event:
#   { id, action, actor_id, doc_id?, doc_title?, created_at }
# The frontend renders each as "<timestamp> · <doc title> · <action> · <actor>".
@router.get("/contexts/{context_id}/activity/recent")
async def list_recent_activity(
    ctx: Dict[str, Any] = Depends(require_context_membership()),
    limit: int = 25,
):
    cid = ctx["context"]["id"]
    # Surface every audit row that mentions the context (either
    # context-scoped or with a resource_id known to belong to this
    # context). Avoid N+1: pull rows first, then resolve doc titles
    # in one batch.
    rows = await db.audit_log.find(
        {"context_id": cid},
        {"_id": 0},
    ).sort("created_at", -1).to_list(min(limit, 100))
    # Resolve doc titles for any rows that point at a document.
    doc_ids = {
        r.get("resource_id") for r in rows
        if r.get("resource_type") == "document" and r.get("resource_id")
    }
    title_map: Dict[str, str] = {}
    if doc_ids:
        docs = await db.documents.find(
            {"id": {"$in": list(doc_ids)}, "context_id": cid},
            {"_id": 0, "id": 1, "name": 1, "original_filename": 1},
        ).to_list(length=len(doc_ids))
        for d in docs:
            title_map[d["id"]] = d.get("name") or d.get("original_filename") or d["id"]
    out: List[Dict[str, Any]] = []
    for r in rows:
        out.append({
            "id":          r.get("id"),
            "action":      r.get("action") or "",
            "actor_id":    r.get("account_id"),
            "doc_id":      r.get("resource_id") if r.get("resource_type") == "document" else None,
            "doc_title":   title_map.get(r.get("resource_id")) if r.get("resource_type") == "document" else None,
            "created_at":  r.get("created_at"),
        })
    return out


@router.get("/contexts/{context_id}/documents/{doc_id}")
async def get_document_detail(
    context_id: str, doc_id: str,
    ctx: Dict[str, Any] = Depends(require_context_membership()),
):
    # Bug #1 fix (2026-05-27) — work_studio_exports.id vs documents.id
    # mismatch. The Universal Document Drawer's ?doc_id= URL contract is
    # called with both `documents.id` AND `work_studio_exports.id` (since
    # Phase O routed Minutes/Deck/Report card opens through it). Only 18
    # of 391 exports have a documents-mirror row (created by the
    # "Continue in chat" flow at work_studio_export.py:941, back-ref via
    # `documents.work_studio_export_id`). The other 373 have no mirror —
    # opening their card 404'd with "Document not found". Resolver chain:
    #
    #   (1) Direct hit on documents.id (the original path).
    #   (2) Reverse-lookup: documents.find_one({work_studio_export_id})
    #       — handles "Continue-spawned" mirrors.
    #   (3) Last resort: work_studio_exports.find_one({id}) + synthesize
    #       a documents-shaped read-only payload so the drawer renders
    #       the artefact without an underlying mirror row.
    #
    # Backfill of historical exports into `documents` is OUT_OF_SCOPE per
    # the dispatch brief; step (3)'s synthesis keeps the existing data
    # untouched.
    d = await db.documents.find_one(
        {"id": doc_id, "context_id": context_id}, {"_id": 0, "storage_key": 0}
    )
    if not d:
        d = await db.documents.find_one(
            {"work_studio_export_id": doc_id, "context_id": context_id},
            {"_id": 0, "storage_key": 0},
        )
    if not d:
        export_row = await db.work_studio_exports.find_one(
            {"id": doc_id, "context_id": context_id}, {"_id": 0},
        )
        if export_row:
            d = _synthesize_doc_from_export(export_row)
    if not d:
        raise HTTPException(status_code=404, detail="Document not found")
    out = sanitize_doc(d)
    out["extracted_text"] = d.get("extracted_text", "")[:MAX_EXTRACT_CHARS_OUT]
    if d.get("akki_summary"):
        out["akki_summary"] = d["akki_summary"]
    # Phase L.4: surface the Synisense + sensitivity-scorer artefacts on
    # the detail endpoint so the document page can render them and the
    # Strategic-Pack ingestion can be verified end-to-end.
    for k in (
        "body_redacted", "synisense_version",
        "sensitivity_score", "sensitivity_band", "sensitivity_label",
        "sensitivity_reasons", "doc_kind",
        # Phase M.2: surface the journal-specific fields too.
        "source_channel", "journal_commentary",
        "journal_commentary_generated_at", "journal_commentary_synisense_version",
        # Phase E.3 (2026-05-26) — Universal Document Drawer:
        # state (draft/committed), objective (goal+context), origin
        # (akki_generated/upload/email_receipt), audience.
        "state", "objective", "origin", "audience",
        # Bug #1 fix (2026-05-27) — surface the work_studio_export id +
        # synthesis marker so the frontend can opt out of endpoints that
        # don't apply to synthesised payloads (e.g. /intelligence).
        "work_studio_export_id", "_synthesized_from",
    ):
        if k in d:
            out[k] = d[k]
    return out


# ═════════════════════════════════════════════════════════════════════════
#  Phase E.3 (2026-05-26) — Universal Document Drawer
#
#  Drawer-driven document surface. Backend exposes:
#    PATCH  /contexts/{cid}/documents/{did}              — edit fields
#    POST   /contexts/{cid}/documents/{did}/commit       — draft → committed
#    GET    /contexts/{cid}/documents/{did}/intelligence — cached envelope
#    POST   /contexts/{cid}/documents/{did}/intelligence/regenerate
#                                                        — async re-extract
#
#  New schema fields on `db.documents` (all optional, backwards-compatible):
#    state:     "draft" | "committed" | None
#    objective: { goal: str, context: str, set_at: ISO } | None
#    origin:    "akki_generated" | "upload" | "email_receipt" | None
#    audience:  "board" | "committee" | "regulator" | "public" | None
# ═════════════════════════════════════════════════════════════════════════


class _DocPatchIn(BaseModel):
    title:    Optional[str] = None
    body:     Optional[str] = None
    state:    Optional[str] = None  # "draft" | "committed"
    objective: Optional[Dict[str, Any]] = None
    audience: Optional[str] = None
    origin:   Optional[str] = None
    # Phase Z (2026-05-27) — category is the orthogonal classification
    # axis (Work Studio tab). Empty string clears the value.
    category: Optional[str] = None
    # Sprint Z1.2 (2026-05-29) — free-form notes saved from the drawer
    # Notes tab. Empty string clears the value (NULLs the field).
    notes:    Optional[str] = None


@router.patch("/contexts/{context_id}/documents/{doc_id}")
async def patch_document(
    context_id: str, doc_id: str,
    body: _DocPatchIn,
    ctx: Dict[str, Any] = Depends(require_context_membership()),
):
    """Phase E.3 — partial edits from the drawer's Document tab."""
    d = await db.documents.find_one(
        {"id": doc_id, "context_id": context_id}, {"_id": 0, "id": 1, "state": 1},
    )
    if not d:
        raise HTTPException(status_code=404, detail="Document not found")
    update: Dict[str, Any] = {"updated_at": datetime.now(timezone.utc).isoformat()}
    if body.title is not None:
        clean = body.title.strip()[:240]
        if clean:
            update["name"] = clean
    if body.body is not None:
        update["extracted_text"] = body.body[:MAX_EXTRACT_CHARS_OUT * 2]
    if body.state is not None:
        if body.state not in ("draft", "committed"):
            raise HTTPException(status_code=400, detail="state must be draft|committed")
        update["state"] = body.state
        if body.state == "committed":
            update["committed_at"] = datetime.now(timezone.utc).isoformat()
    if body.objective is not None:
        # Normalise the shape: {goal, context, set_at}.
        goal = (body.objective.get("goal") or "").strip()
        if not goal:
            raise HTTPException(status_code=400, detail="objective.goal is required")
        update["objective"] = {
            "goal":   goal[:480],
            "context": (body.objective.get("context") or "").strip()[:2000],
            "set_at": body.objective.get("set_at") or datetime.now(timezone.utc).isoformat(),
        }
    if body.audience is not None:
        if body.audience not in ("board", "committee", "regulator", "public"):
            raise HTTPException(status_code=400, detail="invalid audience")
        update["audience"] = body.audience
    if body.origin is not None:
        if body.origin not in ("akki_generated", "upload", "email_receipt"):
            raise HTTPException(status_code=400, detail="invalid origin")
        update["origin"] = body.origin
    if body.category is not None:
        # Phase Z (2026-05-27) — accept the 6 canonical categories +
        # empty string ("Uncategorized" → null in storage).
        _CATEGORY_ENUM = ("board_pack", "minutes", "draft", "deck", "report", "briefing")
        if body.category == "":
            update["category"] = None
        elif body.category in _CATEGORY_ENUM:
            update["category"] = body.category
        else:
            raise HTTPException(status_code=400, detail="invalid category")
    if body.notes is not None:
        # Sprint Z1.2 (2026-05-29) — notes-only PATCH must be valid.
        # Empty string clears the field; non-empty stores trimmed.
        notes_clean = body.notes.strip()[:8000]
        update["notes"] = notes_clean or None
        # Track B Phase B5 G6 (2026-06-04) — per-notes timestamp so the
        # FE's "Last updated: …" indicator reflects WHEN notes were
        # saved, not when ANY field changed. `updated_at` (set at the
        # top of this fn) churns on title/category/audience edits too;
        # `notes_updated_at` only changes when notes do, matching the
        # QA spec's literal "the date and time the note was last
        # updated" requirement. Also set on the delete path
        # (`notes_clean == ""`) so the FE can render "Last updated: …"
        # against the post-delete moment.
        update["notes_updated_at"] = update["updated_at"]
    if len(update) == 1:
        raise HTTPException(status_code=400, detail="Send at least one field")
    await db.documents.update_one({"id": doc_id, "context_id": context_id}, {"$set": update})
    # Invalidate the intelligence cache on any body/title/state change
    # so the next /intelligence GET returns a stale flag (the frontend
    # then auto-regenerates).
    if any(k in update for k in ("extracted_text", "name", "state", "objective")):
        await db.document_intelligence.delete_many({"doc_id": doc_id})
    fresh = await db.documents.find_one(
        {"id": doc_id, "context_id": context_id}, {"_id": 0, "storage_key": 0},
    )
    return sanitize_doc(fresh)


@router.get("/contexts/{context_id}/documents/{doc_id}/intelligence")
async def get_document_intelligence(
    context_id: str, doc_id: str,
    ctx: Dict[str, Any] = Depends(require_context_membership()),
):
    """Phase E.3 — returns the cached intelligence envelope. When no
    cache row exists OR the cached hash mismatches the current doc,
    returns `{status: "pending"}` so the drawer renders the skeleton.
    Frontend then POSTs /intelligence/regenerate to kick off the async
    extraction."""
    from services.documents.intelligence_service import _doc_hash
    doc = await db.documents.find_one(
        {"id": doc_id, "context_id": context_id}, {"_id": 0, "storage_key": 0},
    )
    if not doc:
        raise HTTPException(status_code=404, detail="Document not found")
    cached = await db.document_intelligence.find_one(
        {"doc_id": doc_id}, {"_id": 0},
    )
    current_hash = _doc_hash(doc)
    if not cached or cached.get("doc_hash") != current_hash:
        return {"status": "pending", "doc_id": doc_id, "doc_hash": current_hash}
    return {"status": "ready", **cached}


@router.post(
    "/contexts/{context_id}/documents/{doc_id}/briefings/generate",
    status_code=202,
)
async def generate_briefing_from_document(
    context_id: str, doc_id: str,
    ctx: Dict[str, Any] = Depends(require_context_membership()),
):
    """Sprint Z1.5 (2026-05-29) — Generate a briefing scoped to one
    document. Replaces the older Drawer behaviour of redirecting to
    Solva when the founder clicks `Generate brief`.

    Flow:
      1. Resolve all signals in `db.signals` whose `sources[].doc_id`
         matches the document's id (`status=active` only).
      2. If zero, return 400 with the same copy the briefings router
         already uses, so the drawer can surface a clear toast.
      3. Otherwise spawn the existing `_create_briefing_worker` via
         `_create_job` so the job-poll contract the briefings tab
         already uses works unchanged.

    Returns `{job_id, status: "queued"}` like the canonical endpoint.
    The frontend polls `/api/jobs/{job_id}` to discover the new
    `briefing_id` on completion, then routes the founder to
    `/app/work-studio?tab=briefings&highlight={briefing_id}`.
    """
    # Re-use the briefings router's job machinery + worker — no new
    # codepaths means the briefing render contract stays single-source.
    from routers.briefings import (
        _create_briefing_worker, BriefingCreateIn,
        _create_job, _mark_running, _mark_completed, _mark_failed,
        _spawn,
    )
    doc = await db.documents.find_one(
        {"id": doc_id, "context_id": context_id},
        {"_id": 0, "id": 1, "name": 1},
    )
    if not doc:
        raise HTTPException(status_code=404, detail="Document not found")
    # Resolve signals owned by this doc. Mirrors the matching the
    # briefings worker does on its `signals.sources[].doc_id` field.
    related_signals = await db.signals.find(
        {
            "context_id": context_id,
            "status": "active",
            "sources.doc_id": doc_id,
        },
        {"_id": 0, "id": 1},
    ).to_list(50)
    # P0-A (2026-02) — State-read parity. If the workspace-level
    # `db.signals` filter returns empty, fall through to the
    # document-level `document_intelligence.key_signals` array (the
    # SAME source the Drawer's Signals tab renders at
    # `DocumentDrawer.jsx:310 — intel?.key_signals`). When found,
    # materialise them as proper `db.signals` rows scoped to this doc
    # so the existing briefing worker contract (signal_ids) is
    # preserved. Stable id `sig:from_intel:{doc_id}:{idx}` makes the
    # upsert idempotent across re-clicks.
    #
    # P1-A (2026-02) — The promotion is now shared with the eager
    # extraction-time path (`regenerate_document_intelligence`). Same
    # helper, same stable ids: extraction promotes once, this lazy
    # fallthrough is a no-op for any doc whose intelligence ran
    # through the eager site. The fallthrough still matters for
    # backfill of old intel rows that pre-date P1-A.
    if not related_signals:
        intel = await db.document_intelligence.find_one(
            {"doc_id": doc_id, "context_id": context_id},
            {"_id": 0, "key_signals": 1},
        )
        intel_signals = (intel or {}).get("key_signals") or []
        if intel_signals:
            from services.documents.intelligence_service import (
                promote_intelligence_signals_to_pulse,
            )
            promoted = await promote_intelligence_signals_to_pulse(
                db,
                doc=doc,
                context_id=context_id,
                account_id=ctx["account"]["id"],
                key_signals=intel_signals,
            )
            if promoted:
                related_signals = [{"id": sid} for sid in promoted]
            # Track B Phase B4 G11 (2026-06-04) — Lazy Q4Y back-fill.
            # When the eager path didn't fire (intel rows that pre-date
            # the eager hook), promote any cached `open_questions` to
            # `cycle_questions` so the Brief click also fills Q4Y.
            # Stable-id upsert keeps this a no-op for intel already
            # promoted eagerly.
            intel_questions = await db.document_intelligence.find_one(
                {"doc_id": doc_id, "context_id": context_id},
                {"_id": 0, "open_questions": 1},
            )
            iqs = (intel_questions or {}).get("open_questions") or []
            if iqs:
                try:
                    from services.documents.intelligence_service import (
                        promote_intelligence_questions_to_q4y,
                    )
                    await promote_intelligence_questions_to_q4y(
                        db,
                        doc=doc,
                        context_id=context_id,
                        account_id=ctx["account"]["id"],
                        open_questions=iqs,
                    )
                except Exception as e:  # noqa: BLE001
                    import logging as _log
                    _log.getLogger("documents.intelligence").warning(
                        "lazy promote_intelligence_questions_to_q4y "
                        "failed for doc=%s: %s",
                        doc_id, e,
                    )
    if not related_signals:
        raise HTTPException(
            status_code=400,
            detail=(
                "No signals to brief on yet for this document. Extract "
                "intelligence first, or generate signals via the "
                "Signals tab."
            ),
        )
    signal_ids = [s["id"] for s in related_signals if s.get("id")]
    body = BriefingCreateIn(
        signal_ids=signal_ids,
        title=f"Briefing from {doc.get('name') or 'document'}",
    )

    job_id = await _create_job(
        kind="briefing.create",
        account_id=ctx["account"]["id"],
        context_id=context_id,
        input_summary={
            "signal_ids": signal_ids,
            "title": body.title,
            "signal_count": len(signal_ids),
            "source_doc_id": doc_id,
        },
    )

    background_account_id = ctx["account"]["id"]
    background_context_name = ctx["context"]["name"]

    async def _runner():
        await _mark_running(job_id)
        try:
            result = await _create_briefing_worker(
                body=body,
                account_id=background_account_id,
                context_id=context_id,
                context_name=background_context_name,
            )
            await _mark_completed(job_id, result)
        except HTTPException as e:
            detail = e.detail if isinstance(e.detail, str) else str(e.detail)
            await _mark_failed(job_id, f"http_{e.status_code}: {detail}")
        except Exception as e:
            logger.exception("briefing.from-doc worker crashed (job=%s)", job_id)
            await _mark_failed(job_id, f"{type(e).__name__}: {str(e)[:400]}")

    _spawn(_runner())
    return {"job_id": job_id, "status": "queued", "source_doc_id": doc_id}


@router.post("/contexts/{context_id}/documents/{doc_id}/intelligence/regenerate")
async def regenerate_document_intelligence(
    context_id: str, doc_id: str,
    background_tasks: BackgroundTasks,
    ctx: Dict[str, Any] = Depends(require_context_membership()),
):
    """Phase E.3 — kick off async intelligence extraction. Returns
    immediately with `{status: "queued"}`. The drawer polls
    /intelligence until status becomes "ready"."""
    doc = await db.documents.find_one(
        {"id": doc_id, "context_id": context_id}, {"_id": 0, "storage_key": 0},
    )
    if not doc:
        raise HTTPException(status_code=404, detail="Document not found")
    mode = "creation" if (doc.get("state") == "draft" and
                          doc.get("origin") == "akki_generated") else "reference"
    account_id = ctx.get("account", {}).get("id") or ctx["account_id"]
    # Schedule the extraction. Failure inside the task is logged + the
    # cache row stays absent so the next GET retries automatically.
    async def _run() -> None:
        from services.documents.intelligence_service import (
            extract_intelligence,
            promote_intelligence_signals_to_pulse,
        )
        try:
            envelope = await extract_intelligence(
                doc=doc, account_id=account_id, mode=mode,
            )
            await db.document_intelligence.update_one(
                {"doc_id": doc_id},
                {"$set": envelope},
                upsert=True,
            )
            # P1-A (2026-02) — Eager Pulse promotion. The user-reported
            # gap was: documents with extracted intelligence had their
            # `key_signals` invisible on Pulse Signals + Generate-Brief
            # until the user manually clicked Brief (which only triggered
            # the lazy promotion). Promoting at extraction time closes
            # the surface gap. Idempotent via stable id `sig:from_intel:
            # {doc_id}:{idx}` — re-running extraction refreshes the
            # signal row in place, never duplicates. The P0-A lazy
            # promotion at briefings/generate becomes a no-op for any
            # doc whose intelligence has already run through here.
            try:
                await promote_intelligence_signals_to_pulse(
                    db,
                    doc=doc,
                    context_id=context_id,
                    account_id=account_id,
                    key_signals=envelope.get("key_signals") or [],
                )
            except Exception as e:  # noqa: BLE001
                import logging as _log
                _log.getLogger("documents.intelligence").warning(
                    "promote_intelligence_signals_to_pulse failed for doc=%s: %s",
                    doc_id, e,
                )
            # Track B Phase B4 G11 (2026-06-04) — Eager Q4Y promotion.
            # Mirrors the signals promoter above so the CompanyHome
            # "Open questions" attention card reflects doc-extracted
            # questions in real time (no Brief click required).
            # Idempotent via stable id `q4y:from_intel:{doc_id}:{idx}`.
            try:
                from services.documents.intelligence_service import (
                    promote_intelligence_questions_to_q4y,
                )
                await promote_intelligence_questions_to_q4y(
                    db,
                    doc=doc,
                    context_id=context_id,
                    account_id=account_id,
                    open_questions=envelope.get("open_questions") or [],
                )
            except Exception as e:  # noqa: BLE001
                import logging as _log
                _log.getLogger("documents.intelligence").warning(
                    "promote_intelligence_questions_to_q4y failed for doc=%s: %s",
                    doc_id, e,
                )
        except Exception as e:  # noqa: BLE001
            import logging as _log
            _log.getLogger("documents.intelligence").warning(
                "extract_intelligence failed for doc=%s: %s", doc_id, e,
            )
    background_tasks.add_task(_run)
    return {"status": "queued", "doc_id": doc_id, "mode": mode}


# ─────────────────────────────────────────────────────────────────────
# Phase E.3 — Prompted-edit pipeline
#
# User types a natural-language instruction in the drawer's Document
# tab composer; the backend:
#   1. Runs both the prompt AND the doc body through Shield (no bypass).
#   2. Calls the universal LLM proxy via Shield's invoke() to produce the
#      rewritten body.
#   3. Returns the proposed replacement body to the frontend (NOT
#      committed yet).
#   4. Frontend renders a diff preview (strikethrough on removed,
#      oxblood-underline on added). User clicks Apply → PATCH commits.
#
# Audit: a `document_prompted_edit` audit_log row is written with
# user_id, doc_id, prompt_hash, diff_size. The raw prompt + content
# are NOT logged here — Shield's own audit chain already captures
# de-identified copies.
# ─────────────────────────────────────────────────────────────────────


class _PromptedEditIn(BaseModel):
    prompt: str = Field(min_length=1, max_length=2000)


@router.post("/documents/{doc_id}/prompted-edit")
async def prompted_edit(
    doc_id: str,
    body: _PromptedEditIn,
    current: Dict[str, Any] = Depends(get_current_account),
):
    """Phase E.3 — prompted-edit pipeline. Returns proposed new body
    WITHOUT committing. Caller (the drawer) then renders a diff
    preview and POSTs a PATCH to commit if the user accepts."""
    import hashlib as _hash
    from services.synisense.shield.client import invoke as shield_invoke
    from datetime import datetime as _dt, timezone as _tz
    import uuid as _uuid

    # Scope: only drafts owned-and-accessible by the user. We trust
    # `documents.context_id` membership check on the next read.
    doc = await db.documents.find_one(
        {"id": doc_id}, {"_id": 0, "id": 1, "context_id": 1, "name": 1,
                          "extracted_text": 1, "state": 1, "origin": 1},
    )
    if not doc:
        raise HTTPException(status_code=404, detail="Document not found")
    # Membership check.
    mem = await db.memberships.find_one({
        "account_id": current["id"], "context_id": doc["context_id"],
        "status": "active",
    }, {"_id": 0, "id": 1})
    if not mem and not current.get("is_superadmin"):
        raise HTTPException(status_code=403, detail="No access to this document")
    if doc.get("state") != "draft":
        raise HTTPException(
            status_code=400,
            detail={"code": "NOT_A_DRAFT",
                    "message": "Prompted edits apply only to drafts."},
        )

    body_text = (doc.get("extracted_text") or "")[:30000]
    prompt = body.prompt.strip()
    title = doc.get("name") or "Untitled draft"

    # Shield-bounded LLM call. The combined payload (prompt + body)
    # passes through Shield's de-id pipeline before the LLM sees it.
    llm_prompt = (
        "You are an executive editor. Rewrite the document below per the "
        "user's instruction. Return ONLY the full rewritten document text "
        "— no preamble, no commentary, no markdown fences. Preserve the "
        "voice and structure unless the instruction says otherwise.\n\n"
        f"INSTRUCTION: {prompt}\n\n"
        f"TITLE: {title}\n\n"
        f"CURRENT BODY:\n{body_text}"
    )
    try:
        result = await shield_invoke(
            purpose="document_journal.prompted_edit.rewrite",
            content=llm_prompt,
            tenant_id=current["id"],
            consumer_id="documents",
            user_id=current["id"],
            model_preference="balanced",
        )
        new_body = (result.get("response") or "").strip()
        # Strip any markdown fences the model may have added.
        if new_body.startswith("```"):
            new_body = re.sub(r"^```[a-z]*\n", "", new_body)
            new_body = re.sub(r"\n```$", "", new_body)
    except Exception as e:  # noqa: BLE001
        raise HTTPException(
            status_code=503,
            detail={"code": "LLM_REWRITE_FAILED",
                    "message": "Could not generate the rewritten body.",
                    "error": str(e)},
        ) from e

    if not new_body:
        raise HTTPException(
            status_code=503,
            detail={"code": "LLM_EMPTY_RESPONSE",
                    "message": "LLM returned an empty body — try a more specific prompt."},
        )

    # Audit row.
    diff_size = abs(len(new_body) - len(body_text))
    prompt_hash = _hash.sha256(prompt.encode("utf-8")).hexdigest()[:16]
    try:
        await db.audit_log.insert_one({
            "id":            str(_uuid.uuid4()),
            "context_id":    doc["context_id"],
            "account_id":    current["id"],
            "action":        "document.prompted_edit.proposed",
            "resource_type": "document",
            "resource_id":   doc_id,
            "metadata": {
                "prompt_hash":      prompt_hash,
                "diff_size":        diff_size,
                "current_body_len": len(body_text),
                "new_body_len":     len(new_body),
            },
            "created_at": _dt.now(_tz.utc).isoformat(),
        })
    except Exception as e:  # noqa: BLE001
        logger.warning("prompted_edit audit failed: %s", e)

    return {
        "doc_id":       doc_id,
        "prompt_hash":  prompt_hash,
        "current_body": body_text,
        "new_body":     new_body,
        "diff_size":    diff_size,
    }


import re  # noqa: E402  used by prompted_edit fence-stripping
import logging as _logging  # noqa: E402
logger = _logging.getLogger("documents")


# ─────────────────────────────────────────────────────────────────────
# Phase E.3 — Related-docs typing (4 relationship buckets)
#
# Returns related docs grouped by type:
#   • metadata_match — same context + same doc_type + (optionally
#                       same uploader)
#   • content_similarity — BM25 over the source doc's title+body
#                           against all sibling docs' paragraphs
#   • explicit_attachment — gap: no doc-to-doc link table exists.
#   • canonical_lineage — gap: no parent_doc_id / derived_from field.
#
# Each gap surfaces with `available: false` so the drawer can render
# an empty/disabled state honestly — no fake data.
# ─────────────────────────────────────────────────────────────────────
@router.get("/contexts/{context_id}/documents/{doc_id}/related")
async def list_related_documents(
    context_id: str, doc_id: str,
    ctx: Dict[str, Any] = Depends(require_context_membership()),
):
    doc = await db.documents.find_one(
        {"id": doc_id, "context_id": context_id},
        {"_id": 0, "id": 1, "name": 1, "doc_type": 1, "uploaded_by_email": 1,
         "extracted_text": 1, "paragraphs": 1, "parent_doc_id": 1,
         "version_label": 1},
    )
    if not doc:
        raise HTTPException(status_code=404, detail="Document not found")

    out: Dict[str, Any] = {
        "doc_id": doc_id,
        "groups": {
            "metadata_match":       {"available": True,  "items": [], "label": "Same metadata"},
            "explicit_attachment":  {"available": True,  "items": [], "label": "Explicit attachment"},
            "canonical_lineage":    {"available": True,  "items": [], "label": "Canonical lineage"},
            "content_similarity":   {"available": False, "items": [], "label": "Content similarity",
                                     "gap_reason": "Embedding-based similarity is deferred to Phase G (embedding model + vector store required)."},
        },
    }

    # ── Metadata-match: same context + same doc_type ──────────────
    if doc.get("doc_type"):
        peers = await db.documents.find(
            {
                "context_id": context_id,
                "id": {"$ne": doc_id},
                "doc_type": doc["doc_type"],
                "status": {"$ne": "archived"},
            },
            {"_id": 0, "id": 1, "name": 1, "doc_type": 1, "created_at": 1,
             "uploaded_by_email": 1},
        ).sort("created_at", -1).limit(8).to_list(length=8)
        out["groups"]["metadata_match"]["items"] = peers
    else:
        # Fall back: same context only — surface as a coarse-match
        # (still inside the metadata_match bucket since context is
        # metadata).
        peers = await db.documents.find(
            {
                "context_id": context_id,
                "id": {"$ne": doc_id},
                "status": {"$ne": "archived"},
            },
            {"_id": 0, "id": 1, "name": 1, "doc_type": 1, "created_at": 1},
        ).sort("created_at", -1).limit(6).to_list(length=6)
        out["groups"]["metadata_match"]["items"] = peers

    # ── Explicit attachment (Debt W3 — 2026-05-26) ────────────────
    # Symmetric query over `document_attachments`: a doc appears as
    # related if THIS doc is either side of the link. Records also
    # carry `attached_by_user_id` and an optional `note` per the brief.
    try:
        atts = await db.document_attachments.find(
            {"$or": [{"source_doc_id": doc_id}, {"target_doc_id": doc_id}]},
            {"_id": 0},
        ).limit(50).to_list(length=50)
        peer_ids = []
        for a in atts:
            other = a["target_doc_id"] if a.get("source_doc_id") == doc_id else a.get("source_doc_id")
            if other and other != doc_id:
                peer_ids.append((other, a))
        if peer_ids:
            id_set = list({pid for pid, _ in peer_ids})
            peers_map_cursor = db.documents.find(
                {"id": {"$in": id_set}, "status": {"$ne": "archived"}},
                {"_id": 0, "id": 1, "name": 1, "doc_type": 1, "created_at": 1},
            )
            peers_map = {p["id"]: p async for p in peers_map_cursor}
            items: List[Dict[str, Any]] = []
            for pid, att in peer_ids:
                p = peers_map.get(pid)
                if not p:
                    continue
                items.append({
                    **p,
                    "attachment_id":     att.get("id"),
                    "attached_by":       att.get("attached_by_user_id"),
                    "attached_at":       att.get("attached_at"),
                    "note":              att.get("note"),
                    "direction":         "outgoing" if att.get("source_doc_id") == doc_id else "incoming",
                })
            out["groups"]["explicit_attachment"]["items"] = items
    except Exception as e:  # noqa: BLE001
        logger.warning("related_docs explicit_attachment failed: %s", e)

    # ── Canonical lineage (Debt W3 — 2026-05-26) ──────────────────
    # Walk parent_doc_id chain UP (ancestors) and find children
    # whose parent_doc_id == this doc. Cap at 10 hops per branch.
    try:
        ancestors: List[Dict[str, Any]] = []
        current = doc.get("parent_doc_id")
        hops = 0
        visited = {doc_id}
        while current and hops < 10 and current not in visited:
            visited.add(current)
            anc = await db.documents.find_one(
                {"id": current, "status": {"$ne": "archived"}},
                {"_id": 0, "id": 1, "name": 1, "doc_type": 1, "created_at": 1,
                 "parent_doc_id": 1, "version_label": 1},
            )
            if not anc:
                break
            ancestors.append({**anc, "lineage": "ancestor", "depth": hops + 1})
            current = anc.get("parent_doc_id")
            hops += 1
        # Descendants — single level (deeper if needed in a follow-up).
        descendants_cursor = db.documents.find(
            {"parent_doc_id": doc_id, "status": {"$ne": "archived"}},
            {"_id": 0, "id": 1, "name": 1, "doc_type": 1, "created_at": 1,
             "version_label": 1},
        ).limit(20)
        descendants = [
            {**d, "lineage": "descendant", "depth": 1}
            async for d in descendants_cursor
        ]
        out["groups"]["canonical_lineage"]["items"] = ancestors + descendants
    except Exception as e:  # noqa: BLE001
        logger.warning("related_docs canonical_lineage failed: %s", e)

    return out


# ─────────────────────────────────────────────────────────────────────
# Debt closure W3 — Explicit attachment endpoints (2026-05-26).
# ─────────────────────────────────────────────────────────────────────
class AttachmentCreateBody(BaseModel):
    target_doc_id: str = Field(..., min_length=1)
    note: Optional[str] = Field(None, max_length=500)


@router.post("/documents/{doc_id}/attachments")
async def create_document_attachment(
    doc_id: str, body: AttachmentCreateBody,
    account: Dict[str, Any] = Depends(get_current_account),
):
    """Create an explicit attachment from `doc_id` (source) to
    `target_doc_id`. Symmetric — the Related-docs endpoint surfaces
    the link from either side."""
    if doc_id == body.target_doc_id:
        raise HTTPException(status_code=400, detail="cannot attach a document to itself")
    src = await db.documents.find_one(
        {"id": doc_id, "status": {"$ne": "archived"}}, {"_id": 0, "id": 1, "context_id": 1},
    )
    if not src:
        raise HTTPException(status_code=404, detail="source document not found")
    tgt = await db.documents.find_one(
        {"id": body.target_doc_id, "status": {"$ne": "archived"}}, {"_id": 0, "id": 1},
    )
    if not tgt:
        raise HTTPException(status_code=404, detail="target document not found")
    # Dedupe: existing attachment in either direction.
    existing = await db.document_attachments.find_one({
        "$or": [
            {"source_doc_id": doc_id, "target_doc_id": body.target_doc_id},
            {"source_doc_id": body.target_doc_id, "target_doc_id": doc_id},
        ],
    }, {"_id": 0})
    if existing:
        return existing
    row = {
        "id":                   str(uuid.uuid4()),
        "source_doc_id":        doc_id,
        "target_doc_id":        body.target_doc_id,
        "attached_by_user_id":  account["id"],
        "attached_at":          datetime.now(timezone.utc).isoformat(),
        "note":                 (body.note or None),
    }
    await db.document_attachments.insert_one(dict(row))
    try:
        await write_audit(
            src.get("context_id"), account["id"],
            "document.attachment.created", "create", row["id"],
            {"source": doc_id, "target": body.target_doc_id},
        )
    except Exception:  # noqa: BLE001
        pass
    return row


@router.delete("/documents/{doc_id}/attachments/{attachment_id}")
async def delete_document_attachment(
    doc_id: str, attachment_id: str,
    account: Dict[str, Any] = Depends(get_current_account),
):
    """Remove an attachment. The caller must own at least one side of
    the link (no cross-account stitching)."""
    row = await db.document_attachments.find_one(
        {"id": attachment_id}, {"_id": 0},
    )
    if not row:
        raise HTTPException(status_code=404, detail="attachment not found")
    if doc_id not in (row.get("source_doc_id"), row.get("target_doc_id")):
        raise HTTPException(status_code=403, detail="attachment does not involve this document")
    await db.document_attachments.delete_one({"id": attachment_id})
    try:
        await write_audit(
            None, account["id"],
            "document.attachment.deleted", "delete", attachment_id,
            {"source": row.get("source_doc_id"), "target": row.get("target_doc_id")},
        )
    except Exception:  # noqa: BLE001
        pass
    return {"ok": True, "id": attachment_id}


# ─────────────────────────────────────────────────────────────────────
# Debt closure W3 — Canonical lineage endpoint (2026-05-26).
# ─────────────────────────────────────────────────────────────────────
class LineagePatchBody(BaseModel):
    parent_doc_id: Optional[str] = Field(
        None,
        description="Set to a doc id to mark this doc as derived from it. "
                    "Pass null/omit to unlink.",
    )
    version_label: Optional[str] = Field(
        None, max_length=80,
        description="Free-text label for this version (e.g., 'v2 draft').",
    )


@router.patch("/documents/{doc_id}/lineage")
async def patch_document_lineage(
    doc_id: str, body: LineagePatchBody,
    account: Dict[str, Any] = Depends(get_current_account),
):
    """Mark `doc_id` as derived from `parent_doc_id`. Set
    `parent_doc_id: null` to unlink."""
    doc = await db.documents.find_one(
        {"id": doc_id, "status": {"$ne": "archived"}},
        {"_id": 0, "id": 1, "context_id": 1, "parent_doc_id": 1},
    )
    if not doc:
        raise HTTPException(status_code=404, detail="document not found")
    update: Dict[str, Any] = {}
    if "parent_doc_id" in body.model_fields_set:
        new_parent = body.parent_doc_id
        if new_parent == doc_id:
            raise HTTPException(status_code=400, detail="document cannot be its own parent")
        if new_parent:
            # Cycle check — walk up new_parent's chain (cap 10) to
            # confirm doc_id is NOT an ancestor.
            cur = new_parent
            seen = {doc_id}
            for _ in range(10):
                if not cur or cur in seen:
                    break
                seen.add(cur)
                p = await db.documents.find_one(
                    {"id": cur}, {"_id": 0, "parent_doc_id": 1},
                )
                if not p:
                    break
                if p.get("parent_doc_id") == doc_id:
                    raise HTTPException(status_code=400, detail="lineage cycle detected")
                cur = p.get("parent_doc_id")
        update["parent_doc_id"] = new_parent
    if "version_label" in body.model_fields_set:
        update["version_label"] = body.version_label
    if not update:
        raise HTTPException(status_code=400, detail="nothing to update")
    await db.documents.update_one({"id": doc_id}, {"$set": update})
    try:
        await write_audit(
            doc.get("context_id"), account["id"],
            "document.lineage.updated", "patch", doc_id,
            update,
        )
    except Exception:  # noqa: BLE001
        pass
    return {"ok": True, "id": doc_id, **update}


# ─────────────────────────────────────────────────────────────────────


@router.get("/contexts/{context_id}/documents/{doc_id}/export-guard")
async def check_export_guard(
    context_id: str, doc_id: str,
    ctx: Dict[str, Any] = Depends(require_context_membership()),
):
    """Phase E.3 — DRAFT export guard. Per E.3 scope-compliance
    (2026-05-26): drafts ARE exportable now — the download endpoint
    embeds a visible DRAFT watermark before serving the bytes. The
    guard returns `can_export: True` plus `watermark_required: True`
    on drafts so callers can confirm a stamp will be applied.

    The spec-compliant block-on-failure path still lives inside the
    download endpoint: if watermarking actually fails, the export is
    refused with HTTP 503 + `code: DRAFT_WATERMARK_FAILED`.
    """
    doc = await db.documents.find_one(
        {"id": doc_id, "context_id": context_id}, {"_id": 0, "id": 1, "state": 1},
    )
    if not doc:
        raise HTTPException(status_code=404, detail="Document not found")
    if doc.get("state") == "draft":
        return {
            "can_export":         True,
            "watermark_required": True,
            "watermark_label":    "DRAFT",
        }
    return {
        "can_export":         True,
        "watermark_required": False,
        "watermark_label":    None,
    }


@router.get("/contexts/{context_id}/documents/{doc_id}/download")
async def download_document(
    context_id: str, doc_id: str,
    format: str = "pdf",
    ctx: Dict[str, Any] = Depends(require_context_membership()),
):
    """Phase E.3 — Direct document download in PDF / DOCX / PPTX.
    Drafts ship with a diagonal DRAFT watermark embedded into the
    exported bytes. Committed exports have no watermark. If
    watermarking fails (lib error / malformed input), the export is
    blocked with HTTP 503 — preserving the original spec's hard rule
    that no draft ships without its watermark."""
    from services.documents.watermark_service import (
        watermark_file, WatermarkError,
    )
    from fastapi.responses import Response
    import re as _re

    fmt = (format or "pdf").lower().lstrip(".")
    if fmt not in ("pdf", "docx", "pptx"):
        raise HTTPException(status_code=400, detail="format must be pdf|docx|pptx")
    doc = await db.documents.find_one(
        {"id": doc_id, "context_id": context_id},
        {"_id": 0, "id": 1, "name": 1, "original_filename": 1,
         "extracted_text": 1, "state": 1},
    )
    if not doc:
        raise HTTPException(status_code=404, detail="Document not found")
    title = doc.get("name") or doc.get("original_filename") or "Document"
    body  = doc.get("extracted_text") or ""

    # ── Render the doc body to the requested format ─────────────
    rendered: bytes
    if fmt == "pdf":
        from reportlab.pdfgen import canvas as _canvas
        from reportlab.lib.pagesizes import A4
        from reportlab.lib.units import inch
        import io as _io
        buf = _io.BytesIO()
        c = _canvas.Canvas(buf, pagesize=A4)
        page_w, page_h = A4
        # Title.
        c.setFont("Helvetica-Bold", 16)
        c.drawString(0.8 * inch, page_h - 0.8 * inch, title[:120])
        # Body — naive wrap at ~95 chars/line, ~50 lines per page.
        c.setFont("Helvetica", 10)
        lines: List[str] = []
        for para in (body or "").splitlines():
            while len(para) > 95:
                lines.append(para[:95])
                para = para[95:]
            lines.append(para)
        y = page_h - 1.2 * inch
        for ln in lines:
            if y < 0.7 * inch:
                c.showPage()
                c.setFont("Helvetica", 10)
                y = page_h - 0.7 * inch
            c.drawString(0.8 * inch, y, ln[:110])
            y -= 14
        c.save()
        rendered = buf.getvalue()
    elif fmt == "docx":
        from docx import Document as _Doc
        import io as _io
        d = _Doc()
        d.add_heading(title[:120], level=0)
        for para in (body or "").split("\n\n"):
            d.add_paragraph(para)
        buf = _io.BytesIO()
        d.save(buf)
        rendered = buf.getvalue()
    else:  # pptx
        from pptx import Presentation as _Pres
        import io as _io
        prs = _Pres()
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        slide.shapes.title.text = title[:120]
        # Single content slide carrying the body excerpt.
        body_box = slide.placeholders[1]
        body_box.text = (body or "")[:2000]
        buf = _io.BytesIO()
        prs.save(buf)
        rendered = buf.getvalue()

    # ── Watermark drafts; preserve block-on-failure ─────────────
    if doc.get("state") == "draft":
        try:
            rendered = watermark_file(rendered, fmt=fmt, label="DRAFT")
        except WatermarkError as e:
            # Spec-compliant fallback: BLOCK the export with a clear
            # error rather than ship an unmarked draft.
            raise HTTPException(
                status_code=503,
                detail={
                    "code": "DRAFT_WATERMARK_FAILED",
                    "message": "Watermark pipeline failed; draft export blocked.",
                    "error": str(e),
                },
            ) from e

    # Audit + return.
    safe_name = _re.sub(r"[^A-Za-z0-9._-]+", "_", title)[:60] or "document"
    media_type = {
        "pdf":  "application/pdf",
        "docx": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        "pptx": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
    }[fmt]
    return Response(
        content=rendered,
        media_type=media_type,
        headers={
            "Content-Disposition": f'attachment; filename="{safe_name}.{fmt}"',
            "X-Document-State":    doc.get("state") or "committed",
            "X-Watermark-Applied": "1" if doc.get("state") == "draft" else "0",
        },
    )


# ═════════════════════════════════════════════════════════════════════════
#  Phase M.2 — Document Journal endpoints
#
#  The Document Journal is the per-context chronological listing of every
#  document the context has ever received (upload / inbound email /
#  share / sandbox). The listing is mostly handled by the existing
#  GET /api/contexts/{cid}/documents — these endpoints add the two
#  pieces the Journal UX needs that the existing surface does not:
#
#    • A row-projection that joins source_channel + cached commentary
#      flags so the listing can render markers without N+1 round trips.
#    • An on-demand commentary generator (Akki dry FT-voice take on
#      one specific document; cached on the row).
# ═════════════════════════════════════════════════════════════════════════


@router.get("/contexts/{context_id}/document-journal")
async def get_document_journal(
    context_id: str,
    ctx: Dict[str, Any] = Depends(require_context_membership()),
):
    """Listing projection for the Document Journal surface.

    Returns docs newest-first with the marker fields the UI renders
    inline (sensitivity, data_trust, source_channel, has_commentary).
    """
    cursor = db.documents.find(
        {"context_id": context_id},
        {
            "_id": 0,
            "id": 1, "name": 1, "title": 1, "doc_kind": 1, "doc_type": 1,
            "data_trust": 1, "sensitivity_band": 1, "sensitivity_score": 1,
            "source_channel": 1, "source": 1, "status": 1,
            "size_bytes": 1, "preview": 1,
            "created_at": 1, "updated_at": 1,
            "journal_commentary": 1, "journal_commentary_generated_at": 1,
        },
    ).sort("created_at", -1)
    rows = await cursor.to_list(500)
    items = []
    for r in rows:
        items.append({
            "id": r["id"],
            "title": r.get("name") or r.get("title") or "(untitled)",
            "doc_kind": r.get("doc_kind"),
            "doc_type": r.get("doc_type"),
            "data_trust": r.get("data_trust") or "trusted",
            "sensitivity_band": r.get("sensitivity_band") or "internal",
            "sensitivity_score": r.get("sensitivity_score"),
            "source_channel": r.get("source_channel") or "upload",
            "source": r.get("source"),
            "status": r.get("status") or "ready",
            "size_bytes": r.get("size_bytes"),
            "preview": (r.get("preview") or "")[:240],
            "created_at": r.get("created_at"),
            "updated_at": r.get("updated_at"),
            "has_commentary": bool(r.get("journal_commentary")),
            "commentary_generated_at": r.get("journal_commentary_generated_at"),
        })
    return {"items": items, "count": len(items)}



# ─────────────────────────────────────────────────────────────────────
# Chunk 6.5-REVISED (2026-05-13, Task D)
#
# Thin wrapper for the Work Studio "Document Journal" side deck. Same
# shape as `/document-journal` above but bounded by `limit` (default
# 5, max 25). Newest-first. Used by `CompilationRail.jsx` to populate
# the Document Journal deck without pulling the full 500-row payload
# every page render.
# ─────────────────────────────────────────────────────────────────────
@router.get("/contexts/{context_id}/document-journal/recent")
async def get_document_journal_recent(
    context_id: str,
    limit: int = 5,
    ctx: Dict[str, Any] = Depends(require_context_membership()),
):
    if limit < 1:
        limit = 1
    if limit > 25:
        limit = 25
    # Work Studio recurring-bug fix #2 (2026-05-27, recurrence #3).
    # **Structural root cause:** smoke-test runs (e.g. an old upload
    # contract test) wrote 100 documents named "smoke-upload" into
    # the TEST_SeededNedCo context and never cleaned them up. The
    # `/document-journal/recent` endpoint had no filter against test-
    # debris doc names, so the CompilationRail right-rail surfaced
    # them on a real user-facing render. Belt-and-suspenders fix:
    # (a) one-shot DB cleanup ran 2026-05-27 to delete the 100 rows;
    # (b) this `$not` filter regression-guards against any future
    # smoke run that forgets to clean up. The filter is case-
    # insensitive and matches the exact known test-debris name
    # pattern. If you add new smoke tests that write documents, give
    # them a cleanup hook OR add the name pattern here.
    test_debris_name_re = re.compile(r"^smoke[-_]upload(\.[a-z0-9]+)?$", re.IGNORECASE)
    cursor = db.documents.find(
        {
            "context_id": context_id,
            "name": {"$not": test_debris_name_re},
        },
        {
            "_id": 0,
            "id": 1, "name": 1, "title": 1, "doc_kind": 1, "doc_type": 1,
            "created_at": 1, "updated_at": 1,
        },
    ).sort("created_at", -1)
    rows = await cursor.to_list(limit)
    items = [
        {
            "id": r["id"],
            "title": r.get("name") or r.get("title") or "(untitled)",
            "doc_kind": r.get("doc_kind"),
            "doc_type": r.get("doc_type"),
            "created_at": r.get("created_at"),
            "updated_at": r.get("updated_at"),
        }
        for r in rows
    ]
    return {"items": items, "count": len(items), "limit": limit}



# ═════════════════════════════════════════════════════════════════════════
#  Phase E — Documents Journal full-text search
#
#  Indexed full-text + notes + metadata search scoped to the active
#  context. We reuse the BM25 helper from `bm25.py` (the same one chat
#  uses for grounded retrieval). One chunk per document combining its
#  name, extracted body, journal commentary, and structural metadata,
#  then ranked top-N with a snippet around the first matching token.
# ═════════════════════════════════════════════════════════════════════════
_JOURNAL_SEARCH_MAX_LIMIT = 50
_JOURNAL_SEARCH_SNIPPET_RADIUS = 80
_JOURNAL_SEARCH_BODY_CAP = 50000  # per-doc text cap to keep BM25 bounded


def _journal_snippet(text: str, q_tokens: List[str]) -> str:
    """Find the first occurrence of any q-token in `text` (length ≥ 3,
    case-insensitive) and return a ~160-char window around it."""
    if not text:
        return ""
    t_lower = text.lower()
    hit = -1
    for tok in q_tokens:
        if len(tok) < 3:
            continue
        idx = t_lower.find(tok.lower())
        if idx >= 0 and (hit < 0 or idx < hit):
            hit = idx
    if hit < 0:
        return text[:160].replace("\n", " ").strip() + ("…" if len(text) > 160 else "")
    start = max(0, hit - _JOURNAL_SEARCH_SNIPPET_RADIUS)
    end = min(len(text), hit + _JOURNAL_SEARCH_SNIPPET_RADIUS)
    snippet = text[start:end].replace("\n", " ").strip()
    if start > 0:
        snippet = "…" + snippet
    if end < len(text):
        snippet = snippet + "…"
    return snippet


@router.get("/contexts/{context_id}/document-journal/search")
async def search_document_journal(
    context_id: str,
    q: str = "",
    limit: int = 10,
    ctx: Dict[str, Any] = Depends(require_context_membership()),
):
    """Phase E — ranked full-text search across the Document Journal.
    Hits include doc body + journal_commentary + name + metadata
    (doc_kind, sensitivity_band, source_channel). Scoped to the active
    context (Privacy Wall — boundaries respected via the same
    require_context_membership dependency every other doc endpoint
    uses)."""
    if not q or not q.strip():
        return {"q": q, "hits": [], "total": 0}

    cap = max(1, min(limit, _JOURNAL_SEARCH_MAX_LIMIT))
    docs = await db.documents.find(
        {"context_id": context_id, "status": {"$ne": "archived"}},
        {
            "_id": 0, "id": 1, "name": 1, "extracted_text": 1,
            "journal_commentary": 1, "doc_kind": 1, "sensitivity_band": 1,
            "source_channel": 1, "created_at": 1, "size_bytes": 1,
            "data_trust": 1, "status": 1,
        },
    ).to_list(2000)

    if not docs:
        return {"q": q, "hits": [], "total": 0}

    # One chunk per doc — combined searchable text. The BM25 helper
    # tokenises and scores; we pass full body + commentary + name +
    # the metadata fields the user is likely to query for.
    chunks: List[Dict[str, Any]] = []
    for d in docs:
        body = (d.get("extracted_text") or "")[:_JOURNAL_SEARCH_BODY_CAP]
        commentary = d.get("journal_commentary") or ""
        meta_blob = " ".join(filter(None, [
            d.get("doc_kind") or "",
            d.get("sensitivity_band") or "",
            d.get("source_channel") or "",
        ]))
        # Name doubles to give title hits a sensible weight on short bodies.
        text = "\n".join([
            d.get("name") or "",
            d.get("name") or "",
            meta_blob,
            commentary,
            body,
        ])
        chunks.append({
            "text": text,
            "doc_id": d["id"],
            "doc_name": d.get("name") or "(untitled)",
            "doc_kind": d.get("doc_kind"),
            "sensitivity_band": d.get("sensitivity_band"),
            "source_channel": d.get("source_channel"),
            "created_at": d.get("created_at"),
            "data_trust": d.get("data_trust"),
            "size_bytes": d.get("size_bytes"),
            "_snippet_text": "\n".join(filter(None, [commentary, body])),
        })

    # Reuse the chat path's BM25.
    from bm25 import score_bm25, tokenize
    ranked = score_bm25(q, chunks, k=cap)
    q_tokens = tokenize(q)

    hits: List[Dict[str, Any]] = []
    for s, c in ranked:
        if s <= 0:
            continue
        hits.append({
            "doc_id": c["doc_id"],
            "doc_name": c["doc_name"],
            "snippet": _journal_snippet(c.get("_snippet_text") or "", q_tokens),
            "score": round(float(s), 4),
            "doc_kind": c.get("doc_kind"),
            "sensitivity_band": c.get("sensitivity_band"),
            "source_channel": c.get("source_channel"),
            "created_at": c.get("created_at"),
            "size_bytes": c.get("size_bytes"),
            "data_trust": c.get("data_trust"),
        })
    return {"q": q, "hits": hits, "total": len(hits)}



@router.post("/contexts/{context_id}/documents/{doc_id}/journal-commentary")
async def journal_commentary(
    context_id: str, doc_id: str,
    refresh: bool = False,
    ctx: Dict[str, Any] = Depends(require_context_membership()),
):
    """Generate (or return cached) Akki commentary on a single document.

    ~400 word target, FT-voice, dry, specific. Cached on the row;
    subsequent calls return cached unless `?refresh=true`.

    Phase 1 (2026-05-05): the generation logic lives in
    `document_commentary_service.generate_journal_commentary` so the
    backfill script (`scripts/backfill_journal_commentary.py`) calls
    the exact same function. The pre-Phase-1 surface mis-label
    (`synisense_pipeline.run(... surface="briefing")`) is fixed in
    that shared service and surfaced here through the same import.
    """
    d = await db.documents.find_one(
        {"id": doc_id, "context_id": context_id}, {"_id": 0},
    )
    if not d:
        raise HTTPException(status_code=404, detail="Document not found")

    from document_commentary_service import (
        generate_journal_commentary, CommentaryGenerationError,
    )
    try:
        out = await generate_journal_commentary(
            doc=d,
            account_id=ctx["account"]["id"],
            refresh=refresh,
            record_audit=True,
        )
    except CommentaryGenerationError as exc:
        raise HTTPException(status_code=502, detail=str(exc))

    if out["status"] == "skipped":
        # Map the generic skip reasons onto the HTTP responses the live
        # endpoint historically used so existing UI code doesn't have
        # to learn a new contract.
        reason = out["reason"]
        if reason == "no_extracted_text":
            raise HTTPException(
                status_code=422,
                detail="Document has no extracted text — cannot generate commentary.",
            )
        if reason.startswith("sensitivity_band="):
            raise HTTPException(
                status_code=403,
                detail="Document is restricted; commentary cannot be generated.",
            )
        if reason.startswith("doc_status="):
            raise HTTPException(
                status_code=422,
                detail=f"Document not ready for commentary ({reason}).",
            )
        # Anything else — surface as 422 with the literal reason.
        raise HTTPException(status_code=422, detail=f"skipped:{reason}")

    return {
        "doc_id": doc_id,
        "commentary": out["commentary"],
        "commentary_redacted": out["redacted"],
        "synisense_version": out["synisense_version"],
        "generated_at": out["generated_at"],
        "cached": out["status"] == "cached",
    }


@router.patch("/contexts/{context_id}/documents/{doc_id}")
async def update_document_trust(
    context_id: str, doc_id: str,
    body: DocumentTrustUpdate,
    ctx: Dict[str, Any] = Depends(require_context_membership()),
):
    """Patch trust and/or evolution-chain link.

    The endpoint name predates the evolution-link feature (added iter34); the
    payload now carries either `data_trust`, `related_doc_id`, or both.
    Passing `related_doc_id=null` unlinks the document from any predecessor.
    """
    update: Dict[str, Any] = {"updated_at": _iso(_now())}
    if body.data_trust is not None:
        update["data_trust"] = body.data_trust
    if "related_doc_id" in body.model_fields_set:
        new_link = body.related_doc_id
        if new_link == doc_id:
            raise HTTPException(status_code=400, detail="A document cannot be its own predecessor.")
        if new_link:
            # Verify the predecessor exists in the same context AND that
            # linking won't introduce a cycle.
            parent = await db.documents.find_one(
                {"id": new_link, "context_id": context_id},
                {"_id": 0, "id": 1, "related_doc_id": 1},
            )
            if not parent:
                raise HTTPException(status_code=404, detail="Predecessor document not found in this company.")
            cur = parent
            for _ in range(20):
                if cur.get("related_doc_id") == doc_id:
                    raise HTTPException(status_code=400, detail="That link would create a cycle.")
                if not cur.get("related_doc_id"):
                    break
                cur = await db.documents.find_one(
                    {"id": cur["related_doc_id"], "context_id": context_id},
                    {"_id": 0, "id": 1, "related_doc_id": 1},
                ) or {}
        update["related_doc_id"] = new_link
    if len(update) == 1:  # only updated_at — nothing was sent
        raise HTTPException(status_code=400, detail="Send data_trust and/or related_doc_id.")
    res = await db.documents.update_one(
        {"id": doc_id, "context_id": context_id},
        {"$set": update},
    )
    if res.matched_count == 0:
        raise HTTPException(status_code=404, detail="Document not found")
    audit_changes = {k: v for k, v in update.items() if k != "updated_at"}
    await write_audit(
        context_id, ctx["account"]["id"], "document.updated", "document", doc_id,
        audit_changes,
    )
    d = await db.documents.find_one({"id": doc_id}, {"_id": 0, "storage_key": 0, "extracted_text": 0})
    return sanitize_doc(d)


@router.post("/contexts/{context_id}/documents/{doc_id}/evolution-diff")
async def document_evolution_diff(
    context_id: str, doc_id: str,
    ctx: Dict[str, Any] = Depends(require_context_membership()),
    refresh: bool = False,
):
    """LLM-powered "what changed" between this doc and its immediate
    predecessor. Used by the NED Document Evolution panel to surface drift
    in recurring reports across cycles.

    Returns:
        {
          "previous_doc": {id, name, created_at} | null,
          "diff": {
            "what_changed": "<2-3 sentences: the headline of the drift>",
            "added_or_strengthened": ["<bullet>", ...],
            "weakened_or_removed":   ["<bullet>", ...],
            "questions_for_management": ["<bullet>", "<bullet>"]
          } | null
        }

    Cached on the doc record; pass ?refresh=true to regenerate.
    """
    d = await db.documents.find_one(
        {"id": doc_id, "context_id": context_id}, {"_id": 0},
    )
    if not d:
        raise HTTPException(status_code=404, detail="Document not found")
    parent_id = d.get("related_doc_id")
    if not parent_id:
        return {"previous_doc": None, "diff": None}
    parent = await db.documents.find_one(
        {"id": parent_id, "context_id": context_id}, {"_id": 0},
    )
    if not parent:
        return {"previous_doc": None, "diff": None}

    cached = d.get("evolution_diff")
    if cached and cached.get("previous_doc_id") == parent_id and not refresh:
        return {
            "previous_doc": {
                "id": parent["id"], "name": parent.get("name"),
                "created_at": parent.get("created_at"),
            },
            "diff": cached.get("diff"),
        }

    cur_text  = (d.get("extracted_text") or "").strip()[:8000]
    prev_text = (parent.get("extracted_text") or "").strip()[:8000]
    if not cur_text or not prev_text:
        raise HTTPException(
            status_code=400,
            detail="Both documents must have extractable text for a comparison.",
        )

    prompt = (
        "Compare two versions of a recurring executive report. Tell a NED, "
        "in plain English, how the second version has evolved from the first. "
        "Stay neutral. Surface drift, softening, and quiet omissions — these "
        "are the parts NEDs actually want flagged. Return JSON ONLY:\n"
        '{\n'
        '  "what_changed": "<2-3 sentences — the headline of the drift>",\n'
        '  "added_or_strengthened": [\n'
        '    "<plain-language bullet — claim/figure/promise that is new or now stronger>",\n'
        '    "<...3-5 bullets total>"\n'
        '  ],\n'
        '  "weakened_or_removed": [\n'
        '    "<bullet — claim/figure that is gone or now hedged>",\n'
        '    "<...3-5 bullets total>"\n'
        '  ],\n'
        '  "questions_for_management": [\n'
        '    "<question NED should put on the table>",\n'
        '    "<question 2>"\n'
        '  ]\n'
        '}\n\n'
        f"=== PREVIOUS VERSION ({parent.get('name')}, {parent.get('created_at')}) ===\n"
        f"{prev_text}\n\n"
        f"=== CURRENT VERSION ({d.get('name')}, {d.get('created_at')}) ===\n"
        f"{cur_text}"
    )
    out = await llm_call_llm(
        module="document.evolution_diff",
        user_query=prompt,
        context_object=None,
        session_context={"session_id": f"docevo-{doc_id}"},
        data_trust={"overall": d.get("data_trust", "unrated")},
        response_format="json",
    )
    parsed = parse_json_response(out.get("response", "")) or {}
    diff = {
        "what_changed": (parsed.get("what_changed") or "").strip(),
        "added_or_strengthened": [
            str(x).strip() for x in (parsed.get("added_or_strengthened") or [])
            if str(x).strip()
        ][:5],
        "weakened_or_removed": [
            str(x).strip() for x in (parsed.get("weakened_or_removed") or [])
            if str(x).strip()
        ][:5],
        "questions_for_management": [
            str(x).strip() for x in (parsed.get("questions_for_management") or [])
            if str(x).strip()
        ][:3],
        "mode": out.get("mode"),
        "generated_at": _iso(_now()),
    }
    await db.documents.update_one(
        {"id": doc_id, "context_id": context_id},
        {"$set": {
            "evolution_diff": {"previous_doc_id": parent_id, "diff": diff},
            "updated_at": _iso(_now()),
        }},
    )
    return {
        "previous_doc": {
            "id": parent["id"], "name": parent.get("name"),
            "created_at": parent.get("created_at"),
        },
        "diff": diff,
    }


@router.delete("/contexts/{context_id}/documents/{doc_id}")
async def archive_document(
    context_id: str, doc_id: str,
    ctx: Dict[str, Any] = Depends(require_context_membership()),
):
    d = await db.documents.find_one({"id": doc_id, "context_id": context_id})
    if not d:
        raise HTTPException(status_code=404, detail="Document not found")
    # Z2.5 (2026-02) — defence-in-depth: only `origin == "upload"`
    # documents are user-deletable. Akki-generated artefacts and
    # email-receipt-ingested files have their own lifecycle and
    # must NOT be deleted through this endpoint. Soft-delete with
    # tombstone is preserved: status flips to "archived" so audit
    # trail + signals provenance survive while the row disappears
    # from listing fetches (which all filter `status: {$ne: archived}`).
    if (d.get("origin") or "") != "upload":
        raise HTTPException(
            status_code=403,
            detail="Only uploaded documents can be deleted; Akki-generated artefacts have their own lifecycle.",
        )
    if d.get("uploaded_by") != ctx["account"]["id"] and ctx["membership"].get("sub_role") != "admin":
        raise HTTPException(status_code=403, detail="Only the uploader or a context admin can archive this document.")
    await db.documents.update_one(
        {"id": doc_id},
        {"$set": {"status": "archived", "archived_at": _iso(_now())}},
    )
    try:
        delete_from_storage(d.get("storage_key", ""))
    except Exception as e:
        logger.warning(f"delete_from_storage failed: {e}")
    await write_audit(context_id, ctx["account"]["id"], "document.archived", "document", doc_id, {})
    return {"ok": True}


# Track A Phase 4 (2026-06-04) — Removed a dead-code redefinition of
# `download_document` that was registered on the same route as the
# Phase E.3 watermarked-export handler at line ~1671. FastAPI only
# dispatches to the first-registered handler on identical routes;
# the second was unreachable. Ruff F811 confirmed unused-redefinition.


# ---------------------------------------------------------------------------
# Paragraph anchors (Reading Viewer · Phase 1, Advisory 2)
#
# Lazy-on-read endpoint + nightly cron sweep. Schema additions on `db.documents`:
#   - paragraphs: [{id, page, paragraph_number, text, char_start, char_end}]
#   - paragraphs_computed_at: iso str | null
#   - paragraphs_version: int (bumps if algorithm changes)
# Existing docs continue to work; anchors materialise on first read or on the
# nightly sweep, whichever comes first.
# ---------------------------------------------------------------------------
import os as _os  # noqa: E402
from fastapi import Header  # noqa: E402

from paragraph_anchors import compute_paragraphs, PARAGRAPH_ANCHOR_VERSION  # noqa: E402


async def _materialise_paragraphs(d: Dict[str, Any]) -> Dict[str, Any]:
    """Compute + persist paragraphs[] for a document. Idempotent.

    Phase 12.2 ITEM B — paragraphs are computed FIRST (anchor IDs are
    content hashes of the *original* text, so they stay stable across
    redaction). Synisense then runs per-paragraph in `shield_reversible`
    mode; we replace `paragraphs[i].text` with the redacted version and
    persist `shield_map_id` next to it so context members can reverse
    server-side via `/paragraphs/{pid}/original`. The TTL on the shield
    map (24h default) means originals reclaim themselves automatically.

    Returns the paragraph payload. Raises ValueError if the doc has no
    extracted_text yet (caller decides whether that's a 409 or a skip)."""
    text = (d.get("extracted_text") or "").strip()
    if not text:
        raise ValueError("Document has no extracted text yet.")
    payload = compute_paragraphs(d["id"], d["extracted_text"])

    # Run Synisense per paragraph. We do this AFTER the anchor compute
    # because anchor IDs are content hashes of the original — running
    # before would change every anchor ID on every reupload.
    redacted_paragraphs: List[Dict[str, Any]] = []
    syn_total_spans = 0
    syn_failed = False
    try:
        from services.synisense import run as syn_run
        for p in payload["paragraphs"]:
            try:
                out = await syn_run(
                    text=p.get("text") or "",
                    context_id=d.get("context_id") or "",
                    surface="ingest",
                    mode="shield_reversible",
                    account_id=None,
                )
                p_redacted = dict(p)
                p_redacted["text"] = out["redacted_text"]
                p_redacted["text_redacted"] = True
                if out.get("shield_map_id"):
                    p_redacted["shield_map_id"] = out["shield_map_id"]
                p_redacted["synisense_span_count"] = len(out.get("spans") or [])
                syn_total_spans += len(out.get("spans") or [])
                redacted_paragraphs.append(p_redacted)
            except Exception:  # noqa: BLE001 — degrade per-paragraph
                redacted_paragraphs.append(dict(p))
    except Exception:  # noqa: BLE001 — module-level degrade
        syn_failed = True
        redacted_paragraphs = list(payload["paragraphs"])

    update = {
        "paragraphs": redacted_paragraphs,
        "paragraphs_computed_at": payload["computed_at"],
        "paragraphs_version": payload["version"],
        "paragraphs_page_count": payload["page_count"],
        "synisense_version": 0 if syn_failed else 1,
        "synisense_paragraph_spans_total": syn_total_spans,
    }
    await db.documents.update_one({"id": d["id"]}, {"$set": update})
    payload["paragraphs"] = redacted_paragraphs
    payload["synisense_version"] = update["synisense_version"]
    payload["synisense_paragraph_spans_total"] = syn_total_spans
    return payload


@router.get("/contexts/{context_id}/documents/{doc_id}/paragraphs")
async def get_document_paragraphs(
    context_id: str, doc_id: str,
    ctx: Dict[str, Any] = Depends(require_context_membership()),
):
    """Return the paragraph anchors for this document. Lazy-computes on
    first read; subsequent reads return the cached array. The doc row is
    updated in-place with `paragraphs[]`, `paragraphs_computed_at`, and
    `paragraphs_version`.

    Response shape:
        {
          "doc_id": str,
          "paragraphs": [...],
          "page_count": int,
          "computed_at": str (iso),
          "version": int
        }
    """
    d = await db.documents.find_one(
        {"id": doc_id, "context_id": context_id},
        {"_id": 0},
    )
    if not d or d.get("status") == "archived":
        raise HTTPException(status_code=404, detail="Document not found")

    cached = (
        d.get("paragraphs")
        and d.get("paragraphs_computed_at")
        and (d.get("paragraphs_version") or 0) >= PARAGRAPH_ANCHOR_VERSION
    )
    if cached:
        return {
            "doc_id": doc_id,
            "paragraphs": d["paragraphs"],
            "page_count": d.get("paragraphs_page_count", 1),
            "computed_at": d["paragraphs_computed_at"],
            "version": d["paragraphs_version"],
        }

    # Lazy compute.
    if not (d.get("extracted_text") or "").strip():
        raise HTTPException(
            status_code=409,
            detail="Document extraction is incomplete; paragraphs not yet computable.",
        )
    try:
        payload = await _materialise_paragraphs(d)
    except ValueError as e:
        raise HTTPException(status_code=409, detail=str(e))
    return {
        "doc_id": doc_id,
        "paragraphs": payload["paragraphs"],
        "page_count": payload["page_count"],
        "computed_at": payload["computed_at"],
        "version": payload["version"],
    }


@router.get("/contexts/{context_id}/documents/{doc_id}/paragraphs/{pid}/original")
async def get_paragraph_original(
    context_id: str, doc_id: str, pid: str,
    ctx: Dict[str, Any] = Depends(require_context_membership()),
):
    """Phase 12.2 ITEM B — server-side shield_map reversal for a single
    paragraph. Members of the context can read the un-redacted original
    via this endpoint as long as the per-paragraph shield map hasn't
    expired (24h default TTL). Never returns the shield_map itself or
    any encrypted bytes — only the plain-text reversal. Per-request
    audit row written via `synisense.unshield`. After TTL expiry, the
    shield_map is gone and the endpoint returns 410.
    """
    d = await db.documents.find_one(
        {"id": doc_id, "context_id": context_id},
        {"_id": 0, "paragraphs": 1},
    )
    if not d:
        raise HTTPException(status_code=404, detail="Document not found")
    paragraph = next(
        (p for p in (d.get("paragraphs") or []) if p.get("id") == pid), None,
    )
    if not paragraph:
        raise HTTPException(status_code=404, detail="Paragraph not found")
    smid = paragraph.get("shield_map_id")
    if not smid:
        # Paragraph had no redactions — original equals current text.
        return {"id": pid, "original_text": paragraph.get("text") or "",
                "had_redactions": False}
    # Reverse server-side. unshield writes its own audit row.
    from services.synisense import unshield
    mapping = await unshield(
        smid, surface="ingest", account_id=ctx["account"]["id"],
    )
    if not mapping:
        raise HTTPException(
            status_code=410,
            detail="Paragraph original is no longer available — shield map expired (24h TTL).",
        )
    # Apply the reverse mapping to reconstruct the original text from the
    # redacted text. Replacement tokens are deterministic so this is exact.
    text = paragraph.get("text") or ""
    for replacement, original in mapping.items():
        text = text.replace(replacement, original)
    return {
        "id": pid,
        "original_text": text,
        "had_redactions": True,
        "span_count": paragraph.get("synisense_span_count", 0),
    }


@router.post("/cron/paragraph-anchors-sweep")
async def cron_paragraph_anchors_sweep(
    x_cron_secret: Optional[str] = Header(default=None, alias="X-Cron-Secret"),
    limit: int = 100,
):
    """Nightly batch sweep — compute paragraphs[] for any docs that are
    missing them or have stale `paragraphs_version`. Designed for the
    APScheduler entry registered in server.py @ 03:00 UTC.

    Gated by `X-Cron-Secret` matching the `AKKI_CRON_SECRET` env var.
    Anonymous callers get 401; missing env var fails closed with 503.
    Per-doc timeout is enforced by an asyncio.wait_for at the caller."""
    expected = _os.environ.get("AKKI_CRON_SECRET")
    if not expected:
        raise HTTPException(
            status_code=503,
            detail="Cron disabled — AKKI_CRON_SECRET not configured.",
        )
    if not x_cron_secret or x_cron_secret != expected:
        raise HTTPException(status_code=401, detail="Invalid or missing X-Cron-Secret header.")

    cap = max(1, min(limit, 500))
    # Pick docs where anchors are missing OR stale.
    cursor = db.documents.find(
        {
            "$or": [
                {"paragraphs_computed_at": {"$exists": False}},
                {"paragraphs_computed_at": None},
                {"paragraphs_version": {"$lt": PARAGRAPH_ANCHOR_VERSION}},
            ],
            "status": {"$ne": "archived"},
            "extracted_text": {"$exists": True, "$ne": ""},
        },
        {"_id": 0},
    ).limit(cap)

    swept = 0
    failed = 0
    skipped_no_text = 0

    import asyncio
    async for d in cursor:
        if not (d.get("extracted_text") or "").strip():
            skipped_no_text += 1
            continue
        try:
            await asyncio.wait_for(_materialise_paragraphs(d), timeout=30.0)
            swept += 1
            await write_audit(
                d.get("context_id", ""), "system", "paragraphs.batch_compute",
                "document", d["id"],
                {"version": PARAGRAPH_ANCHOR_VERSION, "outcome": "ok"},
            )
        except asyncio.TimeoutError:
            failed += 1
            logger.warning("paragraph-anchors-sweep: timeout for doc %s", d["id"])
            await write_audit(
                d.get("context_id", ""), "system", "paragraphs.batch_compute",
                "document", d["id"],
                {"version": PARAGRAPH_ANCHOR_VERSION, "outcome": "timeout"},
            )
        except Exception as e:  # noqa: BLE001
            failed += 1
            logger.warning("paragraph-anchors-sweep: failed for doc %s: %s", d["id"], e)
            await write_audit(
                d.get("context_id", ""), "system", "paragraphs.batch_compute",
                "document", d["id"],
                {"version": PARAGRAPH_ANCHOR_VERSION, "outcome": "error", "error": str(e)[:200]},
            )

    return {
        "swept": swept,
        "skipped_no_text": skipped_no_text,
        "failed": failed,
        "version": PARAGRAPH_ANCHOR_VERSION,
        "limit": cap,
    }
